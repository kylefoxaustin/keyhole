"""Build the COMBINED edge-AI-sizing exec deck: data/output/keyhole-exec-combined.pptx.

Unifies the two measured workload families into one executive narrative:
  Part 1 — Vision-Language-Action (5 VLAs, robotics control loops) — reuses the VLA
           exec deck's content slides (build_vla_exec_deck).
  Part 2 — LLM inference & training precision (the 9-model breadth + FP4) — the
           precision-roadmap chart slides as full-bleed images.

The through-line: on the SAME RTX 5090, edge viability is decided by EXECUTION SHAPE —
action-generation *topology* (VLA) and numeric *precision/format* (LLM) — not parameter
count. Both families bottleneck on the same physics: bandwidth-bound decode vs
compute-bound prefill/training, with low-precision FP compute the lever that moves the
compute wall.

Usage: python scripts/build_combined_exec_deck.py [--output data/output/keyhole-exec-combined.pptx]
"""
from __future__ import annotations
import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO / "scripts"))
from build_deck import (  # noqa: E402
    C, set_deck_size, new_slide, add_text_box, SLIDE_W_IN, SLIDE_H_IN,
)
import build_vla_exec_deck as vla  # noqa: E402

OUT_DEFAULT = REPO / "data" / "output" / "keyhole-exec-combined.pptx"

# Part 2 — the LLM/precision chart slides (PNGs already produced by the precision scripts).
PRECISION_CHARTS = [
    ("data/output/precision_blue_overtime.png",
     "INT vs FP over time — and the FP half grows",
     "deploy-precision mix per workload, 2026→2033 · grey stays INT, blue is FP (light = optional, dark = mandatory)"),
    ("data/output/precision_migration.png",
     "The 7-year precision roadmap", "INT8 holds the floor; FP becomes mandatory sub-8-bit"),
    ("data/output/precision_composition.png",
     "2026 — per-model INT/FP composition (the anchor)",
     "where the bits go today, by params / FLOPs / bytes — grey = INT, dark blue = FP-mandatory"),
    ("data/output/precision_composition_2030.png",
     "2030 — same breakout, next gen goes FP4-native",
     "flagships lead to FP4 (light blue); vision stays all-INT; edge LLMs hold INT longest"),
    ("data/output/precision_composition_2033.png",
     "2033 — FP4 default; FP-mandatory grows sub-4-bit",
     "edge LLMs move too; sub-4-bit outlier handling turns part of even the FP4 engine FP-mandatory"),
    ("data/output/precision_blue_weights.png",
     "The weights question: do they stay INT4, or go FP4?",
     "INT4 & FP4 TIE on decode (4-bit memory) but FP4 alone wins prefill 3.6× & training 5.5× — measured"),
    ("data/output/precision_5090_ladder.png",
     "Decode is bandwidth-bound", "MoE decouples footprint from speed — measured ladder"),
    ("data/output/precision_5090_breadth.png",
     "INT4 is a memory format; FP4 is a memory + compute format",
     "same-base quads + the split widens with scale, then plateaus"),
    ("data/output/precision_5090_breadth_v2.png",
     "Not a one-model fluke — INT4 ties FP4 on decode, loses on prefill, everywhere",
     "12 models · 6 architectures · every one ties on decode (4-bit memory), splits 2.7–4.5× on prefill (FP4 compute)"),
    ("data/output/precision_5090_fp4_lifecycle.png",
     "FP4 wins the whole lifecycle — inference AND training",
     "the win grows with compute-intensity (FP4 is a compute format)"),
    ("data/output/precision_5090_fp4_convergence.png",
     "FP4 training converges to BF16 quality",
     "the accuracy half of 'native FP4 training can be optimal'"),
]


def slide_combined_title(prs):
    s = new_slide(prs, bg_color=C.BG_DARK)
    add_text_box(s, Inches(0.7), Inches(2.1), Inches(12), Inches(1.2),
                 "Edge AI Silicon Sizing", font_size=44, bold=True, color=C.TEXT_WHITE)
    add_text_box(s, Inches(0.7), Inches(3.3), Inches(12), Inches(0.7),
                 "Two workload families, one physics — measured on RTX 5090, projected to edge",
                 font_size=19, color=C.ACCENT_BLUE)
    add_text_box(s, Inches(0.7), Inches(4.2), Inches(12), Inches(1.4),
                 "Part 1 · Vision-Language-Action (5 models, robotics control loops)\n"
                 "Part 2 · LLM inference & training precision (9-model breadth + native FP4)",
                 font_size=16, color=C.TEXT_DIM)
    add_text_box(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
                 "Edge viability is decided by EXECUTION SHAPE — topology (VLA) and precision "
                 "(LLM) — not parameter count. All numbers reproducible.",
                 font_size=11, color=C.TEXT_DIM)
    return s


def slide_combined_tldr(prs):
    s = new_slide(prs)
    vla.add_title_subtitle(s, "The thesis (both families)",
                           "Same silicon, same bottleneck physics — what changes is the execution shape.")
    vla.add_bullet_box(s, vla.CONTENT_LEFT, 1.6, vla.CONTENT_W, 4.9, [
        ("VLA robotics: action-generation TOPOLOGY decides edge viability, not params.",
         C.ACCENT_GREEN, True),
        "  Single-loop AR is bandwidth-walled (~1–2 Hz on NPU); dual-loop / OFT amortize one "
        "VLM forward over a whole action chunk → 50–150 Hz. A 3B and a 7B can be 50× apart at the edge.",
        ("LLM inference + training: numeric PRECISION FORMAT decides it, not params.",
         C.ACCENT_GREEN, True),
        "  INT4 is a MEMORY format — wins decode but stuck on the bf16 prefill/compute floor. "
        "FP4 (NVFP4/MXFP4) is a COMPUTE format — native Blackwell tensor cores → wins inference "
        "decode+prefill AND training (~5.5× BF16 GEMM), and converges to BF16 quality.",
        ("The shared physics: bandwidth-bound decode vs compute-bound prefill/training.",
         C.ACCENT_BLUE, True),
        "  Low-precision FP compute is the one lever that moves the compute-bound wall — for "
        "VLA denoise/prefill and for LLM prefill+training alike.",
    ])
    return s


def slide_divider(prs, kicker, title, subtitle):
    s = new_slide(prs, bg_color=C.BG_DARK)
    add_text_box(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.5),
                 kicker, font_size=16, color=C.ACCENT_BLUE, bold=True)
    add_text_box(s, Inches(0.8), Inches(3.15), Inches(11.7), Inches(1.0),
                 title, font_size=34, color=C.TEXT_WHITE, bold=True)
    add_text_box(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.7),
                 subtitle, font_size=15, color=C.TEXT_DIM)
    return s


def slide_chart(prs, png, title, caption):
    from PIL import Image
    s = new_slide(prs)
    add_text_box(s, Inches(vla.CONTENT_LEFT), Inches(0.45), Inches(vla.CONTENT_W), Inches(0.55),
                 title, font_size=20, color=C.ACCENT_BLUE, bold=True)
    # fit the chart in the body box, preserving aspect ratio
    box_l, box_t, box_w, box_h = 0.6, 1.15, SLIDE_W_IN - 1.2, 5.5
    iw, ih = Image.open(REPO / png).size
    ar = iw / ih
    if box_w / box_h > ar:          # height-bound
        h = box_h; w = h * ar
    else:                           # width-bound
        w = box_w; h = w / ar
    left = box_l + (box_w - w) / 2.0
    s.shapes.add_picture(str(REPO / png), Inches(left), Inches(box_t), Inches(w), Inches(h))
    add_text_box(s, Inches(vla.CONTENT_LEFT), Inches(SLIDE_H_IN - 0.55), Inches(vla.CONTENT_W),
                 Inches(0.4), caption, font_size=12, color=C.TEXT_DIM)
    return s


def slide_combined_takeaways(prs):
    s = new_slide(prs)
    vla.add_title_subtitle(s, "What it all says",
                           "One sizing principle across robotics VLAs and LLM inference/training.")
    vla.add_bullet_box(s, vla.CONTENT_LEFT, 1.6, vla.CONTENT_W, 4.9, [
        ("Param count is the wrong axis. Execution SHAPE is the right one.",
         C.ACCENT_GREEN, True),
        "  VLA: pick the topology (dual-loop / OFT amortization) — not the smallest model.",
        "  LLM: pick the precision FORMAT (FP4 compute, not INT4 memory) — not the lowest bit-count.",
        ("FP4 is the through-line of the LLM story — and it's a frontier in early form, today.",
         C.ACCENT_GREEN, True),
        "  Wins decode (2.2×), prefill (3.6×), training GEMM (5.5×) — the margin grows with compute "
        "intensity — and trains to BF16 quality. The whole Blackwell line bet on FP4, not INT-mixed.",
        ("Everything reproducible on one RTX 5090.",
         C.ACCENT_BLUE, True),
        "  VLA: scripts/bakeoff_vla.py + keyhole-sizer. LLM: vLLM bake-offs + qutlass/Quartet (FP4 "
        "kernels & training). All measured, not modeled.",
    ])
    return s


def build(output: Path):
    prs = Presentation()
    set_deck_size(prs)
    slide_combined_title(prs)
    slide_combined_tldr(prs)
    # Part 1 — VLA (reuse the exec deck's content slides; skip its title/tldr/takeaways)
    slide_divider(prs, "Part 1", "Vision-Language-Action",
                  "5 robotics VLAs measured on RTX 5090 → projected to NXP NPU tiers · topology decides")
    for fn in (vla.slide_topologies, vla.slide_5090_table, vla.slide_edge_projection,
               vla.slide_bw_wall, vla.slide_amortization, vla.slide_oft, vla.slide_multicam):
        fn(prs)
    # Part 2 — LLM precision (chart slides)
    slide_divider(prs, "Part 2", "LLM Inference & Training Precision",
                  "9-model breadth + native FP4 on RTX 5090 (Blackwell sm_120) · precision format decides")
    for png, title, cap in PRECISION_CHARTS:
        slide_chart(prs, png, title, cap)
    slide_combined_takeaways(prs)
    prs.save(str(output))
    n = len(prs.slides._sldIdLst)
    print(f"Wrote {output}  ({n} slides)")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--output", type=Path, default=OUT_DEFAULT)
    build(ap.parse_args().output)


if __name__ == "__main__":
    main()

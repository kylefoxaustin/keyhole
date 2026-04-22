"""
Keyhole — Results Deck Generator

Reads profiling data from data/output/ and generates a PowerPoint
presentation with architecture diagrams, profiling results, NPU
projections, and run-over-run comparison.

Usage:
    python scripts/build_deck.py
    python scripts/build_deck.py --output my_deck.pptx
    python scripts/build_deck.py --runs-dir data/output/runs
"""

import json
import math
import io
from datetime import datetime
from pathlib import Path
from typing import Optional

import sys
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import click
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# Color Palette
# ============================================================

class C:
    """Consistent color palette across all slides."""
    BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
    BG_SLIDE = RGBColor(0x16, 0x21, 0x3E)
    ACCENT_BLUE = RGBColor(0x00, 0xD4, 0xFF)
    ACCENT_GREEN = RGBColor(0x00, 0xFF, 0x88)
    ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
    ACCENT_RED = RGBColor(0xFF, 0x44, 0x44)
    ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)
    ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)    # cherry-picked from the Skippy template (storage / projected)
    ACCENT_INDIGO = RGBColor(0x63, 0x66, 0xF1)   # cherry-picked from the Skippy template (transport / emphasis)
    TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEXT_DIM = RGBColor(0xAA, 0xAA, 0xCC)
    TEXT_BRIGHT = RGBColor(0xE0, 0xE0, 0xFF)
    TABLE_HEADER = RGBColor(0x0D, 0x47, 0xA1)
    TABLE_ROW_1 = RGBColor(0x1A, 0x23, 0x40)
    TABLE_ROW_2 = RGBColor(0x22, 0x2B, 0x4A)
    FEASIBLE = RGBColor(0x00, 0xE6, 0x76)
    NOT_FEASIBLE = RGBColor(0xFF, 0x44, 0x44)

# Matplotlib equivalents
MPL_COLORS = {
    "bg": "#1A1A2E",
    "bg_slide": "#16213E",
    "blue": "#00D4FF",
    "green": "#00FF88",
    "orange": "#FF8C00",
    "red": "#FF4444",
    "purple": "#BB86FC",
    "text": "#E0E0FF",
    "dim": "#AAAACC",
    "grid": "#333355",
}


# ============================================================
# Helpers
# ============================================================

def set_slide_bg(slide, color=C.BG_SLIDE):
    """Set slide background to a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=18, color=C.TEXT_WHITE, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a styled text box.

    Writes formatting to the run-level <a:rPr> (via explicit Run) rather
    than the paragraph default <a:pPr><a:defRPr>. The defRPr approach
    renders correctly but PowerPoint's Font Color picker edits run-level
    properties — with no <a:rPr> present on the run, the override
    sometimes fails to apply when a user selects text and tries to
    change color. Putting the formatting on <a:rPr> directly gives
    PowerPoint exactly the element it expects to modify.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_title_bar(slide, title, subtitle=None):
    """Add a consistent title bar across the top of a slide."""
    # Title
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                 title, font_size=28, color=C.ACCENT_BLUE, bold=True)
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.85), Inches(2), Pt(3),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C.ACCENT_BLUE
    line.line.fill.background()

    if subtitle:
        add_text_box(slide, Inches(0.5), Inches(0.95), Inches(9), Inches(0.4),
                     subtitle, font_size=14, color=C.TEXT_DIM)


def add_styled_table(slide, left, top, width, height,
                     headers, rows, col_widths=None, highlight_rows=None,
                     font_size=10, header_font_size=11):
    """Add a table with dark themed styling.

    highlight_rows: iterable of 1-based row indices to emphasize with the indigo
    accent fill + bold white text (Skippy-template-style "target row" highlight).
    font_size: body-cell font size in pt. Drop to 8-9 for dense tables.
    header_font_size: header-row font size in pt (bold, white).
    """
    num_rows = len(rows) + 1
    num_cols = len(headers)
    highlight = set(highlight_rows or [])

    table_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, height
    )
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Style header row. Writes formatting to the run-level <a:rPr> rather
    # than paragraph defRPr so PowerPoint's Font Color picker can override
    # our defaults when a user edits the deck by hand.
    def _set_cell_formatting(cell, text, *, sz, color, bold, align=PP_ALIGN.CENTER):
        cell.text = ""   # clear default paragraph
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(sz)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Segoe UI"

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C.TABLE_HEADER
        _set_cell_formatting(cell, header, sz=header_font_size,
                              color=C.TEXT_WHITE, bold=True)

    # Style data rows
    for r_idx, row in enumerate(rows):
        row_idx = r_idx + 1  # 1-based to match header numbering
        is_hl = row_idx in highlight
        bg = C.ACCENT_INDIGO if is_hl else (C.TABLE_ROW_1 if r_idx % 2 == 0 else C.TABLE_ROW_2)
        fg = C.TEXT_WHITE if is_hl else C.TEXT_BRIGHT
        for c_idx, value in enumerate(row):
            cell = table.cell(row_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            _set_cell_formatting(cell, value, sz=font_size,
                                  color=fg, bold=is_hl)

    return table_shape


# ============================================================
# Widescreen layout helpers (16:9 template modeled on personal-ai-framework)
# ============================================================

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
CONTENT_LEFT = 0.5
CONTENT_W = 12.3
TITLE_TOP = 0.2
SUBTITLE_TOP = 0.8
CONTENT_TOP = 1.4
FOOTER_TOP = 7.1
PROJECT_FOOTER = "Keyhole — Edge AI Video Intelligence"


def set_deck_size(prs: Presentation):
    """Switch the presentation to 16:9 widescreen (13.333 x 7.5 in)."""
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)


def new_slide(prs: Presentation, bg_color=None, accent_stripe: bool = True):
    """Create a blank slide with background + optional top accent stripe."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, bg_color or C.BG_SLIDE)
    if accent_stripe:
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(SLIDE_W_IN), Pt(4),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = C.ACCENT_BLUE
        stripe.line.fill.background()
    return slide


def add_title_subtitle(slide, title, subtitle=None):
    """Standard title (22pt bold) + subtitle (13pt dim) pair.

    22pt keeps all current slide titles single-line at 12.3" width; bump to
    24pt only after verifying against the longest titles in the deck.
    """
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(TITLE_TOP),
                 Inches(CONTENT_W), Inches(0.6),
                 title, font_size=22, color=C.ACCENT_BLUE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(CONTENT_LEFT), Inches(SUBTITLE_TOP),
                     Inches(CONTENT_W), Inches(0.4),
                     subtitle, font_size=13, color=C.TEXT_DIM)


def add_bullet_box(slide, left, top, width, height, items,
                   font_size=13, font_name="Segoe UI"):
    """Put multi-paragraph bulleted content in a single textbox.

    Each item is one of:
      - "" / None: blank separator line
      - str: regular bullet at TEXT_BRIGHT
      - (text, color): colored regular line
      - (text, color, bold): colored bold line
    """
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if item is None or item == "":
            text, color, bold = "", C.TEXT_DIM, False
        elif isinstance(item, str):
            text, color, bold = item, C.TEXT_BRIGHT, False
        else:
            text = item[0]
            color = item[1] if len(item) > 1 else C.TEXT_BRIGHT
            bold = item[2] if len(item) > 2 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        if p.runs:
            r = p.runs[0]
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font_name
    return txBox


def add_pipeline_strip(slide, stages: list[str | tuple[str, bool]], top_in: float = 1.2,
                        accent_color=None):
    """Compact 5-stage pipeline state strip for bake-off slides.

    stages: list of 5 entries — either a plain string (normal/dim), or
    a tuple (label, True) for a highlighted box (what changed on this slide).
    Used to orient the reader on which component the slide is swapping.

    accent_color: override the highlight color (defaults to indigo). Pass
    C.ACCENT_RED on a "before" pipe to flag what got replaced, C.ACCENT_GREEN
    for an "after" emphasis, etc.

    Occupies vertical band [top_in, top_in + 0.55] in. Bake-off slides should
    push their main content_top to ~1.85 in to make room.
    """
    accent = accent_color if accent_color is not None else C.ACCENT_INDIGO
    total_w = CONTENT_W
    n = len(stages)
    gap = 0.12
    arrow_w = 0.22
    box_w = (total_w - (n - 1) * (gap + arrow_w)) / n
    box_h = 0.55
    x = CONTENT_LEFT

    for i, entry in enumerate(stages):
        if isinstance(entry, tuple):
            label, highlighted = entry
        else:
            label, highlighted = entry, False

        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top_in), Inches(box_w), Inches(box_h),
        )
        if highlighted:
            shp.fill.solid(); shp.fill.fore_color.rgb = accent
            shp.line.color.rgb = accent
            text_color = C.TEXT_WHITE
            bold = True
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = C.TABLE_ROW_2
            shp.line.color.rgb = C.TEXT_DIM
            text_color = C.TEXT_DIM
            bold = False
        shp.line.width = Pt(1.0)

        tf = shp.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(40000); tf.margin_right = Emu(40000)
        tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.size = Pt(10)
        r.font.bold = bold
        r.font.color.rgb = text_color
        r.font.name = "Segoe UI"

        if i < n - 1:
            # Arrow between boxes
            ar_x = x + box_w + gap / 2
            ar = add_text_box(slide, Inches(ar_x), Inches(top_in),
                              Inches(arrow_w), Inches(box_h),
                              "→", font_size=14, color=C.TEXT_DIM,
                              alignment=PP_ALIGN.CENTER)
            ar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        x += box_w + gap + arrow_w


def add_footer(slide, n: int, total: int):
    """Add consistent footer with project name + page."""
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(FOOTER_TOP),
                 Inches(CONTENT_W), Inches(0.3),
                 f"{PROJECT_FOOTER}  •  {n}/{total}",
                 font_size=9, color=C.TEXT_DIM)


def finalize_footers(prs: Presentation):
    """Add footers after all slides are built (so we know the total)."""
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        add_footer(slide, i + 1, total)


def fig_to_image_stream(fig) -> io.BytesIO:
    """Convert a matplotlib figure to a BytesIO PNG stream."""
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=200, bbox_inches="tight",
                facecolor=fig.get_facecolor(), edgecolor="none")
    plt.close(fig)
    buf.seek(0)
    return buf


# ============================================================
# Data Loaders
# ============================================================

def load_runs(runs_dir: Path) -> list[dict]:
    """Load all run JSON files, sorted by timestamp."""
    runs = []
    if runs_dir.exists():
        for f in sorted(runs_dir.glob("run_*.json")):
            with open(f) as fh:
                runs.append(json.load(fh))
    return runs


def load_sam3_reference(output_dir: Path) -> Optional[dict]:
    """Load SAM 3 reference architecture data."""
    ref_path = output_dir / "sam3_reference_architecture.json"
    if ref_path.exists():
        with open(ref_path) as f:
            return json.load(f)
    return None


def load_npu_targets() -> list[dict]:
    """Load NPU hardware target definitions."""
    from src.emulate.npu_emulator import PRESET_TARGETS
    return {name: spec.to_dict() for name, spec in PRESET_TARGETS.items()}


# ============================================================
# Chart Builders
# ============================================================

def build_architecture_diagram():
    """Build the Keyhole pipeline architecture diagram."""
    fig, ax = plt.subplots(1, 1, figsize=(10, 3.5), facecolor=MPL_COLORS["bg_slide"])
    ax.set_facecolor(MPL_COLORS["bg_slide"])
    ax.set_xlim(-0.5, 10.5)
    ax.set_ylim(-0.5, 3)
    ax.axis("off")

    boxes = [
        (0.2, 1.0, "Ingest\nFFmpeg", MPL_COLORS["blue"]),
        (2.3, 1.0, "Tier 1\nYOLO 11", MPL_COLORS["green"]),
        (4.4, 1.0, "Tier 2\nSAM 3", MPL_COLORS["orange"]),
        (6.5, 1.0, "Store\nSQLite", MPL_COLORS["purple"]),
        (8.6, 1.0, "Query\nNLQ/LLM", MPL_COLORS["blue"]),
    ]

    for x, y, label, color in boxes:
        rect = FancyBboxPatch(
            (x, y), 1.6, 1.2, boxstyle="round,pad=0.1",
            facecolor=color + "22", edgecolor=color, linewidth=2,
        )
        ax.add_patch(rect)
        ax.text(x + 0.8, y + 0.6, label, ha="center", va="center",
                fontsize=11, fontweight="bold", color=color, family="monospace")

    # Arrows
    for i in range(4):
        x_start = boxes[i][0] + 1.6
        x_end = boxes[i + 1][0]
        y_mid = 1.6
        ax.annotate("", xy=(x_end, y_mid), xytext=(x_start, y_mid),
                     arrowprops=dict(arrowstyle="->", color=MPL_COLORS["dim"],
                                     lw=2, connectionstyle="arc3,rad=0"))

    # Subtitle labels
    subtitles = [
        (1.0, 0.65, "Frame\nextraction"),
        (3.1, 0.65, "Object\ndetection"),
        (5.2, 0.65, "Concept\nenrichment"),
        (7.3, 0.65, "Metadata\npersistence"),
        (9.4, 0.65, "Natural\nlanguage"),
    ]
    for x, y, text in subtitles:
        ax.text(x, y, text, ha="center", va="center",
                fontsize=7, color=MPL_COLORS["dim"], family="monospace")

    fig.tight_layout(pad=0.5)
    return fig


def build_sam3_flop_breakdown(ref_data: dict):
    """Build SAM 3 compute distribution pie/bar chart."""
    categories = ref_data.get("category_summary", {})
    if not categories:
        return None

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 3.5),
                                     facecolor=MPL_COLORS["bg_slide"])

    # Pie chart — FLOP distribution
    ax1.set_facecolor(MPL_COLORS["bg_slide"])
    sorted_cats = sorted(categories.items(), key=lambda x: x[1]["flops"], reverse=True)
    labels = [c[0] for c in sorted_cats]
    sizes = [c[1]["flops"] for c in sorted_cats]
    colors = [MPL_COLORS["blue"], MPL_COLORS["green"], MPL_COLORS["orange"],
              MPL_COLORS["purple"], MPL_COLORS["red"]][:len(labels)]

    wedges, texts, autotexts = ax1.pie(
        sizes, labels=labels, autopct="%1.1f%%", colors=colors,
        textprops={"color": MPL_COLORS["text"], "fontsize": 9},
        pctdistance=0.75, startangle=90,
    )
    for at in autotexts:
        at.set_fontsize(8)
        at.set_color("white")
    ax1.set_title("SAM 3 — FLOP Distribution", color=MPL_COLORS["text"],
                   fontsize=12, fontweight="bold", pad=10)

    # Bar chart — GFLOPs by category
    ax2.set_facecolor(MPL_COLORS["bg_slide"])
    gflops = [c[1]["flops"] / 1e9 for c in sorted_cats]
    bars = ax2.barh(labels[::-1], gflops[::-1], color=colors[::-1], height=0.6)
    ax2.set_xlabel("GFLOPs", color=MPL_COLORS["text"], fontsize=10)
    ax2.set_title("SAM 3 — GFLOPs by Op Type", color=MPL_COLORS["text"],
                   fontsize=12, fontweight="bold", pad=10)
    ax2.tick_params(colors=MPL_COLORS["dim"], labelsize=9)
    ax2.spines[:].set_color(MPL_COLORS["grid"])
    ax2.xaxis.grid(True, color=MPL_COLORS["grid"], alpha=0.3)

    for bar, val in zip(bars, gflops[::-1]):
        ax2.text(bar.get_width() + 20, bar.get_y() + bar.get_height() / 2,
                 f"{val:.0f}", va="center", color=MPL_COLORS["text"], fontsize=8)

    fig.tight_layout(pad=1.5)
    return fig


def build_roofline_chart(targets: dict, sam3_ref: Optional[dict]):
    """Build a roofline model chart for edge hardware targets."""
    fig, ax = plt.subplots(1, 1, figsize=(10, 4.5), facecolor=MPL_COLORS["bg_slide"])
    ax.set_facecolor(MPL_COLORS["bg_slide"])

    target_styles = {
        "rtx5090": (MPL_COLORS["blue"], "o", "RTX 5090"),
        "edge_mpu": (MPL_COLORS["green"], "s", "Edge MPU Target"),
        "edge_mpu_lite": (MPL_COLORS["orange"], "^", "Edge MPU Lite"),
    }

    ai_range = np.logspace(-1, 3, 200)

    for name, spec in targets.items():
        color, marker, label = target_styles.get(name, (MPL_COLORS["dim"], "x", name))
        peak_tops = spec["tops_bf16"] * spec.get("compute_efficiency", 0.65)
        peak_bw = spec["mem_bandwidth_gbs"] * spec.get("bandwidth_efficiency", 0.8)

        # Roofline: performance = min(peak_tops, AI * peak_bw)
        # peak_bw is GB/s, AI is FLOPs/byte, so AI * peak_bw = GFLOP/s when AI in FLOP/byte
        # peak_tops is TOPS = 1000 GFLOP/s
        perf = np.minimum(peak_tops * 1000, ai_range * peak_bw)
        ax.loglog(ai_range, perf, color=color, linewidth=2, label=label)

        # Ridge point
        ridge_ai = (peak_tops * 1000) / peak_bw
        ax.axvline(x=ridge_ai, color=color, linestyle=":", alpha=0.4)

    # Plot workload points
    workloads = [
        ("YOLO 11x", 85.0, 196.0, MPL_COLORS["green"]),
        ("SAM 3 PE", 120.0, 3500.0, MPL_COLORS["orange"]),
        ("SAM 3 DETR", 15.0, 200.0, MPL_COLORS["purple"]),
        ("LLM decode", 1.0, 5.0, MPL_COLORS["red"]),
    ]

    for wl_name, ai, gflops, color in workloads:
        ax.plot(ai, gflops, "D", color=color, markersize=10, markeredgecolor="white",
                markeredgewidth=1.5, zorder=5)
        ax.annotate(wl_name, (ai, gflops), textcoords="offset points",
                    xytext=(8, 8), fontsize=9, color=color, fontweight="bold")

    ax.set_xlabel("Arithmetic Intensity (FLOPs/byte)", color=MPL_COLORS["text"], fontsize=11)
    ax.set_ylabel("Performance (GFLOP/s)", color=MPL_COLORS["text"], fontsize=11)
    ax.set_title("Roofline Model — Workloads vs Hardware Targets",
                  color=MPL_COLORS["text"], fontsize=13, fontweight="bold", pad=12)
    ax.legend(loc="lower right", fontsize=9, facecolor=MPL_COLORS["bg"],
              edgecolor=MPL_COLORS["grid"], labelcolor=MPL_COLORS["text"])
    ax.tick_params(colors=MPL_COLORS["dim"], labelsize=9)
    ax.grid(True, which="both", color=MPL_COLORS["grid"], alpha=0.2)
    ax.spines[:].set_color(MPL_COLORS["grid"])

    fig.tight_layout(pad=1.0)
    return fig


def build_latency_comparison_chart(runs: list[dict], targets: dict):
    """Build a bar chart comparing per-frame latency across runs and targets."""
    if not runs:
        return None

    fig, ax = plt.subplots(1, 1, figsize=(10, 4), facecolor=MPL_COLORS["bg_slide"])
    ax.set_facecolor(MPL_COLORS["bg_slide"])

    run_labels = []
    yolo_times = []
    sam3_times = []

    for run in runs[-8:]:  # Show last 8 runs
        video = run.get("video", {})
        height = video.get("height", 0)
        res = "4K" if height >= 2160 else f"{height}p" if height > 0 else "?"
        is_sp = run.get("pipeline", {}).get("single_pass", False)
        mode = "SP" if is_sp else ("Det" if run.get("pipeline", {}).get("detect_only") else "Seq")

        name = video.get("name", "unknown").replace(".mp4", "").replace("embedded_world_clip", "EW")
        label = f"{name}\n({res}, {mode})"
        run_labels.append(label)

        yolo_times.append(run.get("yolo", {}).get("avg_ms", 0))
        # Handle both single-pass and sequential SAM 3 timing
        sam3_data = run.get("sam3", {})
        sam3_ms = sam3_data.get("avg_inference_ms") or sam3_data.get("avg_enrichment_ms", 0)
        sam3_times.append(sam3_ms)

    x = np.arange(len(run_labels))
    width = 0.35

    bars1 = ax.bar(x - width / 2, yolo_times, width, label="YOLO 11x",
                    color=MPL_COLORS["green"], alpha=0.85)
    bars2 = ax.bar(x + width / 2, sam3_times, width, label="SAM 3",
                    color=MPL_COLORS["orange"], alpha=0.85)

    ax.set_ylabel("Latency (ms)", color=MPL_COLORS["text"], fontsize=11)
    ax.set_title("Per-Frame Inference Latency by Run",
                  color=MPL_COLORS["text"], fontsize=13, fontweight="bold", pad=12)
    ax.set_xticks(x)
    ax.set_xticklabels(run_labels, rotation=30, ha="right")
    ax.legend(facecolor=MPL_COLORS["bg"], edgecolor=MPL_COLORS["grid"],
              labelcolor=MPL_COLORS["text"])
    ax.tick_params(colors=MPL_COLORS["dim"], labelsize=9)
    ax.spines[:].set_color(MPL_COLORS["grid"])
    ax.yaxis.grid(True, color=MPL_COLORS["grid"], alpha=0.3)

    # Add value labels on bars
    for bar in bars1:
        if bar.get_height() > 0:
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5,
                    f"{bar.get_height():.1f}", ha="center", va="bottom",
                    color=MPL_COLORS["green"], fontsize=8)
    for bar in bars2:
        if bar.get_height() > 0:
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.5,
                    f"{bar.get_height():.0f}", ha="center", va="bottom",
                    color=MPL_COLORS["orange"], fontsize=8)

    fig.tight_layout(pad=1.0)
    return fig


# ============================================================
# NPU Projection Engine (mirrors npu_emulator.py logic)
# ============================================================

def project_for_deck(run: dict, targets: dict) -> list[dict]:
    """Project a run's workload onto all hardware targets."""
    from src.emulate.npu_emulator import (
        NPUEmulator, RTX_5090, HardwareSpec, WorkloadProfile,
    )

    results = []
    for name, spec_dict in targets.items():
        spec = HardwareSpec(**{k: v for k, v in spec_dict.items()
                              if k in HardwareSpec.__dataclass_fields__})
        emulator = NPUEmulator(reference=RTX_5090, target=spec)

        yolo_data = run.get("yolo", {})
        sam3_data = run.get("sam3", {})

        yolo_wl = WorkloadProfile(
            stage_name="yolo_detection", model_name="yolo11x",
            param_count=int(yolo_data.get("params_m", 57) * 1e6),
            model_size_bytes=int(yolo_data.get("params_m", 57) * 1e6 * 2),
            precision="fp16", gflops_per_inference=196.0,
            arithmetic_intensity=85.0,
            measured_latency_ms=yolo_data.get("avg_ms", 10.0),
            measured_gpu=RTX_5090.name,
            peak_activation_bytes=int(0.2e9),
        )

        sam3_wl = WorkloadProfile(
            stage_name="sam3_enrichment", model_name="sam3_full",
            param_count=848_000_000,
            model_size_bytes=int(848e6 * 2),
            precision="bf16", gflops_per_inference=350.0,
            arithmetic_intensity=120.0,
            measured_latency_ms=sam3_data.get("avg_enrichment_ms", 30.0),
            measured_gpu=RTX_5090.name,
            peak_activation_bytes=int(1.0e9),
        )

        yolo_r = emulator.project_workload(yolo_wl)
        sam3_r = emulator.project_workload(sam3_wl)

        combined_ms = yolo_r.projected_latency_ms + sam3_r.projected_latency_ms
        combined_fps = 1000.0 / combined_ms if combined_ms > 0 else 0

        results.append({
            "target": spec.name,
            "target_key": name,
            "yolo_projected_ms": yolo_r.projected_latency_ms,
            "yolo_bottleneck": yolo_r.bottleneck,
            "sam3_projected_ms": sam3_r.projected_latency_ms,
            "sam3_bottleneck": sam3_r.bottleneck,
            "combined_ms": combined_ms,
            "combined_fps": combined_fps,
            "feasible_1fps": combined_ms < 1000,
            "feasible_5fps": combined_ms < 200,
            "yolo_fits": yolo_r.fits_in_memory,
            "sam3_fits": sam3_r.fits_in_memory,
            "tops": spec.tops_bf16,
            "bw_gbs": spec.mem_bandwidth_gbs,
            "mem_gb": spec.mem_capacity_gb,
            "tdp_w": spec.tdp_watts,
        })

    return results


# ============================================================
# Slide Builders
# ============================================================

def slide_title(prs: Presentation):
    """Slide 1: Title slide."""
    slide = new_slide(prs, bg_color=C.BG_DARK)

    add_text_box(slide, Inches(0.5), Inches(2.5), Inches(CONTENT_W), Inches(1.2),
                 "KEYHOLE", font_size=64, color=C.ACCENT_BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(3.7), Inches(CONTENT_W), Inches(0.6),
                 "Open-Source AI Key Prototype", font_size=24,
                 color=C.TEXT_WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(CONTENT_W), Inches(0.5),
                 "Edge AI Video Intelligence  •  Model Bake-Off  •  NPU Feasibility",
                 font_size=15, color=C.TEXT_DIM, alignment=PP_ALIGN.CENTER)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(5.1), Inches(2.3), Pt(2),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C.ACCENT_BLUE
    line.line.fill.background()

    timestamp = datetime.now().strftime("%B %d, %Y")
    add_text_box(slide, Inches(0.5), Inches(5.4), Inches(CONTENT_W), Inches(0.4),
                 f"Generated {timestamp}  •  github.com/kylefoxaustin/keyhole",
                 font_size=12, color=C.TEXT_DIM, alignment=PP_ALIGN.CENTER)


def slide_exec_summary(prs: Presentation):
    """The world's best single-slide summary: how to apply SAM 3-like
    capabilities at the edge. Sits right after the title slide."""
    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        "How to Apply SAM 3-like Capabilities at the Edge",
        "Conclusions from every bake-off in this deck, condensed onto one page",
    )

    # ── Hero stat band (3 cards) ───────────────────────────────────────────
    card_y = 1.3
    card_h = 0.95
    card_w = (CONTENT_W - 0.5) / 3
    cards = [
        ("0.4 FPS",    "SAM 3 BF16 baseline",       "Bandwidth-bound on 134.4 GB/s — not feasible", C.ACCENT_RED),
        ("36 FPS",     "Shipping stack (720p)",     "Hybrid V2 + YOLO-seg FP8 + CLIP FP8, all TensorRT",       C.ACCENT_GREEN),
        ("90×",        "Edge FPS improvement",       "Architectural change — not just quantization",        C.ACCENT_INDIGO),
    ]
    for i, (big, label, note, col) in enumerate(cards):
        x = CONTENT_LEFT + i * (card_w + 0.25)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x), Inches(card_y),
                                      Inches(card_w), Inches(card_h))
        shp.fill.solid(); shp.fill.fore_color.rgb = C.BG_SLIDE
        shp.line.color.rgb = col; shp.line.width = Pt(2)
        add_text_box(slide, Inches(x), Inches(card_y + 0.05),
                     Inches(card_w), Inches(0.45),
                     big, font_size=28, color=col, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x), Inches(card_y + 0.5),
                     Inches(card_w), Inches(0.25),
                     label, font_size=11, color=C.TEXT_BRIGHT, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x), Inches(card_y + 0.72),
                     Inches(card_w), Inches(0.25),
                     note, font_size=9, color=C.TEXT_DIM,
                     alignment=PP_ALIGN.CENTER)

    # ── Two pipeline strips: BEFORE (target) vs AFTER (shipping) ───────────
    # BEFORE — SAM 3 highlighted red as the dropped component
    before_label_y = 2.4
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(before_label_y),
                 Inches(CONTENT_W), Inches(0.22),
                 "BEFORE — target pipeline (SAM 3 BF16 baseline, 0.4 FPS)",
                 font_size=10, color=C.ACCENT_RED, bold=True)
    add_pipeline_strip(slide, [
        ("FFmpeg ingest", False),
        ("YOLO 11x", False),
        ("SAM 3 BF16", True),
        ("SQLite + FTS5", False),
        ("NLQ / LLM", False),
    ], top_in=before_label_y + 0.25, accent_color=C.ACCENT_RED)

    # AFTER — shipping recipe with TRT-compiled halves highlighted indigo
    after_label_y = 3.4
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(after_label_y),
                 Inches(CONTENT_W), Inches(0.22),
                 "AFTER — shipping recipe (36 FPS, 90× edge FPS, real-time)",
                 font_size=10, color=C.ACCENT_GREEN, bold=True)
    add_pipeline_strip(slide, [
        ("FFmpeg ingest", False),
        ("YOLO-seg FP8 (TRT)", True),
        ("CLIP FP8 (TRT) @ 1 Hz", True),
        ("SQLite + FTS5", False),
        ("Qwen3-30B-A3B MoE", False),
    ], top_in=after_label_y + 0.25)

    # ── Two-column DO / DON'T ──────────────────────────────────────────────
    col_top = 4.5
    col_h = 1.55
    col_w = (CONTENT_W - 0.25) / 2

    add_bullet_box(slide, CONTENT_LEFT, col_top, col_w, col_h, [
        ("DO — the full recipe", C.ACCENT_GREEN, True),
        ("• Replace SAM 3 with Hybrid V2 — YOLO-seg-s (det+seg, 10M) + OpenCLIP ViT-B-32 (open-vocab tags)",),
        ("• Compile BOTH halves with TensorRT FP8 on Blackwell-class silicon (recall 1.00 / top-1 agree 0.96)",),
        ("• Debounce CLIP at 1 Hz — cheap headroom; multi-stream: YOLO batch=N (4 streams → 26 FPS each)",),
        ("• Co-host Qwen3-30B-A3B MoE (3B active) on the NPU for NLQ — duty-cycle share for short answers",),
    ], font_size=9)

    add_bullet_box(slide, CONTENT_LEFT + col_w + 0.25, col_top, col_w, col_h, [
        ("DON'T — ruled out by bake-off", C.ACCENT_RED, True),
        ("• Run SAM 3 BF16 and hope — 0.4 FPS, FP8 activations only get you to ~1.2 FPS, still dead",),
        ("• INT8 weight-only quant — doesn't touch bandwidth-bound activation traffic, zero edge gain",),
        ("• Cut SAM 3 resolution or prompts — RoPE locked, edge stays bandwidth-bound regardless",),
        ("• torchao FP8 on Conv-only models, or generative LLM on a busy vision NPU (RAG murders it)",),
    ], font_size=9)

    # ── NPU tier sizing ────────────────────────────────────────────────────
    tier_y = 6.15
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(tier_y),
                 Inches(CONTENT_W), Inches(0.25),
                 "Buy the right NPU class for the job",
                 font_size=11, color=C.ACCENT_INDIGO, bold=True)
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(tier_y + 0.3),
                     Inches(CONTENT_W), Inches(0.8),
                     ["NPU tier", "Memory bus", "Vision (1-stream)", "Vision (4-stream, batch=4)",
                      "LLM Q4_K_M decode", "Good for"],
                     [
                         ["NPU Low-LP5",  "64-bit LPDDR5 @ 6.4 GT/s",    "~15 FPS",  "~10 FPS each",
                          "29 tok/s",  "Dense INT8-only silicon (NXP Neutron class)"],
                         ["NPU Mid",  "128-bit LPDDR5X @ 8.4 GT/s",  "36 FPS",   "26 FPS each",
                          "38 tok/s",  "Live multi-stream + occasional LLM"],
                         ["NPU High", "high-bin LPDDR5X",            "~48 FPS",  "~35 FPS each",
                          "50 tok/s",  "Live vision + sustained LLM"],
                     ],
                     highlight_rows=[2],   # NPU Mid — the shipping target
                     font_size=9, header_font_size=10)

    # ── Footer assumption note ─────────────────────────────────────────
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(7.25),
                 Inches(CONTENT_W), Inches(0.2),
                 "Assumes vision/LLM time-slice on the NPU",
                 font_size=8, color=C.TEXT_DIM,
                 alignment=PP_ALIGN.CENTER)


def slide_architecture(prs: Presentation):
    """Slide 2: Pipeline architecture diagram."""
    slide = new_slide(prs)
    add_title_subtitle(slide, "Pipeline Architecture",
                       "5-stage edge AI video intelligence pipeline")

    fig = build_architecture_diagram()
    img_stream = fig_to_image_stream(fig)
    slide.shapes.add_picture(img_stream, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                              width=Inches(CONTENT_W))

    add_bullet_box(slide, CONTENT_LEFT, 4.5, CONTENT_W, 2.0, [
        ("Pipeline stages and scale", C.ACCENT_BLUE, True),
        ("• Tier 1 — YOLO 11x: 56.9M params, ~196 GFLOPs/frame — detection + bounding boxes",),
        ("• Tier 2 — SAM 3 concept segmentation: 840.5M params, ~4,175 GFLOPs/frame — open-vocabulary masks",),
        ("• NLQ — Claude API / Ollama / Skippy: 3B–8B int4 models — natural-language questions over the metadata",),
        ("• Store — SQLite + FTS5 + optional vector embeddings — drives the /api/events search endpoint",),
    ], font_size=13)


def slide_sam3_reference(prs: Presentation, ref_data: dict):
    """Slide 3: SAM 3 reference architecture breakdown."""
    slide = new_slide(prs)

    summary = ref_data.get("model_summary", {})
    total_gflops = summary.get("total_gflops", 0)
    total_params = summary.get("total_params", 0)
    num_layers = len(ref_data.get("layers", []))

    add_title_subtitle(slide, "SAM 3 Reference Architecture",
                       f"{total_params/1e6:.0f}M params  •  {total_gflops:.0f} GFLOPs  •  {num_layers} layers")

    fig = build_sam3_flop_breakdown(ref_data)
    if fig:
        img_stream = fig_to_image_stream(fig)
        slide.shapes.add_picture(img_stream, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                                  width=Inches(CONTENT_W))

    components = ref_data.get("model_summary", {}).get("components", {})
    if components:
        headers = ["Component", "Params", "Role"]
        rows = [[name.replace("_", " ").title(), f"{info['params_m']}M", info["role"]]
                for name, info in components.items()]
        add_styled_table(slide, Inches(CONTENT_LEFT), Inches(5.2),
                         Inches(CONTENT_W), Inches(1.4),
                         headers, rows,
                         col_widths=[Inches(2.5), Inches(1.3), Inches(8.5)])


def slide_roofline(prs: Presentation, targets: dict, sam3_ref: Optional[dict]):
    """Slide 4: Roofline model."""
    slide = new_slide(prs)
    add_title_subtitle(slide, "Roofline Model — Compute vs Bandwidth",
                       "Workload placement determines the bottleneck on each hardware target")

    fig = build_roofline_chart(targets, sam3_ref)
    img_stream = fig_to_image_stream(fig)
    # Center the chart at ~9" wide so it doesn't collide with the bullet box.
    slide.shapes.add_picture(img_stream, Inches(2.2), Inches(CONTENT_TOP),
                              width=Inches(9.0))

    add_bullet_box(slide, CONTENT_LEFT, 5.5, CONTENT_W, 1.4, [
        ("Reading the chart", C.ACCENT_BLUE, True),
        "• SAM 3 PE sits high on the bandwidth ramp — memory-bandwidth-bound on every target",
        "• YOLO 11x is past the ridge point — compute-bound, not traffic-bound",
        "• Edge bandwidth gap (RTX 5090 → Edge MPU Target) is ~15×; compute gap is only ~1.2×",
    ], font_size=12)


def slide_run_results(prs: Presentation, run: dict, run_index: int):
    """Per-run result slide with profiling data."""
    slide = new_slide(prs)

    video = run.get("video", {})
    video_name = video.get("name", "unknown")
    run_id = run.get("run_id", "unknown")
    width = video.get("width", 0)
    height = video.get("height", 0)

    if height >= 2160: res_label = "4K"
    elif height >= 1080: res_label = "1080p"
    elif height >= 720: res_label = "720p"
    elif height > 0: res_label = f"{height}p"
    else: res_label = "?"

    is_single_pass = run.get("pipeline", {}).get("single_pass", False)
    mode_label = ("Single-Pass" if is_single_pass
                  else ("Detect Only" if run.get("pipeline", {}).get("detect_only") else "YOLO+SAM3"))

    add_title_subtitle(slide, f"Test Run — {video_name} ({res_label})",
                       f"{width}x{height}  •  {mode_label}  •  {video.get('extract_fps', '?')} FPS extraction  •  Run: {run_id}")

    yolo = run.get("yolo", {})
    sam3 = run.get("sam3", {})
    pipeline = run.get("pipeline", {})

    # Per-stage YOLO stats aren't captured in single-pass runs (YOLO + SAM 3
    # share one forward). Show "—" instead of fake zeros.
    def _ms(val, fmt=".1f"):
        return f"{val:{fmt}} ms" if val and val > 0 else "—"

    yolo_avg = yolo.get("avg_ms") or yolo.get("avg_inference_ms")
    yolo_p95 = yolo.get("p95_ms") or yolo.get("p95_inference_ms")
    yolo_p99 = yolo.get("p99_ms") or yolo.get("p99_inference_ms")
    yolo_params = yolo.get("params_m") or yolo.get("model_params_m")

    yolo_rows = [
        ["Model",         yolo.get("model", "yolo11x.pt")],
        ["Avg Inference", _ms(yolo_avg)],
        ["P95 Latency",   _ms(yolo_p95)],
        ["P99 Latency",   _ms(yolo_p99)],
        ["Parameters",    f"{yolo_params:.1f}M" if yolo_params else "— (folded into SAM 3 single-pass)"],
    ]

    # Accept both enrichment_ms and inference_ms field naming in SAM 3 runs.
    sam3_avg = sam3.get("avg_enrichment_ms") or sam3.get("avg_inference_ms")
    sam3_p95 = sam3.get("p95_enrichment_ms") or sam3.get("p95_inference_ms")
    sam3_p99 = sam3.get("p99_enrichment_ms") or sam3.get("p99_inference_ms")

    sam3_rows = [
        ["Model",            sam3.get("model", "not loaded")],
        ["Avg Enrichment",   _ms(sam3_avg, ".0f")],
        ["P95 Latency",      _ms(sam3_p95, ".0f")],
        ["P99 Latency",      _ms(sam3_p99, ".0f")],
        ["Parameters",       f"{sam3.get('model_params_m', 0):.1f}M" if sam3.get('model_params_m') else "—"],
        ["Frames Profiled",  str(sam3.get("total_frames", 0))],
    ]
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP), Inches(6), Inches(0.4),
                 "YOLO 11x Detection", font_size=14, color=C.ACCENT_GREEN, bold=True)
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.85),
                     Inches(6), Inches(2.0),
                     ["Metric", "Value"], yolo_rows,
                     col_widths=[Inches(2.3), Inches(3.7)])

    add_text_box(slide, Inches(6.8), Inches(CONTENT_TOP), Inches(6), Inches(0.4),
                 "SAM 3 Enrichment", font_size=14, color=C.ACCENT_ORANGE, bold=True)
    add_styled_table(slide, Inches(6.8), Inches(1.85),
                     Inches(6), Inches(2.0),
                     ["Metric", "Value"], sam3_rows,
                     col_widths=[Inches(2.3), Inches(3.7)])

    # Pipeline summary — single-row wide table
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(4.2), Inches(CONTENT_W), Inches(0.4),
                 "Pipeline Summary", font_size=14, color=C.ACCENT_BLUE, bold=True)
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(4.65),
                     Inches(CONTENT_W), Inches(1.1),
                     ["Frames", "Detections", "Pipeline Time", "Throughput", "Mode"],
                     [[
                         str(pipeline.get("total_frames", 0)),
                         str(pipeline.get("total_detections", 0)),
                         f"{pipeline.get('total_seconds', 0):.1f}s",
                         f"{pipeline.get('fps', 0):.2f} FPS",
                         "Detect Only" if pipeline.get("detect_only") else "Full (YOLO+SAM3)",
                     ]])


def slide_npu_projections(prs: Presentation, run: dict, targets: dict):
    """NPU projection slide for a given run."""
    slide = new_slide(prs)

    video_name = run.get("video", {}).get("name", "unknown")
    add_title_subtitle(slide, "Edge NPU Projections",
                       f"{video_name}  •  Measured on RTX 5090, projected to edge targets")

    projections = project_for_deck(run, targets)

    headers = ["Target", "TOPS", "BW (GB/s)", "YOLO", "SAM 3", "Combined", "FPS",
               "1 FPS?", "5 FPS?", "TDP"]
    rows = []
    for p in projections:
        rows.append([
            p["target"], f"{p['tops']:.0f}", f"{p['bw_gbs']:.0f}",
            f"{p['yolo_projected_ms']:.1f}ms", f"{p['sam3_projected_ms']:.1f}ms",
            f"{p['combined_ms']:.1f}ms", f"{p['combined_fps']:.0f}",
            "YES" if p["feasible_1fps"] else "NO",
            "YES" if p["feasible_5fps"] else "NO",
            f"{p['tdp_w']:.0f}W",
        ])

    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                     Inches(CONTENT_W), Inches(1.8),
                     headers, rows)

    edge_mpu = next((p for p in projections if p["target_key"] == "edge_mpu"), None)
    items = [("Key finding — Edge MPU Target", C.ACCENT_GREEN, True)]
    if edge_mpu:
        feasible_1 = "FEASIBLE" if edge_mpu["feasible_1fps"] else "NOT FEASIBLE"
        feasible_5 = "FEASIBLE" if edge_mpu["feasible_5fps"] else "NOT FEASIBLE"
        items.extend([
            f"• Combined latency: {edge_mpu['combined_ms']:.1f} ms  →  {edge_mpu['combined_fps']:.0f} FPS",
            f"• SAM 3 is {edge_mpu['sam3_bottleneck']}-bound on this hardware",
            (f"• 1 FPS extraction: {feasible_1}",
             C.FEASIBLE if edge_mpu["feasible_1fps"] else C.NOT_FEASIBLE, True),
            (f"• 5 FPS extraction: {feasible_5}",
             C.FEASIBLE if edge_mpu["feasible_5fps"] else C.NOT_FEASIBLE, True),
        ])
    add_bullet_box(slide, CONTENT_LEFT, 3.6, CONTENT_W, 2.8, items, font_size=12)


def slide_run_comparison(prs: Presentation, runs: list[dict], targets: dict):
    """Comparison chart across all runs."""
    slide = new_slide(prs)
    add_title_subtitle(slide, "Run Comparison",
                       f"{len(runs)} test runs  •  Inference latency on RTX 5090")

    fig = build_latency_comparison_chart(runs, targets)
    if fig:
        img_stream = fig_to_image_stream(fig)
        slide.shapes.add_picture(img_stream, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                                  width=Inches(CONTENT_W))


def slide_bandwidth_wall(prs: Presentation):
    """Slide: Why SAM 3 hits a bandwidth wall on edge hardware."""
    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "The Bandwidth Wall",
                       "SAM 3 is deeply memory-bandwidth-bound — TOPS don't matter")

    add_bullet_box(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_W, 5.4, [
        ("MEASURED on RTX 5090 (209 TOPS, 1792 GB/s, 72 MB L2 cache)", C.ACCENT_BLUE, True),
        "• GPU kernel time 102 ms  •  wall clock 107 ms  •  CPU overhead only 5 ms",
        ("• Theoretical compute floor (350 GFLOPs / 146 TOPS effective): 2.4 ms", C.TEXT_DIM),
        ("• Actual GPU time is 42× longer than the compute-only floor", C.ACCENT_ORANGE, True),
        ("• 98% of GPU time is spent waiting on memory, not computing", C.TEXT_DIM),
        "",
        ("WHY — transformer activations stream through VRAM every layer", C.ACCENT_BLUE, True),
        "• 840M params  •  3.71 GB peak activations  •  ~147 GB total memory traffic per frame",
        ("• Arithmetic intensity ~2 FLOPs/byte (ridge point on 5090: 117 FLOPs/byte)", C.TEXT_DIM),
        ("• Even the 5090's 72 MB L2 cache can't absorb the activation stream", C.ACCENT_ORANGE),
        "",
        ("EDGE PROJECTION — Edge MPU Target (200 TOPS, 134.4 GB/s, ~4 MB SRAM)", C.ACCENT_BLUE, True),
        "• Bandwidth ratio: 1523 / 101 = 15.1× less bandwidth than RTX 5090",
        ("• Projected ~2,400 ms per frame (0.4 FPS) — a 14× slowdown",
         C.NOT_FEASIBLE, True),
        ("• 200 TOPS is irrelevant here — compute is only 2% of total time", C.TEXT_DIM),
        ("• Memory capacity: 7.07 GB peak vs 8 GB total DRAM — no headroom", C.ACCENT_ORANGE),
    ], font_size=12)


def slide_bandwidth_requirements(prs: Presentation):
    """Slide: Required bandwidth for target framerates."""
    slide = new_slide(prs)
    add_title_subtitle(slide, "Required Memory Bandwidth for Real-Time SAM 3",
                       "~147 GB memory traffic per frame at 1080p, 9 concept prompts")

    headers = ["Target FPS", "Time Budget", "Required BW (eff)", "Required BW (raw)",
               "Memory Tech", "Feasible at 25W?"]
    rows = [
        ["1 FPS",  "1000 ms", "147 GB/s",   "196 GB/s",   "256-bit LPDDR5X",      "Possible"],
        ["5 FPS",  "200 ms",  "735 GB/s",   "980 GB/s",   "HBM2e or 512-bit",     "Difficult"],
        ["10 FPS", "100 ms",  "1,470 GB/s", "1,960 GB/s", "HBM3 (desktop-class)", "No"],
        ["24 FPS", "42 ms",   "3,528 GB/s", "4,704 GB/s", "Beyond HBM3",          "No"],
        ["30 FPS", "33 ms",   "4,414 GB/s", "5,885 GB/s", "Multi-die HBM3e",      "No"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                     Inches(CONTENT_W), Inches(2.4), headers, rows)

    add_bullet_box(slide, CONTENT_LEFT, 4.2, CONTENT_W, 2.7, [
        ("Current hardware reference", C.ACCENT_BLUE, True),
        ("• Edge MPU Target:  134.4 GB/s (128-bit LPDDR5X)  →  0.4 FPS", C.NOT_FEASIBLE),
        ("• RTX 5090:        1,792 GB/s (512-bit GDDR7)    →  6.0 FPS", C.TEXT_BRIGHT),
        ("• NVIDIA H200:     4,800 GB/s (HBM3e)            →  ~33 FPS (matches Meta's 30 ms report)", C.ACCENT_GREEN),
        "",
        ("Conclusion — real-time SAM 3 requires HBM-class bandwidth.", C.ACCENT_ORANGE, True),
        ("For edge at 25 W, the model must change — not the hardware.", C.TEXT_WHITE, True),
    ], font_size=12)


def slide_prompt_scaling(prs: Presentation):
    """Slide: How concept prompt count affects performance."""
    slide = new_slide(prs)
    add_title_subtitle(slide, "Prompt Count Scaling — Decoder Cost Is Linear",
                       "Vision encoder is fixed cost (~70 ms); each concept prompt adds ~6 ms in decoder")

    # Measured data table
    headers = ["Concepts", "RTX 5090", "FPS", "Edge Projected", "Edge FPS", "Example Prompts"]
    rows = [
        ["1", "72ms", "13.8", "~1,068ms", "0.9", "person"],
        ["3", "90ms", "11.1", "~1,333ms", "0.7", "person, vehicle, dog"],
        ["9 (current)", "121ms", "8.3", "~1,791ms", "0.6", "person, vehicle, car, truck, ..."],
        ["18", "177ms", "5.7", "~2,618ms", "0.4", "full concept set + accessories"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                     Inches(CONTENT_W), Inches(1.5),
                     headers, rows)

    # Cost breakdown chart
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 3), facecolor=MPL_COLORS["bg_slide"])

    # Left: Latency vs prompt count
    ax1.set_facecolor(MPL_COLORS["bg_slide"])
    prompts = [1, 3, 9, 18]
    latencies = [72, 90, 121, 177]
    ax1.plot(prompts, latencies, "o-", color=MPL_COLORS["blue"], linewidth=2, markersize=8)
    ax1.axhline(y=72, color=MPL_COLORS["orange"], linestyle="--", alpha=0.7, label="Vision encoder floor (72ms)")
    ax1.fill_between(prompts, 72, latencies, alpha=0.2, color=MPL_COLORS["green"], label="Decoder cost")
    ax1.set_xlabel("Number of Concept Prompts", color=MPL_COLORS["text"], fontsize=10)
    ax1.set_ylabel("Latency (ms)", color=MPL_COLORS["text"], fontsize=10)
    ax1.set_title("Latency vs Prompt Count (RTX 5090)", color=MPL_COLORS["text"], fontsize=11, fontweight="bold")
    ax1.legend(fontsize=8, facecolor=MPL_COLORS["bg"], edgecolor=MPL_COLORS["grid"],
               labelcolor=MPL_COLORS["text"])
    ax1.tick_params(colors=MPL_COLORS["dim"], labelsize=9)
    ax1.spines[:].set_color(MPL_COLORS["grid"])
    ax1.grid(True, alpha=0.2, color=MPL_COLORS["grid"])

    # Right: Stacked bar showing encoder vs decoder split
    ax2.set_facecolor(MPL_COLORS["bg_slide"])
    encoder_ms = [70, 70, 70, 70]
    decoder_ms = [2, 20, 51, 107]
    labels = ["1", "3", "9", "18"]
    x = np.arange(len(labels))
    ax2.bar(x, encoder_ms, 0.5, label="Vision Encoder (fixed)", color=MPL_COLORS["orange"])
    ax2.bar(x, decoder_ms, 0.5, bottom=encoder_ms, label="Text + DETR Decoder", color=MPL_COLORS["green"])
    ax2.set_xticks(x)
    ax2.set_xticklabels(labels)
    ax2.set_xlabel("Prompts", color=MPL_COLORS["text"], fontsize=10)
    ax2.set_ylabel("Latency (ms)", color=MPL_COLORS["text"], fontsize=10)
    ax2.set_title("Encoder vs Decoder Split", color=MPL_COLORS["text"], fontsize=11, fontweight="bold")
    ax2.legend(fontsize=8, facecolor=MPL_COLORS["bg"], edgecolor=MPL_COLORS["grid"],
               labelcolor=MPL_COLORS["text"])
    ax2.tick_params(colors=MPL_COLORS["dim"], labelsize=9)
    ax2.spines[:].set_color(MPL_COLORS["grid"])

    fig.tight_layout(pad=1.5)
    img_stream = fig_to_image_stream(fig)
    # Shrink chart to 10" wide centered so the takeaway fits under it
    slide.shapes.add_picture(img_stream, Inches(1.7), Inches(3.2),
                             width=Inches(10.0))

    add_bullet_box(slide, CONTENT_LEFT, 6.35, CONTENT_W, 0.6, [
        ("Takeaway — vision encoder (~70 ms) is the hard floor. Even 1 prompt = 1,068 ms on edge; "
         "prompt tuning helps on desktop (14 FPS) but can't close the edge bandwidth gap.",
         C.ACCENT_ORANGE, True),
    ], font_size=11)


def slide_quantization_tested(prs: Presentation):
    """Slide: Weight-only INT8 quantization results — doesn't help."""
    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "Quantization Tested — Weight-Only INT8 Doesn't Help",
                       "Measured with torchao Int8WeightOnlyConfig on RTX 5090 (9 concepts, 720p)")

    headers = ["Metric", "BF16 (baseline)", "INT8 Weight-Only", "Delta"]
    rows = [
        ["Wall clock",      "121 ms",            "121 ms",            "0% (no change)"],
        ["GPU kernel time", "102 ms",            "117 ms",            "15% SLOWER"],
        ["Peak VRAM",       "7.07 GB",           "5.11 GB",           "2 GB saved"],
        ["Edge projection", "1,791 ms (0.6 FPS)", "1,731 ms (0.6 FPS)", "Negligible"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                     Inches(CONTENT_W), Inches(1.5), headers, rows)

    add_bullet_box(slide, CONTENT_LEFT, 3.2, CONTENT_W, 3.5, [
        ("Why it doesn't help", C.ACCENT_BLUE, True),
        "• Weight-only quantization shrinks model params (3.36 GB → 704 MB)",
        ("• But activations stay in BF16 — and activations are 98% of bandwidth traffic", C.ACCENT_ORANGE, True),
        ("• Dequantization overhead (INT8 → BF16 per matmul) adds latency", C.TEXT_DIM),
        ("• Lost Meta's fused addmm_act kernel → unfused path is slower", C.TEXT_DIM),
        "",
        ("What WOULD help — activation quantization (INT8 or FP8 activations)", C.ACCENT_GREEN, True),
        "• Would halve the dominant memory traffic between layers",
        ("• Edge projection: ~1,700 ms → ~850 ms (1.2 FPS) — still not real-time", C.ACCENT_ORANGE),
    ], font_size=12)


def slide_activation_quant_challenges(prs: Presentation):
    """Slide: Why activation quantization is hard for SAM 3."""
    slide = new_slide(prs)
    add_title_subtitle(slide, "Activation Quantization — Why It's Hard for SAM 3",
                       "The one lever that could halve edge latency, but requires research-grade effort")

    headers = ["Challenge", "Impact", "Mitigation"]
    rows = [
        ["Attention score clipping",
         "INT8 clips outlier scores that encode 'attend strongly to this location'",
         "FP8 (E4M3) preserves dynamic range"],
        ["Text-vision cross-attention",
         "Quant errors → false positives, missed detections, concept misclassification",
         "Per-layer sensitivity analysis, mixed-precision"],
        ["Calibration data dependency",
         "Scale factors derived from cal data; mismatch → degraded accuracy",
         "Diverse calibration set matching deployment distribution"],
        ["Flash Attention 3 incompatibility",
         "No INT8 flash attention kernel; fallback = slower unfused attention",
         "FP8 flash attention (future), or accept the unfused penalty"],
        ["Not all layers are equal",
         "LayerNorm outputs, residuals, first/last layers are quantization-sensitive",
         "Mixed-precision: INT8 bulk matmuls, BF16 for sensitive layers"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                     Inches(CONTENT_W), Inches(2.6), headers, rows,
                     col_widths=[Inches(2.8), Inches(5.0), Inches(4.5)])

    add_bullet_box(slide, CONTENT_LEFT, 4.3, CONTENT_W, 2.6, [
        ("Even with perfect activation quantization", C.ACCENT_BLUE, True),
        ("• INT8 activations  →  ~2× traffic reduction  →  edge: ~850 ms (1.2 FPS)", C.TEXT_BRIGHT),
        ("• FP8 activations   →  ~2× traffic reduction  →  edge: ~850 ms (1.2 FPS)", C.TEXT_BRIGHT),
        ("• INT4 activations  →  ~4× traffic reduction  →  edge: ~425 ms (2.4 FPS) — significant accuracy risk", C.TEXT_BRIGHT),
        ("• None of these reach the 5 FPS (200 ms) budget on 134.4 GB/s LPDDR5X.", C.NOT_FEASIBLE, True),
        "",
        ("Viable paths today", C.ACCENT_GREEN, True),
        ("• SmoothQuant — shifts quant difficulty from activations to weights (proven on LLMs)", C.ACCENT_GREEN),
        ("• FP8 (E4M3) — RTX 5090 supports natively, best accuracy/speed tradeoff", C.ACCENT_GREEN),
        ("• Wait for Meta — an official quantized SAM 3 checkpoint would bypass all these issues", C.ACCENT_GREEN),
    ], font_size=11)


def slide_resolution_lock(prs: Presentation):
    """Slide: Why reducing input resolution doesn't help."""
    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "Resolution Is Locked — Input Size Doesn't Matter",
                       "SAM 3 internally processes 1008x1008 regardless of input resolution")

    headers = ["Input Resolution", "Internal Res", "ViT Tokens", "Avg Latency (5090)", "Detections/frame"]
    rows = [
        ["4K (3840x2160)",    "1008x1008", "3,969", "196 ms", "37.2"],
        ["1080p (1920x1080)", "1008x1008", "3,969", "139 ms", "35.0"],
        ["720p (1280x720)",   "1008x1008", "3,969", "117 ms", "32.1"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                     Inches(CONTENT_W), Inches(1.3), headers, rows)

    add_bullet_box(slide, CONTENT_LEFT, 3.0, CONTENT_W, 3.9, [
        ("Why — Rotary Position Embeddings (RoPE) are resolution-locked", C.ACCENT_BLUE, True),
        "• The ViT uses 2D RoPE pre-computed for a 63x63 token grid (1008/16 = 63)",
        ("• Feeding a different resolution → shape mismatch → assertion failure", C.TEXT_DIM),
        ("• rope_interp support exists but requires model rebuild + retraining", C.TEXT_DIM),
        "",
        ("What this means", C.ACCENT_BLUE, True),
        "• The 16 ms savings from 4K → 720p is just pre/post-processing overhead",
        ("• Model compute + memory traffic are IDENTICAL at every input resolution", C.ACCENT_ORANGE, True),
        ("• Token count (3,969), FLOPs, and activation memory are all fixed", C.TEXT_DIM),
        ("• Only detection accuracy changes (fewer small objects at 720p)", C.TEXT_DIM),
        "",
        ("Implication — resolution reduction is NOT a viable optimization lever.", C.NOT_FEASIBLE, True),
        ("Remaining options: quantization, fewer params, or a different model.", C.TEXT_WHITE),
    ], font_size=12)


def slide_model_comparison(prs: Presentation, data_dir: Path):
    """Slide: Speed vs accuracy tradeoff across model variants."""
    comp_path = data_dir / "model_comparison.json"
    if not comp_path.exists():
        return

    with open(comp_path) as f:
        comp = json.load(f)

    slide = new_slide(prs)
    add_title_subtitle(slide, "Speed vs Accuracy — The Model Tradeoff",
                       "Same 5 frames (720p embedded world), same hardware (RTX 5090)")

    profiles = comp.get("profiles", {})
    headers = ["Model", "Params", "5090 ms", "5090 FPS", "VRAM",
               "Dets/frm", "Recall", "Edge ms", "Edge FPS", "Concepts"]
    rows = []
    for key in ["sam3", "fastsam_x", "fastsam_s", "yolo11x"]:
        if key not in profiles:
            continue
        p = profiles[key]
        rows.append([
            p["name"], f"{p['param_count_m']:.0f}M",
            f"{p['avg_inference_ms']:.0f}", f"{1000/p['avg_inference_ms']:.0f}",
            f"{p['peak_vram_gb']:.1f}GB", f"{p['avg_detections']:.0f}",
            f"{p['recall_vs_sam3']:.0%}",
            f"{p['edge_projected_ms']:.0f}", f"{p['edge_fps']:.1f}",
            p["concept_vocabulary"][:15],
        ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                     Inches(CONTENT_W), Inches(1.8), headers, rows)

    add_bullet_box(slide, CONTENT_LEFT, 3.4, CONTENT_W, 3.5, [
        ("The tradeoff isn't speed — it's capability", C.ACCENT_BLUE, True),
        "",
        ("• SAM 3 (840M)  —  4M+ concept vocab, text-prompted, segmentation masks", C.ACCENT_ORANGE),
        ("    e.g. 'person wearing red jacket carrying backpack near delivery truck'", C.TEXT_DIM),
        ("    0.5 FPS on edge (too slow), but irreplaceable concept richness", C.TEXT_DIM),
        ("• FastSAM (11-68M)  —  YOLO-based segmentation, 80 COCO classes only", C.ACCENT_GREEN),
        ("    'person', 'car', 'dog' — no attributes, no relationships", C.TEXT_DIM),
        ("    3-4 FPS on edge, but 0% recall vs SAM 3 (different job entirely)", C.TEXT_DIM),
        ("• YOLO 11x (57M)  —  Detection boxes only, no masks, 80 classes", C.TEXT_BRIGHT),
        ("    5 FPS on edge, real-time, but no segmentation or concepts", C.TEXT_DIM),
        "",
        ("Conclusion — no smaller model replaces SAM 3's concept understanding.", C.NOT_FEASIBLE, True),
        ("The real question: how much concept richness do you actually need?", C.ACCENT_PURPLE, True),
    ], font_size=11)


def slide_hybrid_v2(prs: Presentation):
    """Slide: Hybrid V2 results — the breakthrough."""
    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "Hybrid V2 — YOLO-Seg + CLIP  —  The Breakthrough",
                       "Eliminate MobileSAM entirely. Two models, 33x faster than SAM 3 on edge.")
    add_pipeline_strip(slide, ["FFmpeg", ("YOLO-seg", True), ("CLIP", True),
                                "SQLite", "NLQ / LLM"])

    headers = ["Pipeline", "5090 ms", "5090 FPS", "Params", "Edge ms", "Edge FPS", "vs SAM 3"]
    rows = [
        ["SAM 3 single-pass",          "121", "8.3",  "840M", "~1,700", "0.6",  "baseline"],
        ["Hybrid V1 (YOLO+SAM+CLIP)",  "142", "7.0",  "218M", "~200",   "5.0",  "8.5x"],
        ["V2 medium-seg + CLIP",       "50",  "20",   "174M", "~65",    "15",   "26x"],
        ["V2 small-seg + CLIP",        "44",  "23",   "161M", "~58",    "17",   "29x"],
        ["V2 nano-seg + CLIP",         "39",  "26",   "155M", "~51",    "20",   "33x"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                     Inches(CONTENT_W), Inches(2.0), headers, rows)

    add_bullet_box(slide, CONTENT_LEFT, 4.2, CONTENT_W, 2.7, [
        ("Why it works", C.ACCENT_BLUE, True),
        "• YOLO-seg does detection + segmentation in ONE pass (3-8 ms)",
        ("• Eliminates MobileSAM's 95 ms image encoder entirely", C.ACCENT_GREEN),
        "• CLIP batched classification: text features cached, crops batched (25-34 ms)",
        ("• Total model params: 155-174M vs SAM 3's 840M — 5× less memory traffic", C.TEXT_DIM),
        "",
        ("Edge feasibility", C.ACCENT_BLUE, True),
        ("• Nano-seg + CLIP: ~51 ms/frame = 20 FPS on 134.4 GB/s LPDDR5X", C.ACCENT_GREEN, True),
        "• Fits comfortably in 8 GB DRAM with headroom for OS + runtime",
        ("• Models small enough to benefit from on-chip SRAM caching", C.TEXT_DIM),
        "",
        ("Visual quality — confirmed indistinguishable from SAM 3", C.ACCENT_PURPLE, True),
        ("• Trade-off: 80 COCO classes + CLIP open-vocab vs SAM 3's 4M+ native concepts", C.TEXT_DIM),
    ], font_size=11)


def _load_bakeoff_data():
    """Load bake-off results and edge projections, if present."""
    base = Path("data/output/bakeoff")
    clips = {
        "720p":  "720p_EW_clip",
        "1080p": "embedded_world_clip_1080p",
        "4K":    "embedded_world_clip",
    }
    contestants = ["mobilesam", "efficientsam_tiny", "efficientsam_small", "yolo_seg"]

    per_clip: dict[str, dict[str, dict]] = {}
    for res, stem in clips.items():
        summary_path = base / stem / "summary.json"
        if not summary_path.exists():
            return None, None  # bake-off not yet run
        per_clip[res] = json.loads(summary_path.read_text())["contestants"]

    # Compute per-frame FPS on 5090 from the detailed result files
    per_frame_fps: dict[str, dict[str, float]] = {res: {} for res in clips}
    for res, stem in clips.items():
        for name in contestants:
            data = json.loads((base / stem / "results" / f"{name}.json").read_text())
            frame_ms = [fr["latency_ms"] for fr in data["frames"] if fr["n_boxes"] > 0]
            per_frame_fps[res][name] = 1000.0 / np.mean(frame_ms) if frame_ms else 0.0

    edge_path = base / "edge_projection.json"
    edge = json.loads(edge_path.read_text()) if edge_path.exists() else None

    return {"per_clip": per_clip, "per_frame_fps": per_frame_fps,
            "contestants": contestants, "clips": list(clips.keys())}, edge


def slide_bakeoff_summary(prs: Presentation):
    """Slide: mask-model bake-off — headline numbers."""
    data, edge = _load_bakeoff_data()
    if data is None:
        return

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "Mask Model Bake-Off — Headline Numbers",
                       "MobileSAM vs EfficientSAM-Tiny/Small vs YOLO-seg, scored against SAM 3 references")
    add_pipeline_strip(slide, ["FFmpeg", "YOLO 11x", ("Mask bake-off", True),
                                "SQLite", "NLQ / LLM"])

    display_name = {
        "mobilesam":          "MobileSAM (vit_t)",
        "efficientsam_tiny":  "EfficientSAM-Tiny",
        "efficientsam_small": "EfficientSAM-Small",
        "yolo_seg":           "YOLO-seg (yolo11s-seg)",
    }

    # Table 1: quality + params
    headers = ["Model", "Params", "VRAM @1080p", "IoU @720p", "IoU @1080p", "IoU @4K"]
    rows = []
    for name in data["contestants"]:
        c720 = data["per_clip"]["720p"][name]
        c1080 = data["per_clip"]["1080p"][name]
        c4k = data["per_clip"]["4K"][name]
        rows.append([
            display_name[name],
            f"{c720['params_m']:.1f}M",
            f"{c1080['peak_vram_mb']:.0f} MB",
            f"{c720['mean_iou']:.3f}",
            f"{c1080['mean_iou']:.3f}",
            f"{c4k['mean_iou']:.3f}",
        ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                     Inches(CONTENT_W), Inches(1.5), headers, rows)

    # Table 2: FPS (5090 measured + edge projection)
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(3.6), Inches(CONTENT_W), Inches(0.3),
                 "Full-frame FPS (~13 boxes/frame average)", font_size=12,
                 color=C.ACCENT_PURPLE, bold=True)
    headers2 = ["Model", "5090 @720p", "5090 @1080p", "5090 @4K",
                "Edge @720p", "Edge @1080p", "Edge @4K"]
    rows2 = []
    for name in data["contestants"]:
        edge_720 = edge["projections"]["720p"][name]["projected_fps_edge"] if edge else 0
        edge_1080 = edge["projections"]["1080p"][name]["projected_fps_edge"] if edge else 0
        edge_4k = edge["projections"]["4K"][name]["projected_fps_edge"] if edge else 0
        rows2.append([
            display_name[name],
            f"{data['per_frame_fps']['720p'][name]:.0f}",
            f"{data['per_frame_fps']['1080p'][name]:.0f}",
            f"{data['per_frame_fps']['4K'][name]:.0f}",
            f"{edge_720:.1f}", f"{edge_1080:.1f}", f"{edge_4k:.1f}",
        ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(3.9),
                     Inches(CONTENT_W), Inches(1.4), headers2, rows2)

    add_bullet_box(slide, CONTENT_LEFT, 5.55, CONTENT_W, 1.5, [
        ("Key findings", C.ACCENT_BLUE, True),
        ("• EfficientSAM-Tiny dominates MobileSAM — same params, ~2-8× faster, slightly higher IoU", C.ACCENT_GREEN),
        "• EfficientSAM-Small leads quality: 0.91 IoU vs 0.86 for Tiny and MobileSAM",
        "• YOLO-seg is 3-5× faster than any SAM variant; trades ~0.1 IoU for that speed",
        ("• MobileSAM is obsoleted — no reason to pick it over EfficientSAM-Tiny", C.ACCENT_ORANGE, True),
        ("• EfficientSAM is resolution-invariant per box; MobileSAM latency grows with image area", C.TEXT_DIM),
        ("• Edge projections use the emulator's 15%/85% compute/bandwidth split (conservative for conv-heavy YOLO)", C.TEXT_DIM),
    ], font_size=11)


def slide_bakeoff_visuals(prs: Presentation):
    """Slide: side-by-side mask comparison PNG."""
    vis_path = Path("data/output/bakeoff/visuals/720p_sidebyside.png")
    if not vis_path.exists():
        return

    slide = new_slide(prs)
    add_title_subtitle(slide, "Mask Model Bake-Off — Visual Comparison",
                       "One frame, 720p: YOLO prompts + SAM 3 reference + four contestants")
    add_pipeline_strip(slide, ["FFmpeg", "YOLO 11x", ("Mask bake-off", True),
                                "SQLite", "NLQ / LLM"])

    slide.shapes.add_picture(str(vis_path),
                             Inches(CONTENT_LEFT), Inches(1.9),
                             width=Inches(CONTENT_W))

    add_bullet_box(slide, CONTENT_LEFT, 6.55, CONTENT_W, 0.45, [
        ("EfficientSAM-Small edges closest to SAM 3; YOLO-seg masks are visibly coarser but fast.",
         C.ACCENT_PURPLE, True),
    ], font_size=13)


def slide_fp8_quantization(prs: Presentation):
    """Slide: FP8 activation quantization results on the bake-off winners."""
    fp8_path = Path("data/output/bakeoff/fp8_edge_projection.json")
    if not fp8_path.exists():
        return
    fp8 = json.loads(fp8_path.read_text())
    fp8_proj = fp8.get("fp8_projections", {})

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "FP8 Activation Quantization — Real-World Test",
                       "Does halving activation traffic actually work on our bake-off winners?")
    add_pipeline_strip(slide, ["FFmpeg", "YOLO 11x",
                                ("ES-Small / YOLO-seg  •  FP8 test", True),
                                "SQLite", "NLQ / LLM"])

    display_name = {
        "efficientsam_small": "EfficientSAM-Small",
        "yolo_seg":           "YOLO-seg (yolo11s-seg)",
    }
    headers = ["Model", "Res", "FP8 applied?", "IoU bf16", "IoU FP8", "Δ IoU",
               "Edge FPS bf16", "Edge FPS FP8"]
    rows = []
    for res in ["720p", "1080p", "4K"]:
        for name in ["efficientsam_small", "yolo_seg"]:
            p = fp8_proj.get(res, {}).get(name)
            if not p:
                continue
            applied = p.get("fp8_actually_applied", False)
            applied_text = f"YES ({p['n_fp8_weights_swapped']} Linears)" if applied else "NO (Conv2d not supported)"
            bw_bf16 = p["bandwidth_limited_ms_bf16"]
            comp = p["compute_limited_ms"]
            edge_bf16_fps = 1000.0 / (comp + bw_bf16) if (comp + bw_bf16) > 0 else 0
            rows.append([
                display_name[name], res, applied_text,
                f"{p['mean_iou_bf16']:.3f}", f"{p['mean_iou_fp8']:.3f}", f"{p['iou_delta']:+.3f}",
                f"{edge_bf16_fps:.1f}", f"{p['projected_fps_edge_fp8']:.1f}",
            ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                     Inches(CONTENT_W), Inches(2.4), headers, rows)

    add_bullet_box(slide, CONTENT_LEFT, 4.5, CONTENT_W, 2.5, [
        ("Key findings", C.ACCENT_BLUE, True),
        "• EfficientSAM-Small: 94 of 95 Linear layers quantized to E4M3 via torchao (PerTensor)",
        ("• Quality loss is negligible (< 0.003 IoU) — FP8 activations preserve mask quality", C.ACCENT_GREEN, True),
        ("• Edge FPS projection doubles: 2.5 → 4.9 FPS at 720p (halved activation traffic on LPDDR5X)", C.ACCENT_GREEN, True),
        "",
        ("YOLO-seg — blocked by tool maturity, not by the model", C.ACCENT_ORANGE, True),
        "• torchao 0.17's Float8DynamicActivationFloat8Weight only targets nn.Linear",
        ("• YOLO-seg is 100 Conv2d / 0 Linear — zero layers were actually quantized", C.TEXT_DIM),
        ("• Conv FP8 needs custom kernels or a transformer_engine-style rewrite", C.TEXT_DIM),
        "",
        ("Desktop caveat — desktop latency is NOT predictive of edge.", C.ACCENT_PURPLE, True),
        ("• RTX 5090 FP8 matmul is slower than bf16 for tiny models (torchao kernel overhead dominates).", C.TEXT_DIM),
        ("• Edge silicon's native FP8 MMA paths will realize the bandwidth win.", C.TEXT_DIM),
    ], font_size=11)


def slide_smoothquant(prs: Presentation):
    """Slide: SmoothQuant + plain INT8 bake-off on ES-Small + YOLO-seg."""
    sq_path = Path("data/output/bakeoff/smoothquant_edge_projection.json")
    if not sq_path.exists():
        return
    sq = json.loads(sq_path.read_text())
    proj = sq.get("projections", {})

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "INT8 + SmoothQuant — Activation Quant on the Winners",
                       "Does activation-safe INT8 preserve IoU, and does SmoothQuant smoothing help?")
    add_pipeline_strip(slide, ["FFmpeg", "YOLO 11x",
                                ("ES-Small / YOLO-seg  •  INT8 + SmoothQuant", True),
                                "SQLite", "NLQ / LLM"])

    display_name = {
        "efficientsam_small": "EfficientSAM-Small",
        "yolo_seg":           "YOLO-seg (yolo11s-seg)",
    }
    # Show only 720p results — pattern is identical at 1080p and 4K.
    headers = ["Model", "Recipe", "Applied?", "IoU bf16", "IoU quant", "Δ IoU",
               "Edge FPS bf16", "Edge FPS quant"]
    rows = []
    res = "720p"
    for name in ["efficientsam_small", "yolo_seg"]:
        for recipe in ["int8", "smoothquant"]:
            p = proj.get(res, {}).get(name, {}).get(recipe)
            if not p:
                # ES-Small SmoothQuant fails at CONVERT due to torchao 0.17 API gap.
                rows.append([display_name[name], recipe, "CONVERT FAILED",
                             "—", "—", "—", "—", "—"])
                continue
            applied = "YES" if p["recipe_applied"] else "NO (Conv-only)"
            comp = p["compute_limited_ms"]
            bw_bf16 = p["bandwidth_limited_ms_bf16"]
            edge_bf16 = 1000.0 / (comp + bw_bf16) if (comp + bw_bf16) > 0 else 0
            rows.append([
                display_name[name], recipe, applied,
                f"{p['mean_iou_bf16']:.3f}",
                f"{p['mean_iou_recipe']:.3f}",
                f"{p['iou_delta']:+.3f}",
                f"{edge_bf16:.1f}",
                f"{p['projected_fps_edge_recipe']:.1f}",
            ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                     Inches(CONTENT_W), Inches(2.0), headers, rows)
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(4.05), Inches(CONTENT_W), Inches(0.3),
                 "720p shown; 1080p and 4K yield the same pattern (ΔIoU -0.001 to -0.002, edge FPS doubles on ES-Small).",
                 font_size=10, color=C.TEXT_DIM)

    add_bullet_box(slide, CONTENT_LEFT, 4.45, CONTENT_W, 2.55, [
        ("Key findings", C.ACCENT_BLUE, True),
        "• Plain INT8 on ES-Small: 95/95 Linears quantized; ΔIoU -0.002 across resolutions",
        ("• Edge FPS projection matches FP8: 2.5 → 4.9 FPS at 720p (halved activation traffic)", C.ACCENT_GREEN, True),
        "• INT8 and FP8 deliver the same bandwidth win — pick based on edge silicon's native path",
        "",
        ("SmoothQuant blocked by torchao API maturity", C.ACCENT_ORANGE, True),
        "• CONVERT step asserts weight implements SupportsActivationPreScaling",
        ("• torchao 0.17 Int8DynamicActivationInt8Weight doesn't implement that protocol — open bug", C.TEXT_DIM),
        ("• Given ΔIoU is already < 0.002 with plain INT8, smoothing is unlikely to move the needle here", C.TEXT_DIM),
        "",
        ("YOLO-seg still blocked on Conv2d — 0 layers quantized by any recipe.", C.ACCENT_ORANGE, True),
    ], font_size=11)


def slide_hybrid_v2_bakeoff(prs: Presentation):
    """Slide: CLIP quantization on Hybrid V2 — YOLO-seg stays BF16, CLIP gets FP8/INT8."""
    path = Path("data/output/bakeoff/hybrid_v2_edge_projection.json")
    if not path.exists():
        return
    proj = json.loads(path.read_text()).get("projections", {})

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "Hybrid V2 Bake-off — CLIP Quantization on the Real-Time Path",
                       "YOLO-seg stays Conv-only / BF16; CLIP's 72 Linears are the quantizable surface")
    add_pipeline_strip(slide, ["FFmpeg", "YOLO-seg", ("CLIP FP8 / INT8", True),
                                "SQLite", "NLQ / LLM"])

    headers = ["Res", "Recipe", "CLIP Linears Q'd", "Top-1 agree", "Top-3 Jaccard",
               "5090 YOLO ms", "5090 CLIP ms", "Edge total ms", "Edge FPS"]
    rows = []
    highlight = []
    row_idx = 0
    for res in ["720p", "1080p", "4K"]:
        for recipe in ["bf16", "fp8", "int8"]:
            p = proj.get(res, {}).get(recipe)
            if not p:
                continue
            row_idx += 1
            applied = f"{p['n_quantized']} of {p['n_linear']}" if p["actually_applied"] else "—"
            top1 = f"{p.get('top1_agreement', 0):.3f}" if recipe != "bf16" else "1.000"
            top3 = f"{p.get('top3_jaccard', 0):.3f}" if recipe != "bf16" else "1.000"
            rows.append([
                res, recipe.upper(), applied, top1, top3,
                f"{p['mean_yolo_ms_5090']:.1f}",
                f"{p['mean_clip_ms_5090']:.1f}",
                f"{p['projected_total_ms_edge']:.0f}",
                f"{p['projected_fps_edge']:.1f}",
            ])
            # Highlight the 720p FP8 row as the winning config (best quality-preserving quantized path)
            if res == "720p" and recipe == "fp8":
                highlight.append(row_idx)

    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                     Inches(CONTENT_W), Inches(3.1), headers, rows,
                     highlight_rows=highlight)

    add_bullet_box(slide, CONTENT_LEFT, 5.2, CONTENT_W, 1.8, [
        ("Key findings", C.ACCENT_BLUE, True),
        "• CLIP — not YOLO — dominates per-frame cost. At 720p: YOLO 4.4 ms vs CLIP 22 ms (5× ratio)",
        ("• torchao swapped 48 of 72 CLIP Linears (67%) for both FP8 and INT8; remaining 24 are small projection layers skipped by torchao 0.17", C.TEXT_DIM),
        ("• Edge FPS projection (720p): 2.9 BF16 → 4.9 quantized (+69%) — same win for FP8 and INT8 since both halve CLIP activation bytes", C.ACCENT_GREEN, True),
        ("• FP8 preserves concept tags better than INT8: 86.8% vs 80.2% top-1 agreement at 720p (FP8's wider dynamic range wins on softmax rankings)", C.ACCENT_AMBER, True),
        "",
        ("Reality check on the earlier 20 FPS Hybrid V2 claim", C.ACCENT_ORANGE, True),
        ("• That number assumed YOLO-dominated cost. Measuring CLIP-every-frame on every detection puts the full pipeline at ~5 FPS edge even after quantization.", C.TEXT_DIM),
        ("• Path to actual real-time: debounce CLIP to keyframes only. YOLO stays every frame (~13 FPS); CLIP reruns every Nth frame for concept tags.", C.ACCENT_INDIGO, True),
        ("• YOLO-seg remains BF16 / unquantized — Conv-only, torchao blocked. Custom kernels or transformer_engine required.", C.TEXT_DIM),
    ], font_size=11)


def slide_yolo_conv_quant(prs: Presentation):
    """Slide: YOLO-seg Conv-INT8 via 1x1 swap — FP8 still blocked, INT8 works."""
    path = Path("data/output/bakeoff/yolo_conv_quant_edge_projection.json")
    if not path.exists():
        return
    data = json.loads(path.read_text())
    proj = data["projections"]

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "YOLO-seg Conv Quantization — Partial Unblock via 1×1 Swap",
                       "torchao swap_conv2d_1x1_to_linear captures half the layers (44% of conv weights)")
    add_pipeline_strip(slide, ["FFmpeg", ("YOLO-seg INT8 (torchao 1×1)", True),
                                "CLIP FP8", "SQLite", "NLQ / LLM"])

    headers = ["Res", "Recipe", "1×1 swapped", "Q'd", "% conv wts", "Box recall",
               "Match IoU", "5090 ms", "Edge BF16", "Edge Q", "Edge FPS"]
    rows = []
    highlight = []
    row_i = 0
    for res in ["720p", "1080p", "4K"]:
        for recipe in ["bf16", "int8_1x1_swap", "fp8_1x1_swap"]:
            p = proj.get(res, {}).get(recipe)
            if not p:
                continue
            row_i += 1
            if "error" in p:
                rows.append([
                    res, recipe, "—", "—", "—", "—", "—", "—", "—", "—",
                    "BLOCKED — see note",
                ])
                continue
            recall_str = f"{p.get('box_recall', 0):.3f}" if recipe != "bf16" else "1.000"
            miou_str = f"{p.get('mean_matched_iou', 0):.3f}" if recipe != "bf16" else "1.000"
            rows.append([
                res, recipe,
                str(p['n_swapped_linears']) if p['n_swapped_linears'] else "—",
                str(p['n_quantized']) if p['n_quantized'] else "—",
                f"{100*p['frac_conv_weights_quantized']:.1f}%" if p['n_quantized'] else "—",
                recall_str, miou_str,
                f"{p['mean_frame_ms_5090']:.1f}",
                f"{p['projected_ms_edge_bf16']:.1f}",
                f"{p['projected_ms_edge_recipe']:.1f}",
                f"{p['projected_fps_edge_recipe']:.1f}",
            ])
            if res == "720p" and recipe == "int8_1x1_swap":
                highlight.append(row_i)

    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                     Inches(CONTENT_W), Inches(3.3), headers, rows,
                     highlight_rows=highlight)

    add_bullet_box(slide, CONTENT_LEFT, 5.35, CONTENT_W, 1.7, [
        ("Key findings", C.ACCENT_BLUE, True),
        "• yolo11s-seg: 100 Conv2d, 0 Linear. Half the Convs are 1×1 (44% of conv weights); these can be swapped to equivalent Linears and quantized.",
        ("• INT8 path works — 49/50 swapped 1×1 Convs quantized, 96-98% box recall, matched IoU 0.986-0.988. Detections essentially unchanged.", C.ACCENT_GREEN, True),
        ("• Edge FPS (720p): 18.7 BF16 → 23.8 INT8 (+27%). Hybrid V2 + 1 Hz CLIP + YOLO-INT8 projects to ~20 FPS — real real-time.", C.ACCENT_GREEN, True),
        "",
        ("FP8 path still blocked — second tooling gap below the Linear-only one", C.ACCENT_ORANGE, True),
        ("• torchao 0.17 Float8 PerTensor version=2 asserts 'input_tensor must be 1x128 scaled' inside _float8_addmm_impl — YOLO activations don't satisfy it", C.TEXT_DIM),
        ("• PerRow fails ('Only bf16/fp16 high-precision output types supported'); version=1 is rejected at runtime", C.TEXT_DIM),
        ("• Full Conv-FP8 needs custom kernels, TensorRT INT8/FP8, or transformer_engine — not reachable without platform work", C.ACCENT_AMBER, True),
    ], font_size=11)


def slide_trt_yolo(prs: Presentation):
    """Slide: Proper TensorRT FP8 on YOLO-seg — the full unblock."""
    path = Path("data/output/bakeoff/trt_yolo_edge_projection.json")
    if not path.exists():
        return
    data = json.loads(path.read_text())
    proj = data["projections"]

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "TensorRT YOLO-seg — FP8 Unblocked, Real-Time Ceiling Broken",
                       "TRT 10.16 on RTX 5090 Blackwell: full-model Conv quantization that torchao couldn't reach")
    add_pipeline_strip(slide, ["FFmpeg", ("YOLO-seg FP8 (TRT)", True),
                                "CLIP FP8", "SQLite", "NLQ / LLM"])

    headers = ["Res", "Recipe", "5090 ms", "Dets", "Box recall", "Matched IoU",
               "Edge ms", "Edge FPS"]
    rows = []
    highlight = []
    row_i = 0
    for res in ["720p", "1080p", "4K"]:
        for recipe in ["fp16", "int8", "fp8"]:
            p = proj.get(res, {}).get(recipe)
            if not p or "error" in p:
                continue
            row_i += 1
            rows.append([
                res, recipe.upper(),
                f"{p['mean_frame_ms_5090']:.2f}",
                f"{p.get('n_matched', 0) + p.get('n_fp', 0)}" if recipe != "fp16" else "—",
                f"{p.get('box_recall', 0):.3f}" if recipe != "fp16" else "1.000",
                f"{p.get('mean_matched_iou', 0):.3f}" if recipe != "fp16" else "1.000",
                f"{p['projected_ms_edge']:.1f}",
                f"{p['projected_fps_edge']:.1f}",
            ])
            if res == "720p" and recipe == "fp8":
                highlight.append(row_i)
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                     Inches(CONTENT_W), Inches(3.1), headers, rows,
                     highlight_rows=highlight)

    add_bullet_box(slide, CONTENT_LEFT, 5.25, CONTENT_W, 1.75, [
        ("Key findings", C.ACCENT_BLUE, True),
        ("• FP8 on YOLO-seg WORKS via TRT 10.16 on Blackwell (SM 12.0). Earlier torchao block was a tool-chain gap, not a fundamental one.", C.ACCENT_GREEN, True),
        ("• FP8 quality essentially perfect — 100% box recall, matched IoU 0.998, indistinguishable from FP16.", C.ACCENT_GREEN, True),
        ("• Edge FPS 720p: 18.6 FP16 → 36.8 FP8 (+98%). Full-stack = Hybrid V2 + 1 Hz CLIP + YOLO-FP8 ≈ 36 FPS edge (prior target 20, nearly 2×).", C.ACCENT_INDIGO, True),
        ("• INT8 ships too but drops low-confidence boxes (87-92% recall). FP8's wider range wins on detection-head logits.", C.ACCENT_AMBER, True),
        ("Preprocessing: 640×640 letterbox runs on CPU, not GPU/NPU (not included in the ms/frame above).", C.ACCENT_ORANGE, True),
        ("• Measured on 5090 host (i9-14900KF, cv2.resize bilinear, 1 thread): 0.17 / 0.32 / 0.33 ms at 720p / 1080p / 4K — ~0.5–1% of one core at 30 fps. Flat across source res (output size dominates).",
         C.TEXT_BRIGHT),
        ("• Edge ARM (Cortex-A55 ≈ 10× slower single-thread) → ~2–3 ms/frame, ~6–10% of one edge core at 30 fps. SoCs with a fixed-function ISP / 2D GPU (Qualcomm, MediaTek, NXP, Ambarella, Hailo) move this off-CPU entirely; pure-NPU boards (Coral) pay the full cost.",
         C.TEXT_DIM),
    ], font_size=9)


def slide_yolov8n_comparison(prs: Presentation):
    """Slide: yolo11s-seg vs yolov8n-seg cross-variant comparison for silicon apples-to-apples.

    Consumes the yolov8n-seg bake-off JSONs produced by running the TRT YOLO
    and concurrency bake-offs with KEYHOLE_YOLO_VARIANT=yolov8n-seg.
    """
    trt8n_p = Path("data/output/bakeoff/trt_yolo_yolov8n-seg_summary.json")
    trt11_p = Path("data/output/bakeoff/trt_yolo_summary.json")
    conc8n_p = Path("data/output/bakeoff/concurrency_yolov8n-seg_edge_projection.json")
    conc11_p = Path("data/output/bakeoff/concurrency_edge_projection.json")
    if not (trt8n_p.exists() and trt11_p.exists()):
        return

    trt8 = json.loads(trt8n_p.read_text())
    trt11 = json.loads(trt11_p.read_text())

    def ms_5090(d, res, recipe):
        r = d.get("results", {}).get(res, {}).get(recipe, {})
        return r.get("mean_frame_ms", 0) if isinstance(r, dict) and "frames" in r else 0

    def quality(d, res, recipe):
        """Return (box_recall, matched_iou) at given res/recipe; (nan, nan) if missing."""
        q = d.get("quality", {}).get(res, {}).get(recipe, {})
        return q.get("box_recall", float("nan")), q.get("mean_matched_iou", float("nan"))

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        "yolo11s-seg vs yolov8n-seg — cross-variant comparison for silicon benchmarking",
        "Same bake-off infrastructure, two YOLO generations. Run both with KEYHOLE_YOLO_VARIANT=… to compare.",
    )
    add_pipeline_strip(
        slide,
        ["FFmpeg", ("YOLO-seg FP8 (TRT)", True), "CLIP FP8 @ 1 Hz", "SQLite", "NLQ / LLM"],
        accent_color=C.ACCENT_INDIGO,
    )

    # Table 1: 5090 engine ms × precision, with 720p box-recall to flag INT8 quality story
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(1.9), Inches(CONTENT_W), Inches(0.3),
                 "Pure TRT engine execute() ms on 5090 Blackwell, + 720p box recall (vs FP16 baseline)",
                 font_size=12, color=C.ACCENT_PURPLE, bold=True)
    headers = ["Variant", "Prec", "Params", "720p ms", "1080p ms", "4K ms",
               "720p recall", "720p IoU", "Verdict"]
    rows = []
    highlight = []
    row_i = 0
    for label, data, params, mAP in [
        ("yolo11s-seg (mAP 37.0)", trt11, "10.1 M", "37.0"),
        ("yolov8n-seg (mAP 30.5)",  trt8,  "3.4 M",  "30.5"),
    ]:
        for recipe in ("fp16", "int8", "fp8"):
            row_i += 1
            m_720 = ms_5090(data, "720p", recipe)
            m_1080 = ms_5090(data, "1080p", recipe)
            m_4k = ms_5090(data, "4K", recipe)
            rec, iou = quality(data, "720p", recipe)
            rec_s = f"{rec:.3f}" if rec == rec else "—"   # NaN check
            iou_s = f"{iou:.3f}" if iou == iou else "—"
            fp16_ms = ms_5090(data, "720p", "fp16")
            verdict = ""
            if recipe == "fp16":
                verdict = "reference"
            elif recipe == "int8":
                # Current INT8 numbers reflect the better-calibration PTQ
                # (Ultralytics coco128-seg dataset, ~128 images). If recall
                # is ≥ 0.90, flag the speed cliff but not a quality cliff.
                slow = m_720 > fp16_ms * 1.05
                lossy = rec == rec and rec < 0.90
                if slow and lossy:
                    verdict = "slower + lossy"
                    highlight.append(row_i)
                elif slow:
                    verdict = "slower, ok quality"
                elif lossy:
                    verdict = "lossy"
                else:
                    verdict = "ok"
            elif recipe == "fp8":
                if rec == rec and rec >= 0.99:
                    verdict = "shipping ✓"
                else:
                    verdict = "check recall"
            rows.append([
                label if recipe == "fp16" else "",
                recipe.upper(),
                params if recipe == "fp16" else "",
                f"{m_720:.2f}" if m_720 > 0 else "—",
                f"{m_1080:.2f}" if m_1080 > 0 else "—",
                f"{m_4k:.2f}" if m_4k > 0 else "—",
                rec_s, iou_s,
                verdict,
            ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(2.22),
                     Inches(CONTENT_W), Inches(1.9), headers, rows,
                     highlight_rows=highlight, font_size=10, header_font_size=10)

    # Table 2: concurrency — edge batch ms side by side (both at FP8 shipping)
    if conc8n_p.exists() and conc11_p.exists():
        add_text_box(slide, Inches(CONTENT_LEFT), Inches(4.25), Inches(CONTENT_W), Inches(0.3),
                     "Multi-stream batching (FP8 shipping) — edge batch ms on NPU Mid",
                     font_size=12, color=C.ACCENT_PURPLE, bold=True)
        c8 = json.loads(conc8n_p.read_text())["batches_edge"]
        c11 = json.loads(conc11_p.read_text())["batches_edge"]
        c8_map = {r["batch"]: r["mean_ms_edge"] for r in c8}
        c11_map = {r["batch"]: r["mean_ms_edge"] for r in c11}
        headers2 = ["Batch", "11s ms", "v8n ms", "Speedup",
                    "11s FPS/stream", "v8n FPS/stream"]
        rows2 = []
        for B in (1, 2, 4, 8, 16):
            m11 = c11_map.get(B, 0)
            m8 = c8_map.get(B, 0)
            fps11 = 1000.0 / m11 if m11 > 0 else 0
            fps8 = 1000.0 / m8 if m8 > 0 else 0
            speedup = m11 / m8 if m8 > 0 else 0
            rows2.append([
                f"B = {B}",
                f"{m11:.1f}" if m11 > 0 else "—",
                f"{m8:.1f}" if m8 > 0 else "—",
                f"{speedup:.2f}×" if speedup > 0 else "—",
                f"{fps11:.1f}" if fps11 > 0 else "—",
                f"{fps8:.1f}" if fps8 > 0 else "—",
            ])
        add_styled_table(slide, Inches(CONTENT_LEFT), Inches(4.55),
                         Inches(CONTENT_W), Inches(1.4), headers2, rows2,
                         highlight_rows=[1, 3], font_size=10)

    add_bullet_box(slide, CONTENT_LEFT, 6.05, CONTENT_W, 1.05, [
        ("Silicon comparison + the INT8 calibration story", C.ACCENT_BLUE, True),
        ("• Vendor NPU benchmarks almost always cite yolov8n-seg — nano is the de-facto industry reference. Alongside yolo11s-seg this lets you drop real-silicon numbers into a direct apples-to-apples compare.",
         C.TEXT_BRIGHT),
        ("• INT8 speed cliff (~22% slower than FP16) is STRUCTURAL — kernel-launch overhead > FP8/INT8 compute savings on 3.4M-param nano. Can't be fixed with more calibration or QAT. FP8 sidesteps it entirely.",
         C.ACCENT_AMBER),
        ("• INT8 quality DEPENDS ON CALIBRATION DATASET: 20-frame cache → 0.714 recall. Ultralytics' coco128-seg PTQ → 0.912. Full-COCO PTQ or true QAT → ~0.95-0.98. Vendor INT8 claims are only credible with disclosed calibration methodology.",
         C.ACCENT_GREEN),
    ], font_size=9)


def slide_trt_clip(prs: Presentation):
    """Slide: TRT-compile CLIP visual — FP8 halves the CLIP edge cost."""
    path = Path("data/output/bakeoff/trt_clip_edge_projection.json")
    if not path.exists():
        return
    data = json.loads(path.read_text())
    proj = data["projections"]

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "TensorRT CLIP — Visual Tower Compiled, FP8 Halves Edge Cost",
                       "Completes the Hybrid V2 TRT stack: YOLO-seg FP8 + CLIP ViT-B-32 FP8")
    add_pipeline_strip(slide, ["FFmpeg", "YOLO-seg FP8 (TRT)",
                                ("CLIP FP8 (TRT)", True), "SQLite", "NLQ / LLM"])

    headers = ["Res", "Recipe", "5090 ms", "Top-1 vs BF16",
               "Edge CLIP ms", "CLIP-only FPS"]
    rows = []
    highlight = []
    row_i = 0
    for res in ["720p", "1080p", "4K"]:
        for recipe in ["bf16_torch", "fp16", "fp8"]:
            p = proj.get(res, {}).get(recipe)
            if not p or "error" in p:
                continue
            row_i += 1
            name = {"bf16_torch": "BF16 torch", "fp16": "TRT FP16", "fp8": "TRT FP8"}[recipe]
            rows.append([
                res, name,
                f"{p['mean_frame_ms_5090']:.2f}",
                f"{p['top1_agreement']:.3f}",
                f"{p['projected_clip_ms_edge']:.1f}",
                f"{p['projected_fps_edge_clip_only']:.1f}",
            ])
            if res == "720p" and recipe == "fp8":
                highlight.append(row_i)

    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                     Inches(CONTENT_W), Inches(3.1), headers, rows,
                     highlight_rows=highlight)

    add_bullet_box(slide, CONTENT_LEFT, 5.25, CONTENT_W, 1.85, [
        ("Key findings", C.ACCENT_BLUE, True),
        ("• CLIP visual TRT-compiled cleanly at FP16 + FP8 (180 MB engine). No QDQ nodes needed — TRT auto-selects FP8 layers.", C.ACCENT_GREEN, True),
        ("• FP8 edge CLIP drops from 29.8 ms (BF16/FP16) to 15.1 ms (+120% CLIP-only FPS).", C.ACCENT_GREEN, True),
        ("• Top-1 concept-tag agreement: TRT FP16 0.970, TRT FP8 0.964 — FP8 costs ~0.4 pts, noise-level.", C.TEXT_DIM),
        ("• Earlier hybrid_v2 CLIP measurement (22 ms) included per-crop Python dispatch. Pure visual() kernel time is 2.3 ms BF16 — TRT exposes the honest number.", C.ACCENT_AMBER, True),
        "",
        ("Recalibrated full-stack (720p Edge MPU)", C.ACCENT_INDIGO, True),
        ("• YOLO-FP8 (27.2 ms) + CLIP-FP8 every frame (15.1 ms) = 42.3 ms → 24 FPS  — real-time with ZERO debouncing",),
        ("• YOLO-FP8 + CLIP-FP8 at 1 Hz (0.5 ms amortized) = 27.7 ms → 36 FPS  — at the YOLO-only ceiling; CLIP is effectively free",),
    ], font_size=11)


def slide_concurrency(prs: Presentation):
    """Slide: multi-stream concurrency — YOLO batching + deployment recipes."""
    path = Path("data/output/bakeoff/concurrency_edge_projection.json")
    if not path.exists():
        return
    data = json.loads(path.read_text())

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "Multi-Stream Concurrency — One NPU, Many Cameras",
                       "Naive math says 36 FPS ÷ N streams. YOLO batching breaks that assumption.")
    add_pipeline_strip(slide, ["FFmpeg (×N)", ("YOLO-seg FP8 batch=N", True),
                                "CLIP FP8 @ 1 Hz", "SQLite", "NLQ / LLM"])

    # Table 1: YOLO batching characterization
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(1.9), Inches(CONTENT_W), Inches(0.3),
                 "YOLO-seg FP8 batched inference (same dynamic-batch TRT engine)",
                 font_size=12, color=C.ACCENT_PURPLE, bold=True)
    bt_headers = ["Batch", "5090 ms/batch", "5090 ms/stream",
                  "Edge ms/batch", "Edge ms/stream", "Edge FPS/stream"]
    bt_rows = []
    for r_5090, r_edge in zip(data["batches_5090"], data["batches_edge"]):
        bt_rows.append([
            f"B = {r_5090['batch']}",
            f"{r_5090['mean_ms']:.2f} ms",
            f"{r_5090['per_stream_ms']:.2f} ms",
            f"{r_edge['mean_ms_edge']:.1f} ms",
            f"{r_edge['per_stream_ms_edge']:.1f} ms",
            f"{1000 / r_edge['mean_ms_edge']:.1f}",
        ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(2.25),
                     Inches(CONTENT_W), Inches(1.55), bt_headers, bt_rows,
                     font_size=10)

    # Table 2: Deployment scenarios
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(4.0), Inches(CONTENT_W), Inches(0.3),
                 "Deployment recipes — N streams through one Edge MPU",
                 font_size=12, color=C.ACCENT_PURPLE, bold=True)
    sc_headers = ["Configuration", "Streams", "Batch cycle ms", "FPS/stream", "Total FPS"]
    sc_rows = []
    highlight = []
    scenarios = data["scenarios_edge"]
    target_labels = {"4 streams, YOLO batch=4"}
    for i, s in enumerate(scenarios):
        sc_rows.append([
            s["label"],
            str(s["n_streams"]),
            f"{s['batch_ms_edge']:.1f} ms",
            f"{s['fps_per_stream']:.1f}",
            f"{s['total_system_fps']:.1f}",
        ])
        if s["label"] in target_labels:
            highlight.append(i + 1)
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(4.35),
                     Inches(CONTENT_W), Inches(2.15), sc_headers, sc_rows,
                     highlight_rows=highlight, font_size=10)

    add_bullet_box(slide, CONTENT_LEFT, 6.55, CONTENT_W, 0.55, [
        ("Key findings", C.ACCENT_BLUE, True),
        ("• Batching amortizes kernel overhead big: 4 streams each at 25.9 FPS (not 36/4=9). 480p per stream or a second NPU unlocks even more.", C.ACCENT_GREEN, True),
    ], font_size=10)


def slide_llm_bakeoff(prs: Presentation):
    """Slide: Qwen3-30B-A3B LLM bake-off — 5090 measured + NPU tier actuals."""
    path = Path("data/output/bakeoff/llm_edge_projection.json")
    sumpath = Path("data/output/bakeoff/llm_summary.json")
    if not path.exists() or not sumpath.exists():
        return
    edge = json.loads(path.read_text())
    summary = json.loads(sumpath.read_text())
    tier_proj = edge.get("tier_projections", {})

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "LLM Bake-off — Qwen3-30B-A3B MoE (30B total / 3B active)",
                       "5090 measured; edge numbers from vendor NPU Low/Mid/High benchmarks (Qwen3-30B-A3B Q4_K_M, 1K prompt)")
    add_pipeline_strip(slide, ["FFmpeg", "YOLO-seg FP8 (TRT)",
                                "CLIP FP8 (TRT)", "SQLite", ("Qwen3-30B-A3B MoE", True)])

    # Table 1 — 5090 measured per quant
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(1.9), Inches(CONTENT_W), Inches(0.3),
                 "Measured on RTX 5090 (llama.cpp / llama-cpp-python 0.3.20):",
                 font_size=11, color=C.ACCENT_PURPLE, bold=True)
    headers1 = ["Quant", "GGUF size", "5090 prefill @2K", "5090 decode (256 tok)",
                "5090 RAG (8K+2K)"]
    rows1 = []
    for quant in ("Q4_K_M", "Q5_K_M", "Q8_0"):
        s = summary.get(quant, {})
        if not s or "error" in s:
            rows1.append([quant, "—", "—", "—", "—"])
            continue
        pf2k = next((r["prefill_tok_s"] for r in s["prefill_sweep"] if r["n_prompt"] == 2048), 0)
        dec = s["decode_sweep"][-1]["decode_tok_s"]
        rag_sec = s["rag"]["total_ms"] / 1000 if "rag" in s else 0
        rows1.append([
            quant,
            f"{s['gguf_size_gb']:.1f} GB",
            f"{pf2k:.0f} tok/s",
            f"{dec:.1f} tok/s",
            f"{rag_sec:.1f} s",
        ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(2.22),
                     Inches(CONTENT_W), Inches(1.2), headers1, rows1, font_size=10)

    # Table 2 — vendor NPU tier actuals (Q4_K_M production target)
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(3.65), Inches(CONTENT_W), Inches(0.3),
                 "Edge NPU tier actuals — vendor benchmarks at Q4_K_M (production target):",
                 font_size=11, color=C.ACCENT_INDIGO, bold=True)
    headers2 = ["NPU tier", "Memory bus", "TTFT 1K prompt",
                "Decode tok/s (Q4_K_M)",
                "Short answer (200 tok)", "RAG (8K+2K)"]
    rows2 = []
    highlight = []
    for i, tier in enumerate(("NPU Low-LP5", "NPU Mid", "NPU High")):
        tp = tier_proj.get(tier)
        if not tp:
            rows2.append([tier, "—", "—", "—", "—", "—"])
            continue
        q4 = tp["per_quant"]["Q4_K_M"]
        rows2.append([
            tier,
            tp["bus"],
            f"{tp['reference_ttft_1k_sec']*1000:.0f} ms",
            f"{q4['decode_tok_s']:.1f} tok/s",
            f"{q4['short_answer_ms']/1000:.1f} s",
            f"{q4['rag_total_sec']:.0f} s",
        ])
        if tier == "NPU Mid":
            highlight.append(i + 1)
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(3.97),
                     Inches(CONTENT_W), Inches(1.35), headers2, rows2,
                     highlight_rows=highlight, font_size=10)

    add_bullet_box(slide, CONTENT_LEFT, 5.55, CONTENT_W, 1.6, [
        ("Key findings", C.ACCENT_BLUE, True),
        ("• MoE wins on bandwidth. Cross-reference from a Kyle-merged production Q4_K_M of the same model on the same 5090 host: 155 tok/s sustained / 192 peak (Prometheus, prod traffic) — within 3% of our synthetic 159 tok/s RAG decode. Synthetic numbers generalize.", C.ACCENT_GREEN, True),
        ("• MoE 30B/3B-active beats dense Qwen 2.5 14B Q4_K_M on the same 5090 by ~25-40% decode (155 vs 85-140 tok/s) despite 2× total params — the 3B active footprint is what BW sees per token.", C.ACCENT_AMBER, True),
        ("• Vendor NPU actuals beat our BW-only edge projection by ~2.3×. Purpose-built LLM silicon (memory controllers, expert routing, tiling) > llama.cpp on a desktop GPU.", C.ACCENT_INDIGO, True),
        ("• NPU Mid: 5.3 s for 200-tok answer / 57 s for RAG 8K+2K. NPU High ~25% better. Q4_K_M is the shipping quant.", C.TEXT_DIM),
    ], font_size=9)


def slide_llm_duty_cycle(prs: Presentation):
    """Slide: vision FPS vs LLM query rate across NPU Low/Mid/High tiers."""
    path = Path("data/output/bakeoff/llm_edge_projection.json")
    if not path.exists():
        return
    edge = json.loads(path.read_text())
    tier_proj = edge.get("tier_projections", {})
    if not tier_proj:
        return

    vision_fps = 36.0          # full-stack shipping target (720p Edge MPU = NPU Mid)

    # Q4_K_M short-answer and RAG answer times per tier (milliseconds)
    tiers_data = {}
    for tier in ("NPU Low-LP5", "NPU Mid", "NPU High"):
        tp = tier_proj.get(tier)
        if not tp:
            continue
        q4 = tp["per_quant"]["Q4_K_M"]
        tiers_data[tier] = {
            "short_ms": q4["short_answer_ms"],
            "rag_ms": q4["rag_total_ms"],
        }

    import matplotlib.pyplot as plt
    import numpy as np
    fig, (ax_short, ax_long) = plt.subplots(1, 2, figsize=(11, 4.2),
                                             facecolor=MPL_COLORS["bg_slide"])

    qpm = np.linspace(0, 120, 200)
    qps = qpm / 60.0

    tier_colors = {
        "NPU Low-LP5":  MPL_COLORS["red"],
        "NPU Mid":  MPL_COLORS["orange"],
        "NPU High": MPL_COLORS["green"],
    }

    for tier, col in tier_colors.items():
        if tier not in tiers_data:
            continue
        d = tiers_data[tier]
        duty_short = qps * d["short_ms"] / 1000
        duty_rag   = qps * d["rag_ms"] / 1000
        fps_short = np.clip(vision_fps * (1 - duty_short), 0, vision_fps)
        fps_rag   = np.clip(vision_fps * (1 - duty_rag),   0, vision_fps)
        ax_short.plot(qpm, fps_short, color=col, linewidth=2.2,
                      label=f"{tier} ({d['short_ms']/1000:.1f} s / answer)")
        ax_long.plot(qpm, fps_rag, color=col, linewidth=2.2,
                     label=f"{tier} ({d['rag_ms']/1000:.0f} s / RAG)")

    for ax, title in [(ax_short, "Short answer (200 tokens)"),
                      (ax_long,  "RAG answer (8K prompt + 2K response)")]:
        for thr in (30, 15, 10):
            ax.axhline(y=thr, color=MPL_COLORS["dim"], linestyle=":", alpha=0.5, linewidth=1)
        ax.set_facecolor(MPL_COLORS["bg_slide"])
        ax.set_xlabel("LLM queries per minute", color=MPL_COLORS["text"], fontsize=10)
        ax.set_ylabel("Effective vision FPS", color=MPL_COLORS["text"], fontsize=10)
        ax.set_title(title, color=MPL_COLORS["text"], fontsize=11, fontweight="bold", pad=6)
        ax.set_xlim(0, 120)
        ax.set_ylim(0, vision_fps + 2)
        ax.tick_params(colors=MPL_COLORS["dim"], labelsize=9)
        ax.spines[:].set_color(MPL_COLORS["grid"])
        ax.grid(True, color=MPL_COLORS["grid"], alpha=0.3)
        ax.legend(loc="lower left", facecolor=MPL_COLORS["bg_slide"],
                  edgecolor=MPL_COLORS["grid"], labelcolor=MPL_COLORS["text"], fontsize=9)
        # Annotate threshold labels on the right edge of the RAG chart only
        if ax is ax_long:
            for thr, lbl in [(30, "30 FPS"), (15, "15 FPS"), (10, "10 FPS")]:
                ax.text(118, thr + 0.7, lbl, color=MPL_COLORS["dim"], fontsize=8, ha="right")

    fig.suptitle("Vision FPS under shared-NPU LLM load — three tiers",
                 color=MPL_COLORS["text"], fontsize=13, fontweight="bold")
    fig.tight_layout(rect=[0, 0, 1, 0.95])

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "NPU Duty-Cycle Trade-off — Vision vs LLM, Three Tiers",
                       "Vendor NPU Low/Mid/High actuals for Qwen3-30B-A3B Q4_K_M. One NPU shared — query pauses vision for answer duration.")
    add_pipeline_strip(slide, ["FFmpeg", "YOLO-seg FP8 (TRT)",
                                "CLIP FP8 (TRT)", "SQLite",
                                ("Qwen3-30B-A3B (shared NPU)", True)])

    buf = fig_to_image_stream(fig)
    slide.shapes.add_picture(buf, Inches(CONTENT_LEFT), Inches(1.9),
                              width=Inches(CONTENT_W))

    mid = tiers_data.get("NPU Mid", {})
    short_ms = mid.get("short_ms", 5300)
    rag_ms = mid.get("rag_ms", 57000)
    # Effective vision FPS at 1 Hz short-answer query rate on NPU Mid
    fps_mid_1hz_short = vision_fps * (1 - (60/60) * short_ms / 1000)
    # Effective vision FPS at 10 queries/min of short-answers on NPU Mid
    fps_mid_10qpm_short = vision_fps * max(0, 1 - (10/60) * short_ms / 1000)

    add_bullet_box(slide, CONTENT_LEFT, 5.9, CONTENT_W, 1.2, [
        ("Interpretation (NPU Mid highlighted — the 128-bit LPDDR5X target)", C.ACCENT_INDIGO, True),
        (f"• 1 short-answer query/minute: vision drops to {fps_mid_10qpm_short:.0f} FPS ({vision_fps-fps_mid_10qpm_short:.0f} FPS lost). Absolutely fine for review/batch use.",
         C.ACCENT_GREEN, True),
        (f"• RAG answers at 1/min on NPU Mid: vision obliterated — {rag_ms/1000:.0f} s/query × 1/min = {rag_ms/60/1000*100:.0f}% duty cycle, vision hits 0 FPS almost immediately.",
         C.ACCENT_RED, True),
        ("• NPU High roughly doubles query capacity vs NPU Mid at same vision FPS. Budget tier accordingly to expected query rate.", C.TEXT_DIM),
    ], font_size=10)


def slide_keyframe_debounce(prs: Presentation):
    """Slide: CLIP keyframe debouncing unlocks real-time on Hybrid V2."""
    path = Path("data/output/bakeoff/keyframe_debounce_summary.json")
    if not path.exists():
        return
    data = json.loads(path.read_text())
    per_res = data["per_resolution"]

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "CLIP Keyframe Debouncing — Unlocking Real-Time on Hybrid V2",
                       "Run YOLO every frame; CLIP only every Nth. YOLO is the real-time ceiling.")
    add_pipeline_strip(slide, ["FFmpeg", "YOLO-seg", ("CLIP FP8 @ 1 Hz", True),
                                "SQLite", "NLQ / LLM"])

    # Primary table: 720p sweep
    r720 = per_res["720p"]
    headers = ["N (native frames)", "Interval", "Top-1 stab", "Top-3 stab",
               "Edge ms / frame", "Edge FPS", "% of YOLO ceiling"]
    rows = []
    highlight = []
    ceiling = r720["edge_fps_yolo_only"]
    target_Ns = {30, 60}  # highlight the two sweet-spot candidates
    for i, row in enumerate(r720["rows"]):
        N = row["N_native_frames"]
        is_target = N in target_Ns
        if is_target:
            highlight.append(i + 1)
        rows.append([
            f"N = {N}",
            f"{row['keyframe_interval_sec']:.2f} s",
            f"{row['stability_top1']:.2f}",
            f"{row['stability_top3_jaccard']:.2f}",
            f"{row['eff_edge_ms_per_frame']:.1f}",
            f"{row['eff_edge_fps']:.1f}",
            f"{100 * row['eff_edge_fps'] / ceiling:.0f}%",
        ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                     Inches(CONTENT_W), Inches(2.7), headers, rows,
                     highlight_rows=highlight)
    add_text_box(slide, Inches(CONTENT_LEFT), Inches(4.7), Inches(CONTENT_W), Inches(0.3),
                 f"720p shown. YOLO-only ceiling: 720p {per_res['720p']['edge_fps_yolo_only']:.1f} FPS  "
                 f"•  1080p {per_res['1080p']['edge_fps_yolo_only']:.1f} FPS  "
                 f"•  4K {per_res['4K']['edge_fps_yolo_only']:.1f} FPS. "
                 f"Same debounce shape at all resolutions.",
                 font_size=10, color=C.TEXT_DIM)

    add_bullet_box(slide, CONTENT_LEFT, 5.1, CONTENT_W, 2.0, [
        ("Key findings", C.ACCENT_BLUE, True),
        ("• Debouncing works. 720p edge FPS climbs from 4.9 (N=1) to 16.0 at N=30 (CLIP every 1 s) — 93% of the YOLO-only ceiling", C.ACCENT_GREEN, True),
        "• YOLO-seg is now the real-time ceiling (17.3 FPS @ 720p edge). Further speedup needs Conv-FP8 / custom kernels.",
        ("• Recommended operating point: N = 30 (1 Hz CLIP rerun). Good FPS, top-3 stability 0.55 — new concept tags within one second of a scene change.", C.ACCENT_INDIGO, True),
        ("• N = 60 buys ~4% more FPS at the cost of 2-second stale tags — viable if the UI already displays cached tags on non-keyframes.", C.ACCENT_AMBER, True),
        "",
        ("Caveat on the stability numbers", C.ACCENT_ORANGE, True),
        ("• Bake-off frames are sampled at 1 fps, so N = 1..29 all collapse onto the 1-second-gap stability measurement. Sub-second drift is unmeasured and will be higher than shown.", C.TEXT_DIM),
        ("• N = 90 (3 s) is where top-1 collapses to 0.19 — upper bound for any naive fixed-interval scheme.", C.TEXT_DIM),
    ], font_size=11)


def slide_efficientsam3p1_textprompt(prs: Presentation):
    """Slide: "EfficientSAM3.1 — the text-prompt-capable smaller variant" — SAM 3.1 student."""
    path = Path("data/output/bakeoff/efficientsam3p1_summary.json")
    if not path.exists():
        return
    data = json.loads(path.read_text())
    by_res = data["by_resolution"]

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        "EfficientSAM3.1 ES-EV-S — text-prompt-capable smaller variant",
        "SAM 3.1 distilled student: 106M params (4× smaller than Option A). Keeps SAM 3's text-concept prompting natively.",
    )
    add_pipeline_strip(
        slide,
        ["FFmpeg", "(text concept prompt)", ("EfficientSAM3.1 ES-EV-S BF16", True), "SQLite", "NLQ / LLM"],
        accent_color=C.ACCENT_PURPLE,
    )

    bw_ratio = (1792.0 * 0.85) / (134.4 * 0.70)   # 16.19× (5090 eff BW ÷ NPU Mid eff BW, uniform 0.70 tier eff)

    # Primary table: per-resolution 5090 cost split + NPU Mid totals for n=1/5/20
    headers = ["Resolution",
                "set_image ms (5090)",
                "per-prompt ms (5090)",
                "n=1 5090 ms",
                "n=5 5090 ms",
                "n=20 5090 ms",
                "n=1 NPU Mid FPS",
                "n=5 NPU Mid FPS"]
    rows = []
    for res in ["720p", "1080p", "4K"]:
        if res not in by_res:
            continue
        r = by_res[res]
        n1 = r["per_frame_5090_ms"]["n_1_concept"]
        n5 = r["per_frame_5090_ms"]["n_5_concepts"]
        n20 = r["per_frame_5090_ms"]["n_20_concepts_exhaustive"]
        fps_n1_mid  = 1000.0 / (n1 * bw_ratio)
        fps_n5_mid  = 1000.0 / (n5 * bw_ratio)
        rows.append([
            res,
            f"{r['set_image_5090_p50_ms']:.1f} ms",
            f"{r['per_prompt_5090_p50_ms']:.1f} ms",
            f"{n1:.1f}", f"{n5:.1f}", f"{n20:.1f}",
            f"{fps_n1_mid:.2f}", f"{fps_n5_mid:.2f}",
        ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                      Inches(CONTENT_W), Inches(1.8), headers, rows,
                      font_size=9, header_font_size=10)

    # Comparison row
    comp_headers = ["Variant", "Total params", "Path", "720p NPU Mid (best-case)"]
    comp_rows = [
        ["SAM 3 BF16 (baseline)",                       "840M", "text-concept (native)",    "0.40 FPS"],
        ["EfficientSAM3 ES-EV-S BF16 (Option A)",       "424M", "box-prompt (batched)",     "2.59 FPS (18 boxes)"],
        ["EfficientSAM3.1 ES-EV-S BF16 (SAM3.1 student)", "106M", "text-concept (native)",  "2.33 FPS (1 concept)"],
        ["Keyhole shipping (TRT FP8 + 1 Hz CLIP)",       "~40M total two-stage", "detector + CLIP tagger", "36.1 FPS"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(3.9),
                      Inches(CONTENT_W), Inches(1.3), comp_headers, comp_rows,
                      highlight_rows=[3, 4],  # Option A-smaller + shipping
                      font_size=10)

    add_bullet_box(slide, CONTENT_LEFT, 5.3, CONTENT_W, 2.1, [
        ("What's new vs Option A", C.ACCENT_BLUE, True),
        ("• 4× fewer params (106M vs 424M). Vision encoder 31M (EfficientViT-S), text encoder 43M (MobileCLIP-S0 ctx=16) — both distilled. "
         "Peak VRAM on 5090: 1.8 GB (vs 3.0 GB for Option A).", C.TEXT_BRIGHT),
        ("• Lives on the upstream `stage1_sam3.1` branch with new arg convention: "
         "`model_name='s'`, `text_encoder_type='mobileclip-s0'`. The main-branch builder silently produces a "
         "state_dict mismatch if loaded the old way. That's why we deferred it — now unblocked.",
         C.TEXT_DIM),
        "",
        ("Workload shape matters", C.ACCENT_ORANGE, True),
        ("• Text-prompt latency is set_image (~10 ms, amortized) + 20 ms per concept. For 1-concept queries matches "
         "Option A's box-prompt throughput (2.3 vs 2.6 FPS). For 5-concept queries drops to 0.6 FPS; for exhaustive "
         "20-concept sweeps, 0.2 FPS — text-prompting is a linear-in-N cost, not free.",
         C.TEXT_BRIGHT),
        ("• Right shape for: 'find me all people' (1 concept) or 'find people and cars' (2 concepts). Wrong shape "
         "for exhaustive label-all-objects scans — our TRT-shipping stack does that at 36 FPS by using CLIP on "
         "detected crops, not by re-running the mask model per concept.", C.ACCENT_GREEN),
    ], font_size=10)


def slide_trt_yoloe26(prs: Presentation):
    """Slide: "Does TRT FP8 close the YOLOE-26 gap?" — negative result, structural gap."""
    path = Path("data/output/bakeoff/trt_yoloe26_summary.json")
    if not path.exists():
        return
    data = json.loads(path.read_text())
    recipes = data["recipes"]

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        "Does TRT FP8 close the one-model gap? — The honest answer: no.",
        "YOLOE-26S-PF → TRT FP8 gives ~17% speedup, not 3×. Gap to shipping is structural.",
    )
    add_pipeline_strip(
        slide,
        ["FFmpeg", ("YOLOE-26S-PF TRT FP8", True), "(no CLIP)", "SQLite", "NLQ / LLM"],
        accent_color=C.ACCENT_AMBER,
    )

    bw_ratio = (1792.0 * 0.85) / (134.4 * 0.70)   # 16.19× (5090 eff BW ÷ NPU Mid eff BW, uniform 0.70 tier eff)

    # Per-recipe per-resolution table
    headers = ["Recipe", "720p 5090 (p50)", "1080p 5090", "4K 5090",
                "720p NPU Mid FPS", "VRAM (5090)", "Speedup vs PT"]
    rows = []
    pt_ref_720 = recipes.get("pytorch_fp16", {}).get("by_resolution", {}).get("720p", {}).get("per_frame_ms_5090", {}).get("p50", 0)
    for tag in ["pytorch_fp16", "trt_fp16", "trt_fp8"]:
        r = recipes.get(tag)
        if not r:
            continue
        def ms(res): return r["by_resolution"].get(res, {}).get("per_frame_ms_5090", {}).get("p50", 0)
        p50_720 = ms("720p")
        ms_mid = p50_720 * bw_ratio
        fps_mid = 1000.0 / ms_mid if ms_mid > 0 else 0
        speedup = pt_ref_720 / p50_720 if p50_720 > 0 else 0
        rows.append([
            tag.replace("_", " "),
            f"{p50_720:.2f} ms",
            f"{ms('1080p'):.2f} ms",
            f"{ms('4K'):.2f} ms",
            f"{fps_mid:.2f} FPS",
            f"{r['peak_vram_mb_5090']:.0f} MB",
            f"{speedup:.2f}×" if tag != "pytorch_fp16" else "baseline",
        ])
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(1.9),
                      Inches(CONTENT_W), Inches(1.6), headers, rows,
                      highlight_rows=[3],   # TRT FP8 row
                      font_size=10)

    # Context: where YOLOE-26 TRT FP8 slots in the progression
    comp_headers = ["Stack", "720p NPU Mid FPS", "Note"]
    comp_rows = [
        ["SAM 3 BF16 (baseline)",                    "0.40 FPS",  "Dead-on-arrival."],
        ["EfficientSAM3 ES-EV-S BF16 (Apr 2026)",     "2.59 FPS",  "Community lite; box-prompt."],
        ["EfficientSAM3.1 ES-EV-S (SAM 3.1 student)", "2.33 FPS",  "Text-prompt; 1 concept."],
        ["YOLOE-26S-PF PyTorch FP16",                 "13.25 FPS", "Plain PyTorch, one model."],
        ["YOLOE-26S-PF TRT FP8 (optimized ceiling)",  "14.36 FPS", "One model fully optimized."],
        ["Keyhole shipping (TRT FP8 + 1 Hz CLIP)",    "36.10 FPS", "Two-stage, CLIP 1/30 frames."],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(3.7),
                      Inches(CONTENT_W), Inches(1.7), comp_headers, comp_rows,
                      highlight_rows=[5, 6],
                      font_size=10)

    add_bullet_box(slide, CONTENT_LEFT, 5.5, CONTENT_W, 1.9, [
        ("Why TRT FP8 helped so little", C.ACCENT_BLUE, True),
        ("• YOLOE-26 is 16M params. At that size, on a 5090, the bottleneck is **kernel launch overhead**, "
         "not tensor compute — FP8 matmul speedup doesn't help when matmul isn't the bottleneck.",
         C.TEXT_BRIGHT),
        ("• Compare to YOLO-seg FP8 TRT (our shipping Conv-unblock): that went 4.9 -> 36.8 FPS edge (+98%), "
         "because YOLO-seg's Conv-heavy graph does benefit from FP8 in TRT. Different architecture, different bottleneck.",
         C.TEXT_DIM),
        "",
        ("The REAL TRT win: 73% VRAM reduction (360 -> 99 MB)", C.ACCENT_GREEN, True),
        ("• Latency gain is modest but VRAM footprint drops dramatically. If NPU memory is tight and you're "
         "running concurrent streams, TRT is still worth it — just not for speed.",
         C.TEXT_BRIGHT),
        "",
        ("Takeaway: the gap to shipping is STRUCTURAL, not optimization", C.ACCENT_ORANGE, True),
        ("• Our two-stage stack runs YOLO-seg every frame + CLIP every 30th. YOLOE-26 runs the whole "
         "open-vocab head every frame. That design choice is the 2.4× gap — no amount of kernel "
         "optimization on YOLOE-26 will close it without a similar keyframe-debounce trick.",
         C.ACCENT_AMBER),
    ], font_size=10)


def slide_yoloe26_onemodel(prs: Presentation):
    """Slide: "Ultralytics YOLOE-26 — the one-model open-vocab alternative" — benches the
    Jan 2026 YOLOE-26S-PF release against our two-stage YOLO-seg + CLIP shipping stack.

    Post-ship watch on architectural simplification (Option B). YOLOE-26 collapses
    detector + open-vocab labeler into a single model with 4585-class built-in vocab.
    Worth benching as a pipeline-simplification story even though we ship something
    faster with the two-stage TRT FP8 stack.
    """
    path = Path("data/output/bakeoff/yoloe26_summary.json")
    if not path.exists():
        return
    data = json.loads(path.read_text())

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        "Option-B watch: Ultralytics YOLOE-26 — one-model open-vocab",
        "Single model replaces YOLO-seg + CLIP. 13 FPS NPU Mid @ 720p in PyTorch FP16 — 3× slower than our shipping stack, but 10× simpler.",
    )
    add_pipeline_strip(
        slide,
        ["FFmpeg", ("YOLOE-26S-PF  (one model)", True), "(no CLIP)", "SQLite", "NLQ / LLM"],
        accent_color=C.ACCENT_AMBER,
    )

    # Per-variant, per-resolution numbers
    bw_ratio = (1792.0 * 0.85) / (134.4 * 0.70)   # 16.19× (5090 eff BW ÷ NPU Mid eff BW, uniform 0.70 tier eff)
    headers = ["Variant", "Res", "Params",
                "5090 ms (p50)", "NPU Mid ms (BW-scaled)", "NPU Mid FPS",
                "Box recall vs YOLO11x"]
    rows = []
    tag_display = {
        "text_prompt_s":  "YOLOE-26S (text prompts: 20 concepts)",
        "prompt_free_s":  "YOLOE-26S-PF (prompt-free, 4585 vocab)",
    }
    for tag in ["text_prompt_s", "prompt_free_s"]:
        v = data["variants"].get(tag)
        if not v or "error" in v:
            continue
        for res in ["720p", "1080p", "4K"]:
            r = v["by_resolution"].get(res)
            if not r:
                continue
            ms_5090 = r["per_frame_ms_5090"]["p50"]
            ms_mid = ms_5090 * bw_ratio
            rows.append([
                tag_display[tag] if res == "720p" else "",
                res,
                f"{v['params_m']:.1f}M" if res == "720p" else "",
                f"{ms_5090:.2f} ms",
                f"{ms_mid:.1f} ms",
                f"{1000/ms_mid:.2f} FPS",
                f"{r['box_recall_vs_yolo11x']:.3f}",
            ])
    add_styled_table(
        slide, Inches(CONTENT_LEFT), Inches(1.9),
        Inches(CONTENT_W), Inches(2.2), headers, rows,
        font_size=9,
    )

    # Head-to-head at 720p NPU Mid
    comp_headers = ["Architecture", "NPU Mid FPS @ 720p", "Models on the NPU", "Note"]
    comp_rows = [
        ["SAM 3 BF16",                           "0.40 FPS",  "1 (SAM 3)",           "Dead-on-arrival baseline."],
        ["EfficientSAM3 ES-EV-S BF16 (Apr 2026)", "2.59 FPS",  "1 (community lite)",  "Community SAM 3 Lite, open-vocab preserved."],
        ["YOLOE-26S-PF FP16 (Jan 2026)",          "13.25 FPS", "1 (one-model)",       "Replaces YOLO-seg + CLIP. Simplest stack."],
        ["Keyhole shipping (TRT FP8 + 1 Hz CLIP)", "36.1 FPS",  "2 (YOLO-seg + CLIP)", "Our optimized two-stage pipeline."],
    ]
    add_styled_table(
        slide, Inches(CONTENT_LEFT), Inches(4.3),
        Inches(CONTENT_W), Inches(1.3), comp_headers, comp_rows,
        highlight_rows=[4],
        font_size=10,
    )

    add_bullet_box(slide, CONTENT_LEFT, 5.7, CONTENT_W, 1.7, [
        ("Why this matters: architectural simplification", C.ACCENT_BLUE, True),
        ("• One model vs two. No CLIP pass, no two-stage orchestration, no separate keyframe-debounce logic. "
         "Same output shape as YOLO-seg today: boxes + masks + labels per frame.", C.TEXT_BRIGHT),
        ("• Recall vs our YOLO11x prompt boxes: 65% @ 720p, 86% @ 1080p, 81% @ 4K (prompt-free variant). "
         "Text-prompted with our 20-concept SAM 3 list: lower recall at 720p (0.60) but comparable at higher resolutions.",
         C.TEXT_DIM),
        "",
        ("What's missing vs shipping", C.ACCENT_ORANGE, True),
        ("• No TRT/FP8 path explored yet — this is plain PyTorch FP16. TRT FP8 would potentially close most of the gap to 36 FPS.",
         C.TEXT_BRIGHT),
        ("• License is AGPL-3.0 (same as ultralytics itself). Commercial licensing via Ultralytics if needed for a proprietary product.",
         C.TEXT_DIM),
        ("• Accuracy tradeoff: 65-86% recall depending on resolution — acceptable for a research demo; would need validation against Keyhole's actual use cases before any production swap.",
         C.ACCENT_AMBER),
    ], font_size=10)


def slide_efficientsam3_community(prs: Presentation):
    """Slide: "The community finally shipped a SAM 3 Lite" — EfficientSAM3 ES-EV-S benched against our shipping stack.

    Post-ship watch on roadmap item #9 (SAM 3 Lite). As of April 2026 the
    community has released EfficientSAM3 (ES-EV-S, Apache-2.0, ~424M total /
    26M vision backbone). We bench it against cached frames + YOLO prompt
    boxes + SAM 3 reference masks from our existing bake-off to check whether
    the community caught up.
    """
    path = Path("data/output/bakeoff/efficientsam3_summary.json")
    if not path.exists():
        return
    data = json.loads(path.read_text())

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        "Post-ship watch: the community finally shipped a SAM 3 Lite",
        "EfficientSAM3 ES-EV-S (Apr 2026, Apache-2.0) — 6.5× faster than SAM 3, still 13× behind our shipping stack.",
    )
    add_pipeline_strip(
        slide,
        ["FFmpeg", "YOLO 11x", ("EfficientSAM3 ES-EV-S BF16", True), "SQLite", "NLQ / LLM"],
        accent_color=C.ACCENT_PURPLE,
    )

    # Per-resolution measured + projected numbers
    by_res = data["by_resolution"]
    headers = ["Resolution", "5090 ms (p50)", "NPU Mid ms (BW-scaled)",
                "NPU Mid FPS", "Mean IoU vs SAM 3", "VRAM"]
    rows = []
    for res in ["720p", "1080p", "4K"]:
        if res not in by_res:
            continue
        r = by_res[res]
        ms_5090 = r["per_frame_ms_5090"]["p50"]
        # Bandwidth ratio from npu_model.py (5090 effective BW ÷ NPU Mid effective BW)
        bw_ratio = (1792.0 * 0.85) / (134.4 * 0.70)   # = 16.19x (uniform 0.70 NPU tier efficiency)
        ms_mid = ms_5090 * bw_ratio
        fps_mid = 1000.0 / ms_mid if ms_mid > 0 else 0.0
        iou = r["iou_vs_sam3"]["mean"]
        rows.append([
            res,
            f"{ms_5090:.1f} ms",
            f"{ms_mid:.0f} ms",
            f"{fps_mid:.2f} FPS",
            f"{iou:.3f}",
            f"{data['peak_vram_mb_5090']:.0f} MB (5090 BF16)" if res == "720p" else "—",
        ])
    add_styled_table(
        slide, Inches(CONTENT_LEFT), Inches(1.9),
        Inches(CONTENT_W), Inches(1.8), headers, rows,
    )

    # Head-to-head comparison at 720p NPU Mid
    comp_headers = ["Model / stack", "NPU Mid FPS @ 720p", "Note"]
    comp_rows = [
        ["SAM 3 BF16 (Meta baseline)",    "0.40 FPS",  "The thing we are replacing — DOA at the edge."],
        ["EfficientSAM3 ES-EV-S (BF16)",  "2.59 FPS",  "Community lite, Apr 2026. 6.5× over SAM 3."],
        ["EfficientSAM-Small FP8 (ours)", "4.93 FPS",  "Mask-only, SAM 1/2 era model + FP8 quant."],
        ["Keyhole shipping (TRT FP8, 1 Hz CLIP)", "36.1 FPS", "YOLO-seg FP8 + CLIP FP8, both on TensorRT."],
    ]
    add_styled_table(
        slide, Inches(CONTENT_LEFT), Inches(4.0),
        Inches(CONTENT_W), Inches(1.5), comp_headers, comp_rows,
        highlight_rows=[4],   # shipping stack
    )

    add_bullet_box(slide, CONTENT_LEFT, 5.9, CONTENT_W, 1.7, [
        ("What just happened", C.ACCENT_BLUE, True),
        ("• ES-EV-S (EfficientViT-B0 vision backbone, distilled from SAM 3's 462M encoder to 26M) dropped on HF + GitHub "
         "mid-April 2026, with Apache-2.0 licensing and preserved text-concept prompting.", C.TEXT_BRIGHT),
        ("• We benched it in a separate Python 3.12 / uv venv against our existing cached frames + YOLO prompt boxes — "
         "no re-extraction needed. IoU vs SAM 3 = 0.575 (moderate mask agreement). 5090 BF16, BW-scaled to NPU Mid.",
         C.TEXT_DIM),
        "",
        ("Why we still ship YOLO-seg FP8 TRT + CLIP FP8 TRT", C.ACCENT_INDIGO, True),
        ("• ES-EV-S is 6.5× faster than SAM 3 but 13× slower than our shipping stack. It has no TRT path, no FP8 path, "
         "and the vision backbone is only 6% of the total 424M params — the rest is text encoder + segmentation head.",
         C.TEXT_BRIGHT),
        ("• Takeaway: even when the community ships a credible SAM 3 Lite, an optimized two-stage pipeline (detector + "
         "open-vocab CLIP) on edge-native kernels beats a monolithic open-vocab SAM by an order of magnitude.",
         C.ACCENT_GREEN),
    ], font_size=10)


def slide_trt_takeaways(prs: Presentation):
    """Slide: synthesize the 3 TRT bake-offs into one decision rule + compiler caveat.

    Consumes existing bake-off JSONs — no new measurements required. The
    numbers here must stay in sync with slides 43/44/51; if a bake-off is
    re-run, this slide re-reads and updates automatically.
    """
    # Shipping YOLO-seg TRT FP8 vs FP16 at 720p
    trt_yolo = Path("data/output/bakeoff/trt_yolo_edge_projection.json")
    trt_clip = Path("data/output/bakeoff/trt_clip_edge_projection.json")
    trt_e26  = Path("data/output/bakeoff/trt_yoloe26_summary.json")
    if not (trt_yolo.exists() and trt_clip.exists() and trt_e26.exists()):
        return

    yolo_proj = json.loads(trt_yolo.read_text())["projections"]["720p"]
    clip_proj = json.loads(trt_clip.read_text())["projections"]["720p"]
    e26       = json.loads(trt_e26.read_text())["recipes"]

    # FPS/ms headline numbers — all edge, NPU Mid
    yolo_fp16_edge = yolo_proj.get("fp16", {}).get("projected_fps_edge", 0)
    yolo_fp8_edge  = yolo_proj.get("fp8",  {}).get("projected_fps_edge", 0)
    clip_bf16_edge_ms = clip_proj.get("bf16_torch", {}).get("projected_clip_ms_edge", 0)
    clip_fp8_edge_ms  = clip_proj.get("fp8",        {}).get("projected_clip_ms_edge", 0)
    pt16_720 = e26.get("pytorch_fp16", {}).get("by_resolution", {}).get("720p", {}).get("per_frame_ms_5090", {}).get("p50", 0)
    tr16_720 = e26.get("trt_fp16",     {}).get("by_resolution", {}).get("720p", {}).get("per_frame_ms_5090", {}).get("p50", 0)
    tr8_720  = e26.get("trt_fp8",      {}).get("by_resolution", {}).get("720p", {}).get("per_frame_ms_5090", {}).get("p50", 0)
    e26_vram_pt = e26.get("pytorch_fp16", {}).get("peak_vram_mb_5090", 0)
    e26_vram_fp8 = e26.get("trt_fp8",     {}).get("peak_vram_mb_5090", 0)

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        "TensorRT — where it pays and where it doesn't",
        "Three TRT FP8 bake-offs, one rule of thumb: optimization loves big per-kernel work.",
    )
    add_pipeline_strip(
        slide,
        ["FFmpeg", ("YOLO-seg FP8 TRT", True), ("CLIP FP8 TRT", True), "SQLite", "NLQ / LLM"],
        accent_color=C.ACCENT_GREEN,
    )

    # Decision matrix table
    headers = ["Model", "Arch / bottleneck", "TRT FP8 result (720p edge)", "Verdict"]
    rows = [
        ["YOLO-seg 11s (10M params)",
         "Dense Conv backbone, matmul-bound",
         f"{yolo_fp16_edge:.1f} FPS (FP16) → {yolo_fp8_edge:.1f} FPS (FP8), "
         f"+{(yolo_fp8_edge/yolo_fp16_edge - 1)*100:.0f}% — full model activation halving works",
         "SHIPS — the core FP8 unblock"],
        ["OpenCLIP ViT-B-32 visual (88M params)",
         "ViT attention + MLP, matmul-bound",
         f"{clip_bf16_edge_ms:.1f} ms BF16 → {clip_fp8_edge_ms:.1f} ms FP8, "
         f"{clip_bf16_edge_ms/clip_fp8_edge_ms:.1f}× faster — Top-1 agreement 0.964",
         "SHIPS — halves CLIP cost"],
        ["YOLOE-26S-PF (16M params)",
         "Small model + complex open-vocab head, kernel-launch bound",
         f"PT {pt16_720:.1f} ms → TRT FP16 {tr16_720:.1f} ms ({pt16_720/tr16_720:.2f}×); "
         f"FP8 {tr8_720:.1f} ms — FP16→FP8 gains ~0% on the matmul",
         "DOESN'T close gap to shipping"],
    ]
    add_styled_table(
        slide, Inches(CONTENT_LEFT), Inches(1.9),
        Inches(CONTENT_W), Inches(2.3), headers, rows,
        highlight_rows=[1, 2], font_size=10, header_font_size=11,
    )

    # Takeaways
    add_bullet_box(slide, CONTENT_LEFT, 4.35, CONTENT_W, 1.6, [
        ("Rule of thumb: TRT FP8 pays off when the kernel is big", C.ACCENT_BLUE, True),
        ("• WORKS: dense Conv (YOLO-seg), dense ViT (CLIP) — FP8 matmul throughput is the bottleneck, and TRT kernel fusion amortizes launch cost over real compute.",
         C.ACCENT_GREEN),
        ("• UNDERPERFORMS: small param-count models with complex graphs (YOLOE-26 open-vocab head). "
         "At 16M params the kernel-launch tax dominates; FP8 can't help with work that isn't matmul.",
         C.ACCENT_AMBER),
        (f"• Orthogonal win: TRT FP8 still cuts VRAM ~{(1 - e26_vram_fp8/e26_vram_pt)*100:.0f}% on YOLOE-26 "
         f"({e26_vram_pt:.0f} → {e26_vram_fp8:.0f} MB) — worth it for multi-stream even when latency gain is modest.",
         C.TEXT_DIM),
    ], font_size=10)

    # Edge compiler caveat — mirrors the sizer's new slider
    add_bullet_box(slide, CONTENT_LEFT, 6.05, CONTENT_W, 1.0, [
        ("Reality check: these numbers came from NVIDIA TensorRT on 5090", C.ACCENT_ORANGE, True),
        ("• All three bake-offs used TRT 10.16 — a best-in-class compiler on best-in-class silicon. "
         "Vendor edge-NPU compilers (Qualcomm SNPE, MediaTek NeuroPilot, OpenVINO-NPU, Hailo SDK) typically "
         "extract 50–75% of the same theoretical peak on first-gen toolchains.",
         C.TEXT_BRIGHT),
        ("• The sizer's “Edge compiler quality vs TensorRT” slider applies a 0–50% haircut to projected "
         "vision FPS to model this gap. Default is 1.00 (parity, optimistic); set to 0.75 for a realistic "
         "first-gen deployment plan.",
         C.ACCENT_INDIGO),
    ], font_size=10)


def _load_ncu_bundle() -> Optional[dict]:
    """Load the vendored ncu sizer bundle (16 workloads) if present."""
    path = Path("data/output/ncu/sizer_bundle.json")
    if not path.exists():
        return None
    return json.loads(path.read_text())


def _ncu_family(workload_id: str) -> str:
    """Group workloads by pipeline family for the table."""
    if workload_id.startswith("sam3_bf16"):
        return "Reference (what we replaced)"
    if workload_id.startswith("efficientsam3p1"):
        return "Community SAM 3.1 Lite"
    if workload_id.startswith("efficientsam3"):
        return "Community SAM 3 Lite"
    if workload_id in {"mobilesam", "efficientsam_tiny", "efficientsam_small"}:
        return "Mask-model bake-off"
    if workload_id in {"yolo_seg", "yolo_seg_fp16_trt", "yolo_seg_fp8_trt", "clip_trt"}:
        return "Shipping (TRT two-stage)"
    if workload_id.startswith("yoloe26"):
        return "YOLOE-26 one-model"
    return "Other"


def build_ncu_headline_chart(bundle: dict):
    """Horizontal log-scale bar chart of per-forward DRAM MB for headline workloads.

    Orders workloads worst → best so the 515× SAM3 → shipping gap reads top-to-bottom.
    """
    # Hand-picked representatives so the chart isn't overcrowded.
    picks = [
        ("sam3_bf16_reference",           "SAM 3 BF16 (reference)",                 MPL_COLORS["red"]),
        ("efficientsam_small",            "EfficientSAM-Small (bake-off)",          MPL_COLORS["orange"]),
        ("efficientsam3_es_ev_s",         "EfficientSAM3 ES-EV-S (Community Lite)", MPL_COLORS["orange"]),
        ("mobilesam",                     "MobileSAM (bake-off)",                   MPL_COLORS["orange"]),
        ("efficientsam3p1_es_ev_s__set_image", "EfficientSAM3.1 ES-EV-S (SAM3.1 student)", MPL_COLORS["orange"]),
        ("yoloe26_pytorch_fp16",          "YOLOE-26S-PF PyTorch FP16",              MPL_COLORS["purple"]),
        ("yoloe26_trt_fp8",               "YOLOE-26S-PF TRT FP8",                   MPL_COLORS["purple"]),
        ("yolo_seg",                      "YOLO11s-seg PyTorch FP32",               MPL_COLORS["blue"]),
        ("clip_trt",                      "OpenCLIP visual TRT",                    MPL_COLORS["blue"]),
        ("yolo_seg_fp8_trt",              "YOLO11s-seg TRT FP8 (SHIPPING)",         MPL_COLORS["green"]),
    ]
    by_id = {w["workload_id"]: w for w in bundle["workloads"]}

    labels, mbs, colors = [], [], []
    for wid, label, color in picks:
        if wid not in by_id:
            continue
        labels.append(label)
        mbs.append(by_id[wid]["per_forward"]["dram_mb"])
        colors.append(color)

    fig, ax = plt.subplots(1, 1, figsize=(11, 4.4), facecolor=MPL_COLORS["bg_slide"])
    ax.set_facecolor(MPL_COLORS["bg_slide"])

    # Reverse so shipping lands at the top of the chart (best first)
    y_pos = np.arange(len(labels))
    bars = ax.barh(y_pos[::-1], mbs, color=colors, edgecolor="none", height=0.65)

    ax.set_yticks(y_pos[::-1])
    ax.set_yticklabels(labels, color=MPL_COLORS["text"], fontsize=9)
    ax.set_xscale("log")
    ax.set_xlabel("DRAM bytes moved per forward pass (MB, log scale)",
                  color=MPL_COLORS["text"], fontsize=10)
    ax.set_xlim(100, 300000)
    ax.tick_params(colors=MPL_COLORS["dim"])
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    for spine in ("left", "bottom"):
        ax.spines[spine].set_color(MPL_COLORS["dim"])
    ax.grid(True, axis="x", which="both", alpha=0.15, color=MPL_COLORS["grid"])

    # Value labels on each bar
    for bar, mb in zip(bars, mbs):
        x = bar.get_width()
        y = bar.get_y() + bar.get_height() / 2
        if mb >= 1000:
            txt = f"{mb/1000:,.1f} GB"
        else:
            txt = f"{mb:,.0f} MB"
        ax.text(x * 1.15, y, txt, va="center", ha="left",
                fontsize=8, color=MPL_COLORS["text"])

    # Annotate the 515× gap
    ship_mb = by_id["yolo_seg_fp8_trt"]["per_forward"]["dram_mb"] + \
              by_id["clip_trt"]["per_forward"]["dram_mb"] / 30.0  # CLIP runs at 1 Hz, 30 fps video
    sam3_mb = by_id["sam3_bf16_reference"]["per_forward"]["dram_mb"]
    ratio = sam3_mb / ship_mb
    ax.set_title(
        f"SAM 3 reference moves {sam3_mb/1000:,.0f} GB/frame — "
        f"shipping TRT FP8 + 1 Hz CLIP moves ~{ship_mb:.0f} MB/frame → {ratio:,.0f}× lighter",
        color=MPL_COLORS["text"], fontsize=11, fontweight="bold", pad=10,
    )
    fig.tight_layout(pad=1.0)
    return fig


def slide_ncu_headline(prs: Presentation):
    """Slide: Nsight Compute measured DRAM bandwidth — the 515× gap."""
    bundle = _load_ncu_bundle()
    if bundle is None:
        return

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        "Measured DRAM bandwidth — how the shipping pipeline actually moves memory",
        f"Nsight Compute 2026 per-kernel bytes × NVTX per-workload, {bundle['n_workloads']} workloads on RTX 5090",
    )
    add_pipeline_strip(
        slide,
        [("FFmpeg", False), ("YOLO-seg FP8 TRT", True), ("CLIP FP8 TRT @ 1 Hz", True),
         ("SQLite", False), ("NLQ / LLM", False)],
        accent_color=C.ACCENT_GREEN,
    )

    fig = build_ncu_headline_chart(bundle)
    img = fig_to_image_stream(fig)
    slide.shapes.add_picture(img, Inches(CONTENT_LEFT), Inches(1.9),
                              width=Inches(9.0))

    # Right-side callout column
    add_text_box(slide, Inches(9.8), Inches(1.95), Inches(3.4), Inches(0.35),
                 "NPU Mid budget", font_size=12, color=C.ACCENT_INDIGO, bold=True)
    add_text_box(slide, Inches(9.8), Inches(2.25), Inches(3.4), Inches(1.2),
                 "128-bit LPDDR5X @ 8.4 GT/s\n"
                 "= 134.4 GB/s theoretical\n"
                 "= 94.1 GB/s effective (70%)",
                 font_size=10, color=C.TEXT_BRIGHT)

    add_text_box(slide, Inches(9.8), Inches(3.35), Inches(3.4), Inches(0.35),
                 "Shipping @ 30 fps, 1 stream", font_size=12, color=C.ACCENT_GREEN, bold=True)
    add_text_box(slide, Inches(9.8), Inches(3.65), Inches(3.4), Inches(1.0),
                 "≈ 8.3 GB/s actual DRAM\n"
                 "(9% of budget — 91% headroom\n"
                 "for concurrent LLM + streams)",
                 font_size=10, color=C.TEXT_BRIGHT)

    add_text_box(slide, Inches(9.8), Inches(4.65), Inches(3.4), Inches(0.35),
                 "Shipping × 4 streams", font_size=12, color=C.ACCENT_AMBER, bold=True)
    add_text_box(slide, Inches(9.8), Inches(4.95), Inches(3.4), Inches(1.0),
                 "≈ 23 GB/s actual DRAM\n"
                 "(24% of budget — still leaves\n"
                 "~71 GB/s for LLM duty-cycling)",
                 font_size=10, color=C.TEXT_BRIGHT)

    add_bullet_box(slide, CONTENT_LEFT, 6.3, CONTENT_W, 0.85, [
        ("Why this matters", C.ACCENT_BLUE, True),
        ("• Before ncu, the sizer assumed every pipeline saturates the bus. Measured bytes show "
         "the shipping stack runs at 9% of NPU Mid — the engineering win is real, with massive headroom.", C.ACCENT_GREEN),
        ("• Conversely, EfficientSAM-Small and the community SAM 3 Lites physically cannot fit: per-forward "
         "DRAM already exceeds what the bus can deliver per second at any usable FPS.", C.ACCENT_RED),
    ], font_size=10)


def slide_ncu_workload_table(prs: Presentation):
    """Slide: Full per-workload table — measured DRAM and BW-bound FPS ceilings."""
    bundle = _load_ncu_bundle()
    if bundle is None:
        return

    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        f"All {bundle['n_workloads']} measured workloads — per-forward DRAM and NPU Mid ceiling",
        "Nsight Compute bytes ÷ NVTX-bounded forwards. FPS ceiling = 94.08 GB/s ÷ MB/forward (BW-bound only).",
    )

    # Order by family, best-case first within each family (low MB = good)
    rank_family = {
        "Reference (what we replaced)": 0,
        "Mask-model bake-off":          1,
        "Community SAM 3 Lite":         2,
        "Community SAM 3.1 Lite":       3,
        "YOLOE-26 one-model":           4,
        "Shipping (TRT two-stage)":     5,
        "Other":                        6,
    }
    rows_raw = []
    for w in bundle["workloads"]:
        wid = w["workload_id"]
        fam = _ncu_family(wid)
        rows_raw.append((
            rank_family[fam], fam, wid,
            w["n_kernels_total"],
            w["per_forward"]["dram_mb"],
            w["edge_projection_npu_mid"]["bw_bound_fps_max"],
        ))
    rows_raw.sort(key=lambda r: (r[0], r[4]))

    # Pretty-print labels + highlight shipping row
    pretty = {
        "sam3_bf16_reference":                  "SAM 3 reference (BF16)",
        "mobilesam":                            "MobileSAM",
        "efficientsam_tiny":                    "EfficientSAM Tiny",
        "efficientsam_small":                   "EfficientSAM Small",
        "efficientsam3_es_ev_s":                "EfficientSAM3 ES-EV-S",
        "efficientsam3p1_es_ev_s__set_image":   "EfficientSAM3.1 set_image",
        "efficientsam3p1_es_ev_s__text_prompt": "EfficientSAM3.1 text_prompt",
        "yoloe26_pytorch_fp16":                 "YOLOE-26S-PF PyTorch FP16",
        "yoloe26_text_prompt_s":                "YOLOE-26S text_prompt",
        "yoloe26_prompt_free_s":                "YOLOE-26S prompt_free",
        "yoloe26_trt_fp16":                     "YOLOE-26S-PF TRT FP16",
        "yoloe26_trt_fp8":                      "YOLOE-26S-PF TRT FP8",
        "yolo_seg":                             "YOLO11s-seg PyTorch FP32",
        "yolo_seg_fp16_trt":                    "YOLO11s-seg TRT FP16",
        "yolo_seg_fp8_trt":                     "YOLO11s-seg TRT FP8 ★ shipping",
        "clip_trt":                             "OpenCLIP visual TRT ★ shipping",
    }

    headers = ["Family", "Workload", "Kernels", "DRAM / forward", "BW-bound FPS ceiling"]
    rows = []
    highlight_rows = []
    for i, (_, fam, wid, k, mb, fps) in enumerate(rows_raw):
        label = pretty.get(wid, wid)
        mb_str = f"{mb/1000:,.1f} GB" if mb >= 1000 else f"{mb:,.1f} MB"
        fps_str = f"{fps:,.1f}" if fps >= 10 else f"{fps:,.2f}"
        rows.append([fam, label, f"{k:,}", mb_str, fps_str])
        if "★ shipping" in label:
            highlight_rows.append(i + 1)

    add_styled_table(
        slide, Inches(CONTENT_LEFT), Inches(1.75),
        Inches(CONTENT_W), Inches(4.6), headers, rows,
        highlight_rows=highlight_rows, font_size=9, header_font_size=10,
    )

    add_bullet_box(slide, CONTENT_LEFT, 6.45, CONTENT_W, 0.85, [
        ("Reading the ceiling column", C.ACCENT_BLUE, True),
        ("• BW-bound FPS is the best case — compute may cut it further. For shipping TRT FP8 the two numbers "
         "agree: YOLO @ 465 ceiling vs ~140 measured, CLIP @ 232 ceiling at 1 Hz — either way 30 fps fits "
         "with headroom. For SAM 3 the ceiling is 0.8 FPS: there is no CPU-side optimization that gets past it.",
         C.ACCENT_GREEN),
        ("• Source: data/output/ncu/sizer_bundle.json (vendored to keyhole-sizer/sizer/). Regenerate with "
         "scripts/export_ncu_for_sizer.py after any ncu re-run.", C.TEXT_DIM),
    ], font_size=10)


def slide_optimization_roadmap(prs: Presentation):
    """Slide: Path to real-time on edge hardware."""
    slide = new_slide(prs)
    add_title_subtitle(slide, "Optimization Roadmap — Path to Edge Real-Time",
                       "Model changes required to fit within 134.4 GB/s bandwidth budget")

    headers = ["Optimization", "Traffic Reduction", "Est. Edge FPS", "Status"]
    rows = [
        ["SAM 3 BF16, 9 prompts (baseline)",  "1× (~147 GB)",               "0.4 FPS",   "MEASURED"],
        ["Lower input resolution (720p)",     "~1× (no change)",            "0.6 FPS",   "TESTED — not viable"],
        ["Reduce internal resolution",        "N/A",                        "N/A",       "BLOCKED — RoPE locked"],
        ["INT8 weight-only quantization",     "Weights only (not traffic)", "0.6 FPS",   "TESTED — no speedup"],
        ["Fewer prompts (1 vs 9)",            "~0.6× (decoder only)",       "0.9 FPS",   "TESTED — helps on desktop"],
        ["INT8 activation quantization",      "~2× (halve act traffic)",    "~1.2 FPS",  "Research-grade effort"],
        ["FP8 activation (E4M3)",             "~2× (halve act traffic)",    "~1.2 FPS",  "Measured on ES-Small"],
        ["INT4 activation quantization",      "~4×",                        "~2.4 FPS",  "Significant accuracy risk"],
        ["EfficientSAM / MobileSAM",          "~50-100× (5-50M params)",    "~15-30 FPS","Bake-off completed"],
        ["Hybrid V2 (YOLO-seg + CLIP) BF16",   "1× (CLIP dominates 22 ms)",  "2.9 FPS",   "MEASURED"],
        ["Hybrid V2 + FP8/INT8 on CLIP",       "~2× on CLIP half",           "4.9 FPS",   "MEASURED (48/72 Linears)"],
        ["Hybrid V2 + CLIP @ 1 Hz (N=30)",      "~30× on CLIP amortized",    "16.0 FPS",  "MEASURED — 93% of YOLO ceiling"],
        ["YOLO-seg INT8 via torchao 1×1 swap",   "~22% BW savings on YOLO",   "23.8 FPS",  "MEASURED — 49/50 swapped (44% wts)"],
        ["YOLO-seg FP8 via TensorRT 10.16",       "~50% on YOLO (full model)", "36.8 FPS",  "MEASURED — recall 1.00, IoU 0.998"],
        ["CLIP visual FP8 via TensorRT 10.16",     "~50% on CLIP (full model)", "66.3 FPS",  "MEASURED — top-1 agree 0.964"],
        ["Hybrid V2 + CLIP every-frame (all TRT)", "stacked, no debounce",       "24 FPS",    "PROJECTED — real-time, simplest"],
        ["Hybrid V2 + 1 Hz CLIP (all TRT)",        "stacked, debounced",         "36 FPS",    "PROJECTED — at YOLO ceiling"],
        ["LLM: Qwen3-30B-A3B Q4_K_M (MoE) — NPU Mid",  "3B active / 30B total",     "37.85 tok/s","VENDOR actual (128-bit LPDDR5X @ 8.4 GT/s)"],
        ["LLM: Qwen3-30B-A3B Q4_K_M — NPU Low-LP5",    "64-bit LPDDR5 @ 6.4 GT/s",  "29.27 tok/s","VENDOR actual — lower-bin bus"],
        ["LLM: Qwen3-30B-A3B Q4_K_M — NPU High",       "higher-bin LPDDR5X",        "50.46 tok/s","VENDOR actual — headroom for concurrent load"],
        ["Vision + LLM concurrent (short query)",      "duty-cycle sharing NPU Mid","~30 FPS",    "PROJECTED — 10 queries/min, 200 tok"],
        ["4-stream concurrent (YOLO batch=4)",      "batching amortizes overhead","25.9 FPS/stream","MEASURED batching, edge-projected"],
        ["8-stream concurrent (YOLO batch=8)",      "further batching gain",      "15.1 FPS/stream","MEASURED batching, edge-projected"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                     Inches(CONTENT_W), Inches(3.9), headers, rows,
                     highlight_rows=[17],  # Full-stack 1 Hz all-TRT — shipping target
                     font_size=9, header_font_size=10)

    add_bullet_box(slide, CONTENT_LEFT, 5.45, CONTENT_W, 1.6, [
        ("Bake-off sequence — all measurements complete", C.ACCENT_PURPLE, True),
        ("1-3. ES-Small quant: FP8 (94/95) & plain INT8 both → 4.9 FPS edge; SmoothQuant CONVERT blocked by torchao 0.17",),
        ("4-5. Hybrid V2 CLIP torchao FP8/INT8 (48/72) → 4.9 FPS; 1 Hz keyframe debounce → 16 FPS (93% of YOLO ceiling)",),
        ("6-7. YOLO-seg Conv: torchao 1×1 swap INT8 → 23.8 FPS (partial); TRT 10.16 full Conv-FP8 → 36.8 FPS (+98%, recall 1.00)",),
        ("8. CLIP visual FP8 via TRT → 29.8 → 15.1 ms edge (+120% CLIP FPS); full TRT stack projects 36 FPS shipping",),
        ("9. LLM — Qwen3-30B-A3B MoE (Q4/Q5/Q8): NPU Low/Mid/High vendor actuals: 29 / 38 / 50 tok/s Q4_K_M; duty-cycle chart quantifies vision+LLM coexistence",),
        ("10. Multi-stream concurrency — TRT YOLO dynamic-batch: 4 streams @ 26 FPS each (not 9), 8 @ 15, batching amortizes kernel overhead",),
        ("11. [PARTIAL] SAM 3 Lite watch — community shipped EfficientSAM3 ES-EV-S Apr 2026; benched at 2.59 FPS @ 720p NPU Mid (6.5× over SAM 3, still 13× behind our shipping stack). Meta SAM 3.1 / distillations still data-center-only.",),
    ], font_size=9)


def gather_platform_specs() -> dict:
    """Read host machine specs directly from /proc, /sys, lscpu, nvidia-smi, and torch."""
    import subprocess, shutil
    def _read(path):
        try:
            return Path(path).read_text().strip()
        except Exception:
            return ""
    def _sh(cmd):
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=5)
            return r.stdout.strip()
        except Exception:
            return ""

    # CPU
    cpu_info = {}
    lscpu = _sh(["lscpu"])
    for line in lscpu.splitlines():
        if ":" in line:
            k, v = line.split(":", 1)
            cpu_info[k.strip()] = v.strip()

    # RAM
    meminfo = _read("/proc/meminfo")
    mem_total_kb = 0
    for line in meminfo.splitlines():
        if line.startswith("MemTotal:"):
            mem_total_kb = int(line.split()[1])
            break
    ram_gb = round(mem_total_kb / (1024 * 1024), 0)

    # System / motherboard (no sudo)
    sys_vendor = _read("/sys/devices/virtual/dmi/id/sys_vendor") or "?"
    product = _read("/sys/devices/virtual/dmi/id/product_name") or "?"
    board_vendor = _read("/sys/devices/virtual/dmi/id/board_vendor") or "?"
    board_name = _read("/sys/devices/virtual/dmi/id/board_name") or "?"
    bios_ver = _read("/sys/devices/virtual/dmi/id/bios_version") or "?"

    # OS
    os_pretty = "?"
    for line in _read("/etc/os-release").splitlines():
        if line.startswith("PRETTY_NAME="):
            os_pretty = line.split("=", 1)[1].strip().strip('"')
            break
    kernel = _sh(["uname", "-r"])

    # Storage (real disks only)
    storage_lines = []
    try:
        r = subprocess.run(["lsblk", "-d", "-n", "-o", "NAME,SIZE,MODEL,TYPE"],
                           capture_output=True, text=True, timeout=5)
        for line in r.stdout.splitlines():
            parts = line.split(None, 3)
            if len(parts) >= 4 and parts[3].strip() == "disk":
                storage_lines.append(f"/dev/{parts[0]}  {parts[1]}  {parts[2]}")
    except Exception:
        pass

    # GPU via nvidia-smi + torch
    gpu_specs = {}
    try:
        r = subprocess.run(
            ["nvidia-smi", "--query-gpu=name,driver_version,compute_cap,"
             "memory.total,power.max_limit,clocks.max.sm,clocks.max.memory",
             "--format=csv,noheader,nounits"],
            capture_output=True, text=True, timeout=5)
        parts = [p.strip() for p in r.stdout.split(",")]
        if len(parts) >= 7:
            gpu_specs = {
                "name": parts[0], "driver": parts[1], "compute_cap": parts[2],
                "vram_mib": int(parts[3]), "tdp_w": float(parts[4]),
                "sm_clock_mhz": int(parts[5]), "mem_clock_mhz": int(parts[6]),
            }
    except Exception:
        pass
    try:
        import torch
        p = torch.cuda.get_device_properties(0)
        gpu_specs["sm_count"] = p.multi_processor_count
        gpu_specs["l2_cache_mb"] = round(getattr(p, "L2_cache_size", 0) / 1e6, 1)
        gpu_specs["cuda_runtime"] = torch.version.cuda
        gpu_specs["cudnn"] = torch.backends.cudnn.version()
    except Exception:
        pass

    return {
        "cpu_model": cpu_info.get("Model name", "?"),
        "cpu_cores": cpu_info.get("Core(s) per socket", "?"),
        "cpu_threads": cpu_info.get("CPU(s)", "?"),
        "cpu_max_mhz": cpu_info.get("CPU max MHz", "?"),
        "cpu_l3": cpu_info.get("L3 cache", "?"),
        "ram_gb": ram_gb,
        "system": f"{sys_vendor} {product}",
        "motherboard": f"{board_vendor} {board_name}  (BIOS {bios_ver})",
        "os": os_pretty,
        "kernel": kernel,
        "storage": storage_lines,
        "gpu": gpu_specs,
    }


def slide_npu_tier_specs(prs: Presentation):
    """Slide: canonical NPU tier assumptions — what each 'edge FPS' number in
    the deck is projected against. Four tiers × (memory, BW, TOPS, capacity,
    TDP, vendor LLM benchmark). Data mirrors keyhole-sizer/sizer/npu_model.py
    (the TIERS dict) so the deck and sizer website agree — if the sizer gets
    a new tier or a spec update, update both in the same commit."""
    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(
        slide,
        "Edge NPU tier assumptions",
        "Silicon class each 'edge FPS' / 'edge tok/s' number in this deck projects against",
    )

    # Row format mirrors the sizer's describe_hw output — single source of truth.
    headers = ["Tier", "Memory bus", "BW theoretical", "BW effective (70%)",
               "Tensor TOPS", "DRAM", "TDP", "LLM Q4 decode", "LLM TTFT @ 1K"]
    rows = [
        ["NPU Low-LP5",  "64-bit LPDDR5 @ 6.4 GT/s",  "51.2 GB/s",  "35.84 GB/s",
         "2 INT8 (dense)",                     "16 GB", "10 W", "29.27 tok/s", "1.67 s"],
        ["NPU Low-LP5X", "64-bit LPDDR5X @ 8.4 GT/s", "67.2 GB/s",  "47.04 GB/s",
         "50 BF16 / 100 INT8 / 100 FP8",       "16 GB", "10 W", "— (projected)", "— (projected)"],
        ["NPU Mid",      "128-bit LPDDR5X @ 8.4 GT/s","134.4 GB/s", "94.08 GB/s",
         "200 BF16 / 400 INT8 / 400 FP8",      "24 GB", "25 W", "37.85 tok/s", "0.351 s"],
        ["NPU High",     "128-bit LPDDR5X @ 11.2 GT/s","179.2 GB/s","125.44 GB/s",
         "275 BF16 / 550 INT8 / 550 FP8",      "32 GB", "40 W", "50.46 tok/s", "0.176 s"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
                     Inches(CONTENT_W), Inches(2.5), headers, rows,
                     highlight_rows=[3],   # NPU Mid — shipping target
                     font_size=10, header_font_size=11)

    # Context bullets + assumptions callout
    add_bullet_box(slide, CONTENT_LEFT, 4.2, CONTENT_W, 2.6, [
        ("How edge FPS numbers in this deck are derived", C.ACCENT_BLUE, True),
        ("• Reference measurement happens on the RTX 5090 (1792 GB/s theo × 0.85 eff = 1523.2 GB/s realized). "
         "Edge ms projects via bandwidth ratio: edge_ms = 5090_ms × (5090_eff_bw / edge_eff_bw). "
         "No TOPS-based compute ceiling in the current math — every tier is treated as bandwidth-bound.",
         C.TEXT_BRIGHT),
        ("• 70% bandwidth efficiency is uniform across all four edge tiers — removes tier-specific efficiency games so cross-tier comparisons reflect silicon differences, not modeling assumptions. "
         "Derived from published NPU vendor benchmarks on Qwen3-30B-A3B Q4_K_M (135 ms TTFT vs 1K prompt on Mid).",
         C.TEXT_DIM),
        "",
        ("Precision + silicon-class notes", C.ACCENT_BLUE, True),
        ("• NPU Low-LP5 is INT8-ONLY silicon (NXP i.MX 95 Neutron N3-1024S class, 2 TOPS dense INT8). "
         "Floating-point pipelines (BF16 / FP8) will either fail to load on this tier or fall back to CPU with catastrophic slowdown. "
         "Real-world FP-pipeline comparison should target Low-LP5X and above.",
         C.ACCENT_AMBER),
        ("• NPU Low-LP5X through High have native BF16/FP8 tensor cores. Tensor TOPS column lists BF16 / INT8 / FP8 peaks per NVIDIA-class spec conventions.",
         C.TEXT_BRIGHT),
        ("• LLM decode + TTFT are vendor-published measurements on Qwen3-30B-A3B Q4_K_M (1K prompt, short response). "
         "Low-LP5X has no vendor LLM benchmark yet — sizer falls back to bandwidth-ratio projection from Low-LP5.",
         C.TEXT_DIM),
    ], font_size=10)


def slide_platform_specs(prs: Presentation):
    """Slide: detail the measurement workstation so readers have a grounded ref."""
    s = gather_platform_specs()
    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "Measurement Platform — Keyhole Development Workstation",
                       "Every 5090 number in this deck was measured here. Edge projections project FROM this.")

    gpu = s["gpu"]
    vram_gb = round(gpu.get("vram_mib", 0) / 1024, 0) if gpu.get("vram_mib") else "?"
    sm_ghz = gpu.get("sm_clock_mhz", 0) / 1000

    # ─── System column (left) ───
    sys_items = [
        ("Host", C.ACCENT_BLUE, True),
        f"System:  {s['system']}",
        f"Motherboard:  {s['motherboard']}",
        f"OS:  {s['os']}  (kernel {s['kernel']})",
        "",
        ("CPU", C.ACCENT_BLUE, True),
        f"{s['cpu_model']}",
        f"{s['cpu_cores']} cores / {s['cpu_threads']} threads  •  boost {s['cpu_max_mhz']} MHz  •  L3 {s['cpu_l3']}",
        "",
        ("Memory", C.ACCENT_BLUE, True),
        f"{int(s['ram_gb'])} GB system RAM (DDR5)",
        "",
        ("Storage", C.ACCENT_BLUE, True),
    ] + [f"• {line}" for line in s["storage"]]
    add_bullet_box(slide, CONTENT_LEFT, CONTENT_TOP, CONTENT_W / 2 - 0.15, 3.2,
                    sys_items, font_size=11)

    # ─── GPU column (right) ───
    gpu_items = [
        ("GPU — the measurement reference", C.ACCENT_GREEN, True),
        f"{gpu.get('name', '?')}",
        f"Compute capability {gpu.get('compute_cap', '?')} (Blackwell)  •  {gpu.get('sm_count', '?')} SMs",
        f"{vram_gb} GB GDDR7  •  1792 GB/s bandwidth (512-bit @ 28 Gbps)",
        f"Boost clock {sm_ghz:.3f} GHz  •  Mem {gpu.get('mem_clock_mhz', 0)} MHz",
        f"TDP {int(gpu.get('tdp_w', 0))} W  •  L2 cache {gpu.get('l2_cache_mb', '?')} MB",
        f"Driver {gpu.get('driver', '?')}  •  CUDA {gpu.get('cuda_runtime', '?')}  •  cuDNN {gpu.get('cudnn', '?')}",
        "",
        ("Peak tensor throughput (dense)", C.ACCENT_GREEN, True),
        "• FP32 shader:    104.8 TFLOPS",
        "• BF16 / FP16:    209.5 TFLOPS (tensor)",
        "• INT8:           419 TOPS (tensor)",
        "• FP8  (E4M3):    419 TFLOPS (tensor — Blackwell)",
        "• INT4 / FP4:     838 TFLOPS (tensor — Blackwell 5th-gen)",
        "• Sparse 2:4:     values above × 2",
    ]
    add_bullet_box(slide, CONTENT_LEFT + CONTENT_W / 2 + 0.15, CONTENT_TOP,
                    CONTENT_W / 2 - 0.15, 3.4, gpu_items, font_size=11)

    # ─── Precision support table (full width) ───
    headers = ["Precision", "Tensor-core support", "Used in Keyhole", "Notes"]
    rows = [
        ["FP32",          "Yes (TF32 tensor)",          "—",                  "Default torch dtype; no quantization payoff"],
        ["BF16",          "Yes",                         "SAM 3 baseline, all PyTorch runs", "Preserves FP32 range, half the bytes"],
        ["FP16",          "Yes",                         "TRT baseline",       "Same bytes as BF16, narrower range"],
        ["INT8",          "Yes (dynamic act)",           "torchao + TRT YOLO + ES-Small",    "96-99% box recall; halves activation bytes"],
        ["FP8 (E4M3)",    "Yes — Blackwell native",      "TRT YOLO + TRT CLIP + ES-Small",  "Best quality/speed; halves activation bytes"],
        ["INT4 / FP4",    "Yes — Blackwell 5th-gen",     "Not yet exercised",  "~2× over INT8/FP8; accuracy risk on detection heads"],
        ["GGUF  Q4_K_M",  "Runs via llama.cpp CUDA",     "—",                  "Practical ~70B LLM ceiling on 32 GB VRAM"],
        ["GGUF  Q8_0",    "Runs via llama.cpp CUDA",     "—",                  "Practical ~14B LLM at long context"],
    ]
    add_styled_table(slide, Inches(CONTENT_LEFT), Inches(5.0),
                     Inches(CONTENT_W), Inches(2.0), headers, rows,
                     highlight_rows=[5])  # FP8 is the headline


def slide_summary(prs: Presentation, runs: list[dict], targets: dict):
    """Final summary slide — the 90× journey from 0.4 to 36 FPS."""
    slide = new_slide(prs, bg_color=C.BG_DARK)
    add_title_subtitle(slide, "Summary & Key Findings — 0.4 FPS → 36 FPS (90× Edge Improvement)",
                       "What Keyhole proved, ruled out, and what ships")

    # Hero stat bar (indigo)
    hero = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(CONTENT_LEFT), Inches(CONTENT_TOP),
        Inches(CONTENT_W), Inches(0.7))
    hero.fill.solid(); hero.fill.fore_color.rgb = C.ACCENT_INDIGO
    hero.line.fill.background()
    add_text_box(slide, Inches(CONTENT_LEFT + 0.2), Inches(CONTENT_TOP + 0.1),
                 Inches(CONTENT_W - 0.4), Inches(0.5),
                 "Shipping stack  •  Hybrid V2 + YOLO-seg FP8 + CLIP FP8 (all TensorRT)  •  720p Edge MPU (134.4 GB/s LPDDR5X)  •  36 FPS projected",
                 font_size=14, color=C.TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    items = [
        ("The journey", C.ACCENT_BLUE, True),
        ("0.4 FPS",   C.ACCENT_RED, True),
        ("  SAM 3 BF16 — bandwidth-bound on 134.4 GB/s LPDDR5X. 840M params, 3.71 GB activations.",),
        ("  Cheap levers tested & ruled out: weight-only INT8 (no act savings), lower resolution (RoPE locked), fewer prompts (encoder floor).",),
        ("4.9 FPS",   C.ACCENT_ORANGE, True),
        ("  Model-family switch → EfficientSAM-Small + FP8/INT8 activation quant (torchao, 94/95 Linears). ΔIoU < 0.002.",),
        ("16 FPS",    C.ACCENT_AMBER, True),
        ("  Hybrid V2 (YOLO-seg + CLIP) + FP8 CLIP + 1 Hz keyframe debounce. YOLO becomes the ceiling.",),
        ("24 FPS",    C.ACCENT_GREEN, True),
        ("  YOLO-seg FP8 via TensorRT 10.16 on Blackwell (recall 1.00, IoU 0.998). CLIP FP8 every frame — no debounce needed.",),
        ("36 FPS  ← shipping", C.ACCENT_INDIGO, True),
        ("  Full TRT stack + CLIP @ 1 Hz. YOLO-only ceiling. Room for INT4/FP4 if the edge NPU exposes them.",),
    ]
    add_bullet_box(slide, CONTENT_LEFT, 2.2, CONTENT_W / 2 - 0.15, 4.6, items, font_size=11)

    right = [
        ("What we proved",  C.ACCENT_GREEN, True),
        "• Bandwidth — not compute — sets the edge ceiling for vision transformers",
        "• FP8 activation quant is near-lossless on ViT + detection heads on Blackwell",
        "• Hybrid pipelines (YOLO-seg + CLIP) beat a monolithic big-ViT mask model for edge",
        "• CLIP keyframe debouncing is an optional lever, not a hard requirement once TRT-compiled",
        "",
        ("What we ruled out", C.ACCENT_RED, True),
        "• SAM 3 BF16 on 134.4 GB/s — not feasible without model replacement",
        "• INT8 weight-only quantization — doesn't touch activation traffic, no edge gain",
        "• torchao FP8 on Conv-only models (YOLO-seg) — tool-chain gap, use TensorRT",
        "",
        ("What remains open", C.ACCENT_AMBER, True),
        "• Meta releasing quantized SAM 3 / SAM 3 Lite (passive watch)",
        "• INT4/FP4 on detection head — accuracy risk; warrants targeted study if NPU ships with it",
        "• Live streaming subsystem — architecture locked (MJPEG+WS, YOLO-FP8), ~500-line build when prioritized",
    ]
    add_bullet_box(slide, CONTENT_LEFT + CONTENT_W / 2 + 0.15, 2.2,
                    CONTENT_W / 2 - 0.15, 4.6, right, font_size=11)


# ============================================================
# Main Build
# ============================================================

@click.command()
@click.option("--output", "-o", default="data/output/keyhole_results.pptx",
              help="Output PPTX path")
@click.option("--runs-dir", default="data/output/runs",
              help="Directory containing run JSON files")
@click.option("--data-dir", default="data/output",
              help="Directory containing reference architecture exports")
def build_deck(output, runs_dir, data_dir):
    """Generate the Keyhole results PowerPoint deck."""
    from rich.console import Console
    console = Console()

    console.print("\n[bold]Keyhole — Building Results Deck[/]\n")

    runs_dir = Path(runs_dir)
    data_dir = Path(data_dir)
    output = Path(output)

    # Load data
    runs = load_runs(runs_dir)
    sam3_ref = load_sam3_reference(data_dir)
    targets = load_npu_targets()

    console.print(f"  Runs found:     {len(runs)}")
    console.print(f"  SAM 3 ref:      {'YES' if sam3_ref else 'NO'}")
    console.print(f"  NPU targets:    {len(targets)}")

    # Build presentation (widescreen 16:9)
    prs = Presentation()
    set_deck_size(prs)

    # Slide 1: Title
    console.print("  Building: Title slide")
    slide_title(prs)

    # Slide 2: Executive summary — the world's best single slide
    # (If you read nothing else, read this one)
    console.print("  Building: Executive summary (front of deck)")
    slide_exec_summary(prs)

    # Slide 3: Platform specs (read live from /proc, /sys, nvidia-smi, torch)
    console.print("  Building: Platform specs")
    slide_platform_specs(prs)

    # Slide 4: Edge NPU tier assumptions — canonical spec table for readers
    # who want to understand what each "edge FPS" projection implies about
    # silicon. Mirrors keyhole-sizer/sizer/npu_model.py::TIERS.
    console.print("  Building: NPU tier assumptions")
    slide_npu_tier_specs(prs)

    # Slide 5: Architecture
    console.print("  Building: Architecture diagram")
    slide_architecture(prs)

    # Slide 3: SAM 3 Reference
    if sam3_ref:
        console.print("  Building: SAM 3 reference breakdown")
        slide_sam3_reference(prs, sam3_ref)

    # Slide 4: Roofline model
    console.print("  Building: Roofline model")
    slide_roofline(prs, targets, sam3_ref)

    # Per-run slides
    for i, run in enumerate(runs):
        video_name = run.get("video", {}).get("name", f"Run {i+1}")
        console.print(f"  Building: Run results — {video_name}")
        slide_run_results(prs, run, i)
        slide_npu_projections(prs, run, targets)

    # Comparison slide (if multiple runs)
    if len(runs) >= 1:
        console.print("  Building: Run comparison chart")
        slide_run_comparison(prs, runs, targets)

    # Bandwidth wall analysis
    console.print("  Building: Bandwidth wall analysis")
    slide_bandwidth_wall(prs)

    # Bandwidth requirements
    console.print("  Building: Bandwidth requirements")
    slide_bandwidth_requirements(prs)

    # Quantization results
    console.print("  Building: Quantization tested (weight-only INT8)")
    slide_quantization_tested(prs)

    # Activation quantization challenges
    console.print("  Building: Activation quantization challenges")
    slide_activation_quant_challenges(prs)

    # Prompt scaling analysis
    console.print("  Building: Prompt count scaling")
    slide_prompt_scaling(prs)

    # Resolution lock finding
    console.print("  Building: Resolution lock analysis")
    slide_resolution_lock(prs)

    # Model comparison
    console.print("  Building: Model comparison (speed vs accuracy)")
    slide_model_comparison(prs, data_dir)

    # Hybrid V2 breakthrough
    console.print("  Building: Hybrid V2 breakthrough")
    slide_hybrid_v2(prs)

    # Mask-model bake-off (only if data present)
    if Path("data/output/bakeoff/720p_EW_clip/summary.json").exists():
        console.print("  Building: Mask bake-off summary")
        slide_bakeoff_summary(prs)
        console.print("  Building: Mask bake-off visuals")
        slide_bakeoff_visuals(prs)

    # FP8 activation quantization (only if data present)
    if Path("data/output/bakeoff/fp8_edge_projection.json").exists():
        console.print("  Building: FP8 activation quantization")
        slide_fp8_quantization(prs)

    # SmoothQuant + plain INT8 (only if data present)
    if Path("data/output/bakeoff/smoothquant_edge_projection.json").exists():
        console.print("  Building: SmoothQuant + INT8")
        slide_smoothquant(prs)

    # Hybrid V2 CLIP-quantization bake-off (only if data present)
    if Path("data/output/bakeoff/hybrid_v2_edge_projection.json").exists():
        console.print("  Building: Hybrid V2 CLIP quantization bake-off")
        slide_hybrid_v2_bakeoff(prs)

    # CLIP keyframe debouncing (only if data present)
    if Path("data/output/bakeoff/keyframe_debounce_summary.json").exists():
        console.print("  Building: CLIP keyframe debouncing")
        slide_keyframe_debounce(prs)

    # YOLO-seg conv quantization (INT8 via 1x1 swap; FP8 blocked)
    if Path("data/output/bakeoff/yolo_conv_quant_edge_projection.json").exists():
        console.print("  Building: YOLO-seg conv quantization")
        slide_yolo_conv_quant(prs)

    # TRT proper FP8/INT8 on YOLO-seg (real unblock)
    if Path("data/output/bakeoff/trt_yolo_edge_projection.json").exists():
        console.print("  Building: TensorRT YOLO FP8/INT8")
        slide_trt_yolo(prs)

    # yolo11s-seg vs yolov8n-seg cross-variant comparison (silicon apples-to-apples)
    if Path("data/output/bakeoff/trt_yolo_yolov8n-seg_summary.json").exists():
        console.print("  Building: yolo11s-seg vs yolov8n-seg comparison")
        slide_yolov8n_comparison(prs)

    # TRT CLIP visual tower
    if Path("data/output/bakeoff/trt_clip_edge_projection.json").exists():
        console.print("  Building: TensorRT CLIP visual")
        slide_trt_clip(prs)

    # LLM bake-off + duty-cycle trade-off (only if data present)
    if Path("data/output/bakeoff/llm_edge_projection.json").exists():
        console.print("  Building: LLM bake-off (Qwen3-30B-A3B)")
        slide_llm_bakeoff(prs)
        console.print("  Building: NPU duty-cycle trade-off")
        slide_llm_duty_cycle(prs)

    # Multi-stream concurrency (YOLO batching)
    if Path("data/output/bakeoff/concurrency_edge_projection.json").exists():
        console.print("  Building: Multi-stream concurrency")
        slide_concurrency(prs)

    # Community SAM 3 Lite post-ship watch (roadmap #9)
    if Path("data/output/bakeoff/efficientsam3_summary.json").exists():
        console.print("  Building: EfficientSAM3 community bake-off (SAM 3 Lite watch)")
        slide_efficientsam3_community(prs)

    # EfficientSAM3.1 text-prompt-capable smaller variant (SAM 3.1 student)
    if Path("data/output/bakeoff/efficientsam3p1_summary.json").exists():
        console.print("  Building: EfficientSAM3.1 text-prompt variant (SAM 3.1 student)")
        slide_efficientsam3p1_textprompt(prs)

    # YOLOE-26 one-model open-vocab (Option B watch)
    if Path("data/output/bakeoff/yoloe26_summary.json").exists():
        console.print("  Building: YOLOE-26 one-model open-vocab (Option B watch)")
        slide_yoloe26_onemodel(prs)

    # TRT YOLOE-26 — does FP8 close the gap? (negative result)
    if Path("data/output/bakeoff/trt_yoloe26_summary.json").exists():
        console.print("  Building: TRT YOLOE-26 — does FP8 close the gap?")
        slide_trt_yoloe26(prs)

    # TRT takeaways — synthesize the 3 TRT bake-offs + edge-compiler caveat
    if (Path("data/output/bakeoff/trt_yolo_edge_projection.json").exists()
        and Path("data/output/bakeoff/trt_clip_edge_projection.json").exists()
        and Path("data/output/bakeoff/trt_yoloe26_summary.json").exists()):
        console.print("  Building: TRT takeaways — where it pays and where it doesn't")
        slide_trt_takeaways(prs)

    # Nsight Compute measured DRAM bandwidth (16 workloads)
    if Path("data/output/ncu/sizer_bundle.json").exists():
        console.print("  Building: ncu measured DRAM — headline 515× gap")
        slide_ncu_headline(prs)
        console.print("  Building: ncu measured DRAM — full workload table")
        slide_ncu_workload_table(prs)

    # Optimization roadmap
    console.print("  Building: Optimization roadmap")
    slide_optimization_roadmap(prs)

    # Summary slide
    console.print("  Building: Summary & findings")
    slide_summary(prs, runs, targets)

    # Add consistent footer (project + page number) to every slide
    finalize_footers(prs)

    # Save
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))

    total_slides = len(prs.slides)
    console.print(f"\n  [bold green]Deck generated: {output}[/]")
    console.print(f"  Total slides: {total_slides}")
    console.print(f"  Runs included: {len(runs)}\n")


if __name__ == "__main__":
    build_deck()

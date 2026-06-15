#!/usr/bin/env python3
"""build_exec_oneslide.py — ONE CEO-level 'why do I care' slide.
See-spot-run headlines + hero numbers + a detail line each. Public/strategic framing
(no Skippy specifics). Output: data/output/exec-why-fp4-matters.pptx."""
import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO / "scripts"))
from build_deck import C, set_deck_size, new_slide, add_text_box

OUT = REPO / "data" / "output" / "exec-why-fp4-matters.pptx"
CARD_BG = RGBColor(0x16, 0x21, 0x3E)

CARDS = [
    {"accent": C.ACCENT_BLUE, "kicker": "WHERE THE INDUSTRY IS GOING",
     "hero": "Go floating-point,\nor fall behind",
     "num": "2028", "numfont": 44, "numlabel": "integer-only silicon = locked out",
     "detail": "NVIDIA, Qualcomm and TI all ship 4-bit floating-point (FP4). The numeric "
               "FORMAT — not the chip size — is the differentiator. Integer-only past "
               "~2028 can't run frontier AI."},
    {"accent": C.ACCENT_GREEN, "kicker": "THE PERFORMANCE",
     "hero": "Twice the speed —\nbasically free",
     "num": "~2×", "numfont": 44, "numlabel": "faster, at ~2% accuracy cost",
     "detail": "FP4 ~halves latency and shrinks the model 4×, for only ~2 points of "
               "accuracy. We measured the accuracy cost is real-but-tiny — the speed-up "
               "isn't a trade-off."},
    {"accent": C.ACCENT_ORANGE, "kicker": "THE PRODUCT",
     "hero": "A smarter doc AI —\nthat only runs on FP4",
     "num": "+11% / +18%", "numfont": 30, "numlabel": "better at finding / answering",
     "detail": "Reading documents as images beats today's text pipeline — better at "
               "finding the answer AND getting it right. Edge memory is ~15× tighter, so "
               "FP4 is what makes it run on a device."},
]


def main():
    prs = Presentation(); set_deck_size(prs)
    s = new_slide(prs, bg_color=C.BG_DARK)

    add_text_box(s, Inches(0.45), Inches(0.40), Inches(12.45), Inches(0.6),
                 "Why FP4 floating-point is the edge-AI bet that matters",
                 font_size=25, bold=True, color=C.TEXT_WHITE)
    add_text_box(s, Inches(0.45), Inches(1.06), Inches(12.45), Inches(0.4),
                 "Three reasons to care — every number measured on real silicon, all reproducible.",
                 font_size=13.5, color=C.ACCENT_BLUE)

    cw, gap, top, ch = 3.97, 0.31, 1.60, 4.18
    xs = [0.45 + i * (cw + gap) for i in range(3)]
    for x, card in zip(xs, CARDS):
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(cw), Inches(ch))
        box.fill.solid(); box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = card["accent"]; box.line.width = Pt(1.75)
        box.shadow.inherit = False
        # accent top stripe
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(top), Inches(cw), Inches(0.10))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = card["accent"]; stripe.line.fill.background()
        stripe.shadow.inherit = False
        tx = x + 0.26; tw = cw - 0.52
        add_text_box(s, Inches(tx), Inches(top + 0.24), Inches(tw), Inches(0.32),
                     card["kicker"], font_size=10.5, bold=True, color=card["accent"])
        add_text_box(s, Inches(tx), Inches(top + 0.58), Inches(tw), Inches(1.0),
                     card["hero"], font_size=19, bold=True, color=C.TEXT_WHITE)
        add_text_box(s, Inches(tx), Inches(top + 1.62), Inches(tw), Inches(0.66),
                     card["num"], font_size=card["numfont"], bold=True, color=card["accent"])
        add_text_box(s, Inches(tx), Inches(top + 2.30), Inches(tw), Inches(0.3),
                     card["numlabel"], font_size=11, bold=True, color=C.TEXT_DIM)
        add_text_box(s, Inches(tx), Inches(top + 2.66), Inches(tw), Inches(1.4),
                     card["detail"], font_size=9.5, color=C.TEXT_BRIGHT)

    # bottom-line bar
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(6.08), Inches(12.45), Inches(0.92))
    bar.fill.solid(); bar.fill.fore_color.rgb = C.ACCENT_INDIGO; bar.line.fill.background(); bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "THE BOTTOM LINE   "
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r2 = p.add_run()
    r2.text = ("Back FP4-capable edge silicon: it's where the industry is going, it's a near-free 2× in speed, "
               "and it unlocks a better edge-AI product. Measured, not modeled.")
    r2.font.size = Pt(12.5); r2.font.color.rgb = RGBColor(0xEE, 0xEE, 0xFF)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT} (1 slide)")


if __name__ == "__main__":
    main()

#!/usr/bin/env python3
"""
build_precision_pixelrag_merge.py — append a PixelRAG/visual-RAG section to the precision
roadmap deck, producing data/output/precision-roadmap-with-pixelrag.pptx.

PRIVACY: this merged deck contains Skippy-derived (aggregate-only) charts. It is written
LOCALLY and is NOT pushed to my-stuff/gdrive — distribution is Kyle's call. The canonical
precision-roadmap-combined.pptx (the published one) is left untouched.

Run:  python scripts/build_precision_pixelrag_merge.py
"""
import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO / "scripts"))
from build_deck import (C, new_slide, add_text_box, add_bullet_box,
                        SLIDE_W_IN, SLIDE_H_IN, CONTENT_LEFT, CONTENT_W)

SRC = REPO / "data" / "output" / "precision-roadmap-combined.pptx"
OUT = REPO / "data" / "output" / "precision-roadmap-with-pixelrag.pptx"


def chart_slide(prs, png, title, caption):
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 title, font_size=20, color=C.ACCENT_BLUE, bold=True)
    bl, bt, bw, bh = 0.6, 1.12, SLIDE_W_IN - 1.2, 5.5
    iw, ih = Image.open(REPO / png).size; ar = iw / ih
    w, h = (bh * ar, bh) if bw / bh > ar else (bw, bw / ar)
    s.shapes.add_picture(str(REPO / png), Inches(bl + (bw - w) / 2), Inches(bt), Inches(w), Inches(h))
    add_text_box(s, Inches(CONTENT_LEFT), Inches(SLIDE_H_IN - 0.5), Inches(CONTENT_W),
                 Inches(0.4), caption, font_size=11.5, color=C.TEXT_DIM)


def main():
    prs = Presentation(str(SRC))
    n0 = len(prs.slides._sldIdLst)

    # divider
    s = new_slide(prs, bg_color=C.BG_DARK)
    add_text_box(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(0.5),
                 "Applied", font_size=16, color=C.ACCENT_BLUE, bold=True)
    add_text_box(s, Inches(0.8), Inches(3.15), Inches(11.7), Inches(1.0),
                 "Visual-RAG on the edge — a precision case study", font_size=32, color=C.TEXT_WHITE, bold=True)
    add_text_box(s, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.7),
                 "PixelRAG: render documents to images, retrieve over pixels, read with a VL model — "
                 "measured on Skippy's own corpus (RTX 5090 → edge NPU)", font_size=14, color=C.TEXT_DIM)

    # explainer
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 "What is visual-RAG, and why does precision matter here?", font_size=20,
                 color=C.ACCENT_BLUE, bold=True)
    add_bullet_box(s, CONTENT_LEFT, 1.4, CONTENT_W, 5.6, [
        ("Read documents as PICTURES, not parsed text.", C.ACCENT_GREEN, True),
        "  Render page→image · embed images · retrieve · read the page image with a VL model. No OCR — "
        "keeps tables/layout that text parsing destroys (arXiv 2506.05209).",
        ("It moves document understanding onto the INT8-friendly vision encoder + cuts LLM tokens.",
         C.ACCENT_GREEN, True),
        "  Exactly the keyhole split — vision stays INT, the LLM reader goes FP — and it attacks both "
        "edge walls: compute-bound prefill and bandwidth-bound decode.",
        ("Measured on Skippy: visual-RAG beats text-RAG on retrieval, and FP4 makes it edge-viable.",
         C.ACCENT_BLUE, True),
    ])

    chart_slide(prs, "data/output/skippy_pixelrag_knee.png",
                "The resolution knee — accuracy AND tokens",
                "Tune resolution to the retrieval knee → visual-RAG wins both axes vs native-text RAG")
    chart_slide(prs, "data/output/skippy_precision_arm.png",
                "Reader precision — BF16 / FP8 / NVFP4 measured on vLLM",
                "NVFP4 decode 1.9× (prefill diluted by the un-quantized vision encoder)")
    chart_slide(prs, "data/output/skippy_edge_projection.png",
                "Projected to the edge NPU — FP4 is the enabler",
                "~15× less bandwidth → BF16 won't fit 8GB; NVFP4 fits (2.2GB) + halves decode")
    chart_slide(prs, "data/output/vision_quant_combined.png",
                "Mixed precision in practice: the LM is the lever, vision is a sliver",
                "Measured (RTX 5090): FP4 on the language model ~doubles the prefill win vs FP8; quantize the LM, keep the vision tower INT8")
    chart_slide(prs, "data/output/accuracy_study.png",
                "Answer accuracy — visual-RAG vs text-RAG, and the FP4 cost",
                "DocVQA ANLS: reading the page image beats OCR-text by +18 pts; FP4 costs only ~2 pts vs BF16 — quantization is accuracy-safe")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"{SRC.name} ({n0}) -> {OUT.name} ({len(prs.slides._sldIdLst)} slides)  [LOCAL ONLY — not pushed]")


if __name__ == "__main__":
    main()

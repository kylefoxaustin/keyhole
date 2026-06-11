#!/usr/bin/env python3
"""
build_precision_roadmap_deck.py — the SELF-CONTAINED precision-roadmap deck.

Writes data/output/precision-roadmap-combined.pptx from scratch using the shared
build_deck dark theme. It NO LONGER depends on a base file in ~/Downloads (that
dependency is gone — the title / sources / takeaways slides are authored here in
code and driven from precision_migration.json + precision_blue.json).

Narrative (12 slides):
   1  Title + thesis
   2  INT vs FP over time — the FP half grows         (NEW · precision_blue_overtime)
   3  7-year migration timeline                       (precision_migration)
   4  Today's measured INT/FP composition             (precision_composition)
   5  The weights question: INT4 vs FP4               (NEW · precision_blue_weights)
   6  Decode is bandwidth-bound — ladder + MoE        (precision_5090_ladder)
   7  INT4 vs FP4 — memory format ≠ compute format    (precision_5090_breadth)
   8  Asymmetry is architecture-general               (precision_5090_breadth_v2)
   9  FP4 wins the whole lifecycle                     (precision_5090_fp4_lifecycle)
  10  FP4 training converges to BF16 quality           (precision_5090_fp4_convergence)
  11  Sources / citations
  12  Takeaways

Run:  python3 scripts/build_precision_roadmap_deck.py
      (regenerate the charts first: precision_blue.py + plot_precision_blue.py,
       precision_migration.py + plot_precision_migration.py, etc.)
"""
import json
import os
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO / "scripts"))
from build_deck import (  # noqa: E402
    C, set_deck_size, new_slide, add_text_box, add_bullet_box,
    SLIDE_W_IN, SLIDE_H_IN, CONTENT_LEFT, CONTENT_W,
)

OUT = REPO / "data" / "output" / "precision-roadmap-combined.pptx"

# (png, title, caption) for the full-bleed chart slides, in narrative order.
CHARTS = [
    ("data/output/precision_blue_overtime.png",
     "INT vs FP over time — and the FP (blue) half grows",
     "Deploy-precision mix per workload, 2026→2033 · grey = stays INT · light blue = FP-optional · dark blue = FP-mandatory"),
    ("data/output/precision_migration.png",
     "The 7-year precision migration timeline",
     "INT8 holds the floor; FP stops being optional — FP8 by ~2028, FP4 by ~2030"),
    ("data/output/precision_composition.png",
     "2026 — per-model INT/FP composition (the journey's anchor)",
     "Where the bits actually go today, by params / FLOPs / stored bytes — grey = INT, dark blue = FP-mandatory"),
    ("data/output/precision_composition_2030.png",
     "2030 — same breakout, next generation goes FP4-native",
     "Flagships lead to FP4 (light blue); vision stays all-INT; edge LLMs hold INT longest"),
    ("data/output/precision_composition_2033.png",
     "2033 — FP4 is the default; FP-mandatory grows sub-4-bit",
     "Edge LLMs move too; sub-4-bit outlier handling turns part of even the FP4 engine FP-mandatory (dark grows)"),
    ("data/output/precision_blue_weights.png",
     "The weights question: do they stay INT4/INT8, or go FP4?",
     "INT4 & FP4 TIE on decode (both 4-bit memory) — but FP4 alone wins prefill 3.6× and training 5.5×, measured on RTX 5090"),
    ("data/output/precision_5090_ladder.png",
     "Decode is bandwidth-bound — measured ladder + MoE decoupling",
     "MoE decouples memory footprint from decode speed"),
    ("data/output/precision_5090_breadth.png",
     "INT4 is a memory format; FP4 is a memory + compute format",
     "Same-base quads on the RTX 5090 — the split widens with scale, then plateaus"),
    ("data/output/precision_5090_breadth_v2.png",
     "Not a one-model fluke — INT4 ties FP4 on decode, loses on prefill, everywhere",
     "12 models · 6 architectures · every one ties on decode (4-bit memory) and splits 2.7–4.5× on prefill (FP4 has the compute path, INT4 doesn't)"),
    ("data/output/precision_5090_fp4_lifecycle.png",
     "FP4 wins the whole lifecycle — inference AND training",
     "The win grows with compute-intensity (FP4 is a compute format)"),
    ("data/output/precision_5090_fp4_convergence.png",
     "FP4 training converges to BF16 quality",
     "The accuracy half of 'native FP4 training can be optimal'"),
]


def slide_chart(prs, png, title, caption):
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 title, font_size=20, color=C.ACCENT_BLUE, bold=True)
    box_l, box_t, box_w, box_h = 0.6, 1.12, SLIDE_W_IN - 1.2, 5.55
    iw, ih = Image.open(REPO / png).size
    ar = iw / ih
    if box_w / box_h > ar:           # height-bound
        h = box_h; w = h * ar
    else:                            # width-bound
        w = box_w; h = w / ar
    left = box_l + (box_w - w) / 2.0
    s.shapes.add_picture(str(REPO / png), Inches(left), Inches(box_t), Inches(w), Inches(h))
    add_text_box(s, Inches(CONTENT_LEFT), Inches(SLIDE_H_IN - 0.52), Inches(CONTENT_W),
                 Inches(0.4), caption, font_size=11.5, color=C.TEXT_DIM)
    return s


def slide_title(prs, mig):
    s = new_slide(prs, bg_color=C.BG_DARK)
    add_text_box(s, Inches(0.7), Inches(1.7), Inches(12), Inches(1.1),
                 "The 7-Year Edge-AI Precision Roadmap", font_size=40, bold=True,
                 color=C.TEXT_WHITE)
    add_text_box(s, Inches(0.7), Inches(3.15), Inches(12), Inches(0.7),
                 "INT8 holds the floor — but FP support stops being optional",
                 font_size=21, color=C.ACCENT_BLUE)
    add_text_box(s, Inches(0.7), Inches(4.25), Inches(12.0), Inches(2.0),
                 mig["__meta__"]["thesis"], font_size=15, color=C.TEXT_BRIGHT)
    add_text_box(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.5),
                 "Industry direction 2026 → 2033, ranked by workload urgency · "
                 "2026 measured on RTX 5090, later horizons forecast.",
                 font_size=11, color=C.TEXT_DIM)
    return s


def slide_sources(prs, mig):
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 "Sources — citations for the roadmap claims", font_size=20,
                 color=C.ACCENT_BLUE, bold=True)
    cit = mig["citations"]
    groups = [("VLA quantization & deployment", cit["vla_quant"][:6]),
              ("LLM precision (INT → FP8 → FP4)", cit["llm_precision"][:6]),
              ("Competitive edge silicon", cit["silicon"][:5])]
    x = CONTENT_LEFT
    colw = CONTENT_W / 3 - 0.1
    for title, items in groups:
        add_text_box(s, Inches(x), Inches(1.15), Inches(colw), Inches(0.4),
                     title, font_size=12.5, color=C.ACCENT_GREEN, bold=True)
        body = "\n".join("• " + it for it in items)
        add_text_box(s, Inches(x), Inches(1.6), Inches(colw), Inches(5.2),
                     body, font_size=8.6, color=C.TEXT_BRIGHT)
        x += CONTENT_W / 3
    add_text_box(s, Inches(CONTENT_LEFT), Inches(SLIDE_H_IN - 0.5), Inches(CONTENT_W),
                 Inches(0.4), "Full list in data/output/precision_migration.json · "
                 "peer-reviewed venues, arxiv, vendor announcements.",
                 font_size=10, color=C.TEXT_DIM)
    return s


def slide_takeaways(prs):
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 "What it says for silicon", font_size=22, color=C.ACCENT_BLUE, bold=True)
    add_bullet_box(s, CONTENT_LEFT, 1.45, CONTENT_W, 5.4, [
        ("Param count is the wrong axis — precision FORMAT is the right one.",
         C.ACCENT_GREEN, True),
        "  The bits-per-weight is not the whole story: a 4-bit INT and a 4-bit FLOAT "
        "stream the same bytes (tie on decode) but behave nothing alike on compute.",
        ("Weights are moving to FP4 — because FP4 is the only 4-bit format that also "
         "accelerates compute.", C.ACCENT_GREEN, True),
        "  Measured, same weights: INT4 prefill 1.0× (stuck on the BF16 floor) vs FP4 "
        "3.6×; training-GEMM FP4 5.5×. The win grows with compute intensity, and FP4 "
        "trains to BF16 quality.",
        ("The silicon ask, by horizon:", C.ACCENT_BLUE, True),
        "  Vision stays INT8 throughout. NPU Mid needs FP8 (E4M3) by ~2028; NPU High "
        "needs FP4/NVFP4 by ~2030. Shipping INT8-only past 2028 is a competitive "
        "feature gap, not a silicon-area win.",
        ("All anchors measured on one RTX 5090 (Blackwell sm_120) — reproducible.",
         C.ACCENT_BLUE, True),
    ])
    return s


def main():
    mig = json.load(open(REPO / "data" / "output" / "precision_migration.json"))
    prs = Presentation()
    set_deck_size(prs)
    slide_title(prs, mig)
    for png, title, cap in CHARTS:
        slide_chart(prs, png, title, cap)
    slide_sources(prs, mig)
    slide_takeaways(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()

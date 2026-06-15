#!/usr/bin/env python3
"""
build_pixelrag_deck.py — package the PixelRAG-on-Skippy bake-off into a result deck.

Self-contained (dark build_deck theme). Embeds the two aggregate charts (knee +
precision) — both contain METRICS ONLY, no personal content. The deck is written to
data/output/pixelrag-skippy-result.pptx; do NOT auto-publish (Skippy-derived) — hand
to Kyle to decide distribution.

Run:  python scripts/build_pixelrag_deck.py
"""
import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO / "scripts"))
from build_deck import (C, set_deck_size, new_slide, add_text_box, add_bullet_box,
                        SLIDE_W_IN, SLIDE_H_IN, CONTENT_LEFT, CONTENT_W)

OUT = REPO / "data" / "output" / "pixelrag-skippy-result.pptx"


def chart_slide(prs, png, title, caption):
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 title, font_size=20, color=C.ACCENT_BLUE, bold=True)
    bl, bt, bw, bh = 0.6, 1.12, SLIDE_W_IN - 1.2, 5.5
    iw, ih = Image.open(REPO / png).size; ar = iw / ih
    w, h = (bh * ar, bh) if bw / bh > ar else (bw, bw / ar)
    s.shapes.add_picture(str(REPO / png), Inches(bl + (bw - w) / 2), Inches(bt), Inches(w), Inches(h))
    add_text_box(s, Inches(CONTENT_LEFT), Inches(SLIDE_H_IN - 0.5), Inches(CONTENT_W),
                 Inches(0.4), caption, font_size=11.5, color=C.TEXT_DIM)
    return s


def main():
    prs = Presentation(); set_deck_size(prs)

    # 1 — title
    s = new_slide(prs, bg_color=C.BG_DARK)
    add_text_box(s, Inches(0.7), Inches(1.9), Inches(12), Inches(1.1),
                 "Visual-RAG (PixelRAG) on Skippy", font_size=40, bold=True, color=C.TEXT_WHITE)
    add_text_box(s, Inches(0.7), Inches(3.05), Inches(12), Inches(0.7),
                 "Render documents to images, retrieve over pixels — measured on Skippy's own corpus, RTX 5090",
                 font_size=18, color=C.ACCENT_BLUE)
    add_text_box(s, Inches(0.7), Inches(4.2), Inches(12), Inches(1.6),
                 "1,218 pages (datasheets / documents / writing) · Qwen3-VL-Embedding-2B + Qwen3-VL-4B · "
                 "240 local-generated queries · all local, aggregate metrics only.",
                 font_size=14, color=C.TEXT_DIM)
    add_text_box(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
                 "Source: github.com/StarTrail-org/PixelRAG (arXiv 2506.05209) · keyhole bake-off harness",
                 font_size=11, color=C.TEXT_DIM)

    # 2 — WHAT IS IT (explainer): what / why / promise / how
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 "What is visual-RAG (PixelRAG)?", font_size=22, color=C.ACCENT_BLUE, bold=True)
    add_bullet_box(s, CONTENT_LEFT, 1.4, CONTENT_W, 5.6, [
        ("What it is — RAG that reads documents as PICTURES, not parsed text.", C.ACCENT_GREEN, True),
        "  Render each page to an image, retrieve over the images, and feed the page image straight to a "
        "vision-language model that reads layout + text together. (Berkeley/Princeton/EPFL/Databricks, "
        "arXiv 2506.05209; lineage: ColPali → DeepSeek-OCR optical compression.)",
        ("Why it exists — text parsing throws away structure.", C.ACCENT_GREEN, True),
        "  OCR / HTML-to-text flattens tables, pinout diagrams, charts and reading order — so the reader "
        "literally can't find answers that live in the layout. Pixels keep it intact.",
        ("What it promises — higher retrieval accuracy + (claimed) ~10× fewer tokens.", C.ACCENT_GREEN, True),
        "  vs legacy text-RAG, across document QA. The token claim is the contentious part — we measured it.",
        ("How it works — two VL models, one pixel pipeline.", C.ACCENT_BLUE, True),
        "  render page→image · embed images (Qwen3-VL-Embedding-2B) → FAISS · retrieve top-k · "
        "read the page image with a VL reader (Qwen3-VL-4B) → answer. No OCR anywhere.",
    ])

    # 3 — RESULTS AT A GLANCE
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 "Results at a glance — measured on Skippy's corpus", font_size=22, color=C.ACCENT_BLUE, bold=True)
    add_bullet_box(s, CONTENT_LEFT, 1.4, CONTENT_W, 5.6, [
        ("Retrieval: pixel beats text — recall@5 0.70 vs 0.59 (and survives every control).", C.ACCENT_GREEN, True),
        ("Token cost: the '10×' is a KNOB — at the ~768px knee, pixel is CHEAPER than text "
         "(423 vs 538 reader tokens); at full res it's 3× MORE.", C.ACCENT_GREEN, True),
        ("Precision: NVFP4 reader = 1.9× decode / 1.4× prefill (vision-diluted), 62% answer-match.",
         C.ACCENT_GREEN, True),
        ("Fusion: pixel-only wins — adding text just dilutes it.", C.ACCENT_BLUE, True),
        ("Edge: ~15× less bandwidth → BF16 won't fit 8GB; NVFP4 is the enabler (fits + halves decode).",
         C.ACCENT_BLUE, True),
        ("Bottom line: visual-RAG works on Skippy's own knowledge, and FP4 makes it edge-viable.",
         C.ACCENT_ORANGE, True),
    ])

    # 4 — the five findings (detail)
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 "Five measured findings", font_size=22, color=C.ACCENT_BLUE, bold=True)
    add_bullet_box(s, CONTENT_LEFT, 1.4, CONTENT_W, 5.6, [
        ("1 · Visual-RAG beats text-RAG on retrieval — on every content type.", C.ACCENT_GREEN, True),
        "  Pixel recall@5 0.70 vs native-text 0.59 (datasheets, documents, writing, even emails).",
        ("2 · There is a resolution knee (~768px) — and it's a DOUBLE win.", C.ACCENT_GREEN, True),
        "  At the knee, the page image is BOTH more accurate AND cheaper than text (423 vs 538 reader "
        "tokens). Above it, image tokens balloon (1536px = 3x text). The '10x tokens' is a knob, not free.",
        ("3 · FP4 reader: 1.9x decode, but prefill only 1.4x (vision-diluted).", C.ACCENT_GREEN, True),
        "  Measured NVFP4 vs BF16: decode 1.88x (bandwidth-bound). Prefill only 1.44x because the vision "
        "encoder runs un-quantized — NOT the deck's pure-LLM 3.6x. Quality: 62.5% answer-agreement vs BF16.",
        ("4 · Pixel-only wins — hybrid (pixel+text) just dilutes it.", C.ACCENT_BLUE, True),
        ("5 · The win is REAL, not embedder bias — a dedicated text retriever (BGE) doesn't close it.",
         C.ACCENT_BLUE, True),
    ])

    # 3 — knee chart
    chart_slide(prs, "data/output/skippy_pixelrag_knee.png",
                "The resolution knee — accuracy AND tokens",
                "Tune tile resolution to the retrieval knee → visual-RAG wins both axes vs native-text RAG")

    # 4 — precision chart
    chart_slide(prs, "data/output/skippy_precision_arm.png",
                "Reader precision — measured directly on vLLM",
                "BF16 / FP8 / NVFP4 at the knee; NVFP4 prefill diluted by the un-quantized vision encoder")

    # 5 — edge projection chart
    chart_slide(prs, "data/output/skippy_edge_projection.png",
                "Projected to the edge NPU — bandwidth-starved",
                "~same compute (prefill ×1.2) but ~15× less bandwidth (decode ×15): BF16 won't fit 8GB; NVFP4 is the enabler")

    chart_slide(prs, "data/output/vision_quant_combined.png",
                "Mixed precision: the LM is the lever, the vision tower is a sliver",
                "Measured: FP4 on the language model ~doubles the prefill win vs FP8; quantizing the vision tower is minor (1–8%) — keep it INT8")

    # rigor & controls (the polish follow-ups)
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 "Rigor — we tried to break the result", font_size=22, color=C.ACCENT_BLUE, bold=True)
    add_bullet_box(s, CONTENT_LEFT, 1.4, CONTENT_W, 5.6, [
        ("Confound: is it just that the VL embedder is better at images than text?", C.ACCENT_ORANGE, True),
        "  Control with a DEDICATED text retriever (BGE) + chunking to fix truncation: BGE-chunked rises "
        "to recall@5 0.59 (beats the VL-text arm) — but pixel still wins 0.69. The advantage is REAL.",
        ("Caveat: the 'documents' class had recall@1 ≈ 0.", C.ACCENT_ORANGE, True),
        "  Cause: podcast transcripts — pages within one doc are near-duplicates, so exact-page R@1 is "
        "unwinnable. At DOCUMENT granularity it's healthy: pixel doc-recall@5 0.83 vs text 0.78.",
        ("We used PixelRAG's ACTUAL retriever — verified by running their own code.", C.ACCENT_GREEN, True),
        "  PixelRAG's retriever IS Qwen3-VL-Embedding-2B (their default) = the model we used. Their "
        "4B 'screenshot-LoRA' is a READER/VQA adapter (task=CAUSAL_LM), not a retriever — run faithfully "
        "as an embedder it scores far below the 2B (page R@5 0.05–0.15 vs 0.69). So our study is faithful.",
        ("Open / honest limits.", C.ACCENT_PURPLE, True),
        "  Synthetic queries (local-generated); modest absolute recall (1,218-page pool); NVFP4 answer-"
        "match 62%. All reproducible, all aggregate, no personal content.",
    ])

    # final — the edge recipe + takeaways
    s = new_slide(prs)
    add_text_box(s, Inches(CONTENT_LEFT), Inches(0.42), Inches(CONTENT_W), Inches(0.55),
                 "The edge recipe", font_size=22, color=C.ACCENT_BLUE, bold=True)
    add_bullet_box(s, CONTENT_LEFT, 1.45, CONTENT_W, 5.4, [
        ("Visual-RAG is viable at the edge — if you tune two knobs.", C.ACCENT_GREEN, True),
        "  (a) Resolution → the retrieval knee (~768px): max accuracy at minimum image tokens.",
        "  (b) Precision → vision-encode INT8 (the keyhole vision thesis) + LLM reader FP4 (decode 1.9x).",
        ("This is the natural keyhole split: vision stays INT, the LLM reader goes FP.", C.ACCENT_BLUE, True),
        "  PixelRAG moves document understanding onto the INT8-friendly vision encoder and cuts LLM "
        "input tokens — attacking both edge walls (compute-bound prefill + bandwidth-bound KV).",
        ("Honest caveats kept in view.", C.ACCENT_ORANGE, True),
        "  Synthetic queries; modest absolute recall (1,218-page pool); NVFP4 quality flag (62.5%); "
        "BGE-control truncates long pages. All reproducible; all aggregate, no personal content.",
    ])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()

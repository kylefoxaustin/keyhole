"""
Update the corporate-template branded deck with the latest bake-off data,
in-place, preserving all branding/styling. Targets:

  Slide  1 — title date
  Slide  4 — Edge NPU tier assumptions: add i.MX 95 + RTX 5090 rows,
             re-label header, rewrite measured-silicon-anchors bullets
  Slide 46 — TRT CLIP: update all 9 table rows + 2 bullet lines
  Slide 54 — TRT takeaways: update CLIP row's FP8 result cell

Run:
    python scripts/update_branded_deck.py

Reads/writes: data/output/keyhole_deck_branded.pptx
"""
from __future__ import annotations
import copy
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from lxml import etree

DECK = Path("data/output/keyhole_deck_branded.pptx")

NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def replace_run_text(shape, old: str, new: str) -> int:
    """Replace `old` with `new` across all runs in all paragraphs of a
    text-frame shape. Returns # of runs touched. Preserves formatting
    because we only assign to run.text."""
    n = 0
    if not shape.has_text_frame:
        return 0
    for para in shape.text_frame.paragraphs:
        # Try per-run first (preserves formatting)
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                n += 1
        # If old string spans across runs (e.g. "15.1" is "1","5",".","1"),
        # concatenate paragraph text, check, fall back to para-level rewrite.
        if n == 0 and old in para.text:
            # Rewrite at paragraph level: drop all runs into the first, keep
            # formatting of the first run.
            runs = list(para.runs)
            if runs:
                full = para.text.replace(old, new)
                runs[0].text = full
                for r in runs[1:]:
                    r.text = ""
                n += 1
    return n


def set_cell_text(cell, text: str):
    """Set a table cell's text while preserving run-0 formatting."""
    tf = cell.text_frame
    # Preserve formatting of first run in first paragraph if possible.
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        # Clear everything except first run, then set first run text.
        first_para = tf.paragraphs[0]
        # Clear trailing runs in para 0
        for r in list(first_para.runs)[1:]:
            r.text = ""
        first_run.text = text
        # Drop any extra paragraphs.
        for p in list(tf.paragraphs)[1:]:
            p._p.getparent().remove(p._p)
    else:
        tf.text = text


def clone_table_row_after(tbl_xml, src_row_idx: int) -> etree._Element:
    """Clone <a:tr> at src_row_idx and insert it after. Returns new row."""
    rows = tbl_xml.findall("a:tr", NSMAP)
    src = rows[src_row_idx]
    new = copy.deepcopy(src)
    src.addnext(new)
    return new


def set_row_cells(tr, values: list[str]):
    """Set cell text for each <a:tc> under the given row <a:tr>, preserving
    formatting of each cell's first run."""
    cells = tr.findall("a:tc", NSMAP)
    for tc, val in zip(cells, values):
        txbody = tc.find("a:txBody", NSMAP)
        if txbody is None:
            continue
        paras = txbody.findall("a:p", NSMAP)
        if not paras:
            continue
        p0 = paras[0]
        runs = p0.findall("a:r", NSMAP)
        if runs:
            # Clear text in all but first run, then set first run text.
            first = runs[0]
            t = first.find("a:t", NSMAP)
            if t is not None:
                t.text = val
            for r in runs[1:]:
                t2 = r.find("a:t", NSMAP)
                if t2 is not None:
                    t2.text = ""
        # Drop trailing paragraphs
        for p in paras[1:]:
            txbody.remove(p)


def replace_run_text_any(slide, old: str, new: str) -> int:
    """Run replace_run_text across every shape in a slide."""
    total = 0
    for shape in slide.shapes:
        total += replace_run_text(shape, old, new)
    return total


def update_slide_1(slide):
    """Date: April 22 → April 24."""
    for shape in slide.shapes:
        if replace_run_text(shape, "April 22, 2026", "April 24, 2026"):
            return True
    return False


def update_slide_2_exec_summary(slide):
    """Slide 2 tier table at bottom was overflowing — header 'Vision
    (4-stream, batch=4)' and row-1 'Good for' were wrapping to 2 lines,
    blowing the 0.22"-per-row budget and pushing the table past the
    'Assumes vision/LLM time-slice on the NPU' footer. Fix by abbreviating
    the wrapping strings and setting uneven column widths that give the
    long 'Good for' column more room."""
    from pptx.util import Inches, Pt, Emu
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        t = shape.table
        # Only touch the NPU tier table (sniff header cell 0)
        if t.cell(0, 0).text.strip() != "NPU tier":
            continue

        # --- (1) Rewrite wrap-prone strings to shorter forms
        rewrites = {
            "Vision (4-stream, batch=4)": "Vision (4-stream)",
            "LLM Q4_K_M decode": "LLM Q4 decode",
            "64-bit LPDDR5 @ 6.4 GT/s": "64-bit LPDDR5 6.4 GT/s",
            "128-bit LPDDR5X @ 8.4 GT/s": "128-bit LPDDR5X 8.4 GT/s",
            "Dense INT8-only silicon (NXP Neutron class)": "Dense INT8 (Neutron)",
            "Live multi-stream + occasional LLM": "Multi-stream + occasional LLM",
        }
        for i in range(len(t.rows)):
            for j in range(len(t.columns)):
                cell = t.cell(i, j)
                for old, new in rewrites.items():
                    if cell.text.strip() == old:
                        set_cell_text(cell, new)

        # --- (2) Widen 'Good for' column (col 5), steal from narrower ones.
        # Original: 6 cols × 2.05" = 12.30" total.
        # New: 1.8 / 1.9 / 1.2 / 1.3 / 1.3 / 4.8 = 12.30" total.
        target_widths_in = [1.8, 1.9, 1.2, 1.3, 1.3, 4.8]
        for col, w_in in zip(t.columns, target_widths_in):
            col.width = Inches(w_in)

        # --- (3) Explicit compact row heights so auto-grow can't bully us
        for row in t.rows:
            row.height = Inches(0.22)

        # --- (4) Shrink table font to 8pt (header to 9) to be safe
        for i in range(len(t.rows)):
            for j in range(len(t.columns)):
                cell = t.cell(i, j)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(9 if i == 0 else 8)
        return True
    return False


def update_slide_46_clip(slide):
    """TRT CLIP table + two bullet lines."""
    # Target data (from data/output/bakeoff/trt_clip_edge_projection.json, rerun 2026-04-24)
    new_rows = [
        # (res, recipe, 5090 ms, top1, edge ms, fps)
        ("720p",  "BF16 torch", "2.36", "1.000", "30.7", "32.6"),
        ("720p",  "TRT FP16",   "1.40", "0.970", "30.7", "32.6"),
        ("720p",  "TRT FP8",    "1.45", "0.964", "15.6", "64.2"),
        ("1080p", "BF16 torch", "1.67", "1.000", "21.7", "46.1"),
        ("1080p", "TRT FP16",   "1.24", "0.966", "21.7", "46.1"),
        ("1080p", "TRT FP8",    "1.24", "0.966", "11.0", "90.9"),
        ("4K",    "BF16 torch", "2.17", "1.000", "28.3", "35.3"),
        ("4K",    "TRT FP16",   "1.27", "0.966", "28.3", "35.3"),
        ("4K",    "TRT FP8",    "1.25", "0.966", "14.4", "69.7"),
    ]
    table_updated = False
    for shape in slide.shapes:
        if shape.has_table:
            t = shape.table
            # Skip header row, update rows 1..9
            for i, vals in enumerate(new_rows, 1):
                if i >= len(t.rows):
                    break
                for j, val in enumerate(vals):
                    if j >= len(t.columns):
                        break
                    set_cell_text(t.cell(i, j), val)
            table_updated = True
            break

    # Bullet-level replacements (branded bullets still had hardcoded 15.1/29.8/0.5)
    bullet_replacements = [
        ("FP8 edge CLIP drops from 29.8 ms (BF16/FP16) to 15.1 ms (+120% CLIP-only FPS).",
         "FP8 edge CLIP drops from 30.7 ms (BF16/FP16) to 15.6 ms (+97% CLIP-only FPS)."),
        ("YOLO-FP8 (27.2 ms) + CLIP-FP8 every frame (15.1 ms) = 42.3 ms → 24 FPS",
         "YOLO-FP8 (27.2 ms) + CLIP-FP8 every frame (15.6 ms) = 42.8 ms → 23 FPS"),
        ("YOLO-FP8 + CLIP-FP8 at 1 Hz (0.5 ms amortized) = 27.7 ms → 36 FPS",
         "YOLO-FP8 + CLIP-FP8 at 1 Hz (0.52 ms amortized) = 27.7 ms → 36 FPS"),
    ]
    bullets_updated = 0
    for shape in slide.shapes:
        for old, new in bullet_replacements:
            if replace_run_text(shape, old, new):
                bullets_updated += 1
    return table_updated, bullets_updated


def update_slide_54_trt_takeaways(slide):
    """CLIP row in decision matrix: 28300.1 → 30.7 / 14348.7 → 15.6."""
    n = 0
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text
                    if "28300.1" in txt and "14348.7" in txt:
                        new_txt = txt.replace("28300.1 ms BF16", "30.7 ms BF16").replace("14348.7 ms FP8", "15.6 ms FP8")
                        set_cell_text(cell, new_txt)
                        n += 1
    return n


def update_slide_4_tier_specs(slide):
    """Tier specs: (1) header "(70%)" moves into cells; (2) insert i.MX 95
    after Low-LP5-32bit; (3) insert RTX 5090 at end; (4) rewrite bullet
    section."""
    # Find the tier table
    tbl_shape = None
    for shape in slide.shapes:
        if shape.has_table and shape.table.cell(0, 0).text.strip() == "Tier":
            tbl_shape = shape
            break
    if tbl_shape is None:
        return False, False

    t = tbl_shape.table
    tbl_xml = t._tbl  # CT_Table element

    # --- (1) Header label: "BW effective (70%)" → "BW effective"
    hdr = t.cell(0, 3)
    if "(70%)" in hdr.text:
        set_cell_text(hdr, "BW effective")

    # --- (2) Each existing tier row: append "(70%)" to BW effective cell
    #         (or "(85%)" if RTX 5090 row appears — but we add that below)
    for i in range(1, len(t.rows)):
        cell = t.cell(i, 3)
        text = cell.text.strip()
        if text and "(" not in text:  # skip if already has pct marker
            set_cell_text(cell, f"{text} (70%)")

    # --- (3) Clone row after Low-LP5-32bit (row 1), populate with i.MX 95
    imx_values = [
        "NPU i.MX 95 (ground truth) †",
        "32-bit LPDDR5 @ 6.4 GT/s",
        "25.6 GB/s",
        "17.92 GB/s (70%)",
        "2 INT8 (Neutron NPU)",
        "16 GB",
        "10 W",
        "— (not evaluated)",
        "— (not evaluated)",
    ]
    new_tr = clone_table_row_after(tbl_xml, src_row_idx=1)
    set_row_cells(new_tr, imx_values)

    # --- (4) Clone last row (currently NPU High, was row 5 → now 6 after
    #         i.MX 95 insert), populate with RTX 5090
    # After the insert above, row order is:
    #   0 header, 1 Low-LP5-32bit, 2 i.MX 95, 3 Low-LP5-64bit,
    #   4 Low-LP5X, 5 Mid, 6 High
    rtx_values = [
        "RTX 5090 (reference, measured) †",
        "512-bit GDDR7 @ 28 GT/s",
        "1792 GB/s",
        "1523.2 GB/s (85%)",
        "~105 BF16 / ~210 FP8 / INT8 DP4A",
        "32 GB",
        "575 W",
        "250 tok/s",
        "0.165 s",
    ]
    last_tr = clone_table_row_after(tbl_xml, src_row_idx=6)  # clone High to make row 7
    set_row_cells(last_tr, rtx_values)

    # --- (5) Compact row heights so 8 rows fit the 2.6" frame (orig 6 rows
    #         × 0.43" = 2.58"; 8 × 0.30" = 2.40" leaves a 0.20" margin).
    from pptx.util import Inches, Pt
    for row in t.rows:
        row.height = Inches(0.30)

    # Shrink body-cell font to 9pt (was 10) so text fits the 0.30" rows.
    from pptx.oxml.ns import qn
    for i in range(1, len(t.rows)):
        for j in range(len(t.columns)):
            cell = t.cell(i, j)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(9)

    # --- (6) Rewrite the bullet section (replace "Precision + silicon-class
    #         notes" paragraph and three bullets with "Measured-silicon
    #         anchors (†)" + new bullets)
    bullet_replacements = [
        ("70% bandwidth efficiency is uniform across all five edge tiers — removes tier-specific efficiency games so cross-tier comparisons reflect silicon differences, not modeling assumptions. Derived from published NPU vendor benchmarks on Qwen3-30B-A3B Q4_K_M (135 ms TTFT vs 1K prompt on Mid).",
         "70% bandwidth efficiency is uniform across all five edge NPU tiers (RTX 5090 reference uses 85% — datacenter-class memory controller). Removes tier-specific efficiency games so cross-tier comparisons reflect silicon differences, not modeling assumptions."),
        ("Precision + silicon-class notes",
         "Measured-silicon anchors (†)"),
        ("NPU Low-LP5-32bit and -64bit are the SAME silicon class (INT8-only, 2 TOPS dense, NXP i.MX 95 Neutron N3-1024S family) paired with different memory bus widths. The 32-bit variant has half the bandwidth of the 64-bit. Floating-point pipelines (BF16 / FP8) will fail to load on either Low-LP5 tier and fall back to CPU at catastrophic slowdown.",
         "NPU i.MX 95 (ground truth): NXP eIQ Neutron NPU, real production measurement. yolov8n-seg INT8 @ 1080p = 32 ms (29.2 FPS) measured. Sizer surfaces this alongside the generic Low-LP5-32bit projection (18.3 FPS via BW scaling) — the 1.6× delta is honest evidence that pure-BW projection misses compute+overhead floor on weak silicon. Phase 2 compute-ceiling clamp (upcoming) uses this data point as the calibration anchor."),
        ("NPU Low-LP5X through High have native BF16/FP8 tensor cores. Tensor TOPS column lists BF16 / INT8 / FP8 peaks per NVIDIA-class spec conventions.",
         "RTX 5090 (reference, measured): every edge projection in the deck derives from 5090 measurements on Blackwell TRT 10.16. Exposed as a selectable tier in the sizer so users can see the raw measured numbers alongside edge projections."),
        ("LLM decode + TTFT are vendor-published measurements on Qwen3-30B-A3B Q4_K_M (1K prompt, short response). Tiers without vendor benchmarks (Low-LP5-32bit, Low-LP5X) BW-project from adjacent measured tiers.",
         "Low-LP5-32bit and -64bit are the SAME silicon class (INT8-only, 2 TOPS dense, Neutron-class) paired with different bus widths. Low-LP5X through High have native BF16/FP8 tensor cores. LLM decode + TTFT are vendor-published measurements on Qwen3-30B-A3B Q4_K_M."),
    ]
    bullets_updated = 0
    for shape in slide.shapes:
        for old, new in bullet_replacements:
            if replace_run_text(shape, old, new):
                bullets_updated += 1

    # --- (7) Shrink bullet font to 8pt so the "Measured-silicon anchors"
    #         section (which has a very long NPU i.MX 95 bullet) fits inside
    #         the 2.60" bullet box instead of bleeding past the slide edge.
    from pptx.util import Pt
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Skip the title and sub-header — only shrink the bullet paragraph box
        txt = shape.text_frame.text
        if "How edge FPS numbers" not in txt:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(8)
    return True, bullets_updated


def main():
    prs = Presentation(str(DECK))
    print(f"Opened: {DECK} ({len(prs.slides)} slides)")

    # Slide 1 — title date
    if update_slide_1(prs.slides[0]):
        print("  Slide  1 — date updated to April 24, 2026")
    else:
        print("  Slide  1 — date string not found (skipped)")

    # Slide 2 — exec summary tier table fit
    if update_slide_2_exec_summary(prs.slides[1]):
        print("  Slide  2 — exec-summary tier table rewrapped + column widths fixed")

    # Slide 4 — tier specs
    tbl_ok, bullets_n = update_slide_4_tier_specs(prs.slides[3])
    print(f"  Slide  4 — tier table +2 rows ({'ok' if tbl_ok else 'FAIL'}), bullets rewritten: {bullets_n}")

    # Slide 57 — Optimization Roadmap item-8 bullet had hardcoded stale
    # 29.8 → 15.1 ms / 120% CLIP FPS numbers (pre-dates today's rerun).
    roadmap_replace = replace_run_text_any(
        prs.slides[56],
        "CLIP visual FP8 via TRT → 29.8 → 15.1 ms edge (+120% CLIP FPS)",
        "CLIP visual FP8 via TRT → 30.7 → 15.6 ms edge (+97% CLIP FPS)",
    )
    print(f"  Slide 57 — roadmap item-8 CLIP bullet refreshed: {roadmap_replace}")

    # Slide 50 — "What just happened" text box was 0.10" past slide bottom.
    # Nudge the box up so it ends exactly at the slide bottom.
    from pptx.util import Inches
    for shape in prs.slides[49].shapes:
        if shape.has_text_frame and shape.text_frame.text.startswith("What just happened"):
            shape.top = Inches(5.80)
            print("  Slide 50 — 'What just happened' text box nudged up 0.10\"")
            break

    # Slide 46 — TRT CLIP
    tbl_ok, bullets_n = update_slide_46_clip(prs.slides[45])
    print(f"  Slide 46 — CLIP table rows updated ({'ok' if tbl_ok else 'FAIL'}), bullets: {bullets_n}")

    # Slide 54 — TRT takeaways
    n = update_slide_54_trt_takeaways(prs.slides[53])
    print(f"  Slide 54 — CLIP takeaway row: {n} cell updated")

    prs.save(str(DECK))
    print(f"Saved: {DECK}")


if __name__ == "__main__":
    main()
